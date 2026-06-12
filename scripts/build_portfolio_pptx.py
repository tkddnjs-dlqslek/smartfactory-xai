# -*- coding: utf-8 -*-
"""포트폴리오용 PPTX 빌더 — 실제 배포 사이트 캡처 중심.
수치는 슬라이드_데이터.md(SSOT, results/slide_data.json) 확정값 그대로 사용."""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
SHOTS = os.path.join(ROOT, "output", "_pptx_shots")
OUT = os.path.join(ROOT, "output", "스마트공장XAI_포트폴리오.pptx")

# 팔레트 (웹사이트 다크 테마와 통일)
BG = RGBColor(0x0B, 0x0D, 0x12)
CARD = RGBColor(0x14, 0x18, 0x21)
CARD_LINE = RGBColor(0x2A, 0x2F, 0x3A)
WHITE = RGBColor(0xF4, 0xF6, 0xF8)
GRAY = RGBColor(0x9C, 0xA3, 0xAF)
BODY = RGBColor(0xC9, 0xCE, 0xD6)
RED = RGBColor(0xEF, 0x44, 0x44)
CYAN = RGBColor(0x22, 0xD3, 0xEE)
AMBER = RGBColor(0xF5, 0x9E, 0x0B)
FONT = "맑은 고딕"

SW, SH = Inches(13.333), Inches(7.5)
prs = Presentation()
prs.slide_width, prs.slide_height = SW, SH
BLANK = prs.slide_layouts[6]

LIVE_URL = "smartfactory-xai.vercel.app"
GH_URL = "github.com/tkddnjs-dlqslek/smartfactory-xai"


def slide():
    s = prs.slides.add_slide(BLANK)
    r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    r.fill.solid(); r.fill.fore_color.rgb = BG
    r.line.fill.background()
    r.shadow.inherit = False
    return s


def _set_font(run, size, color, bold):
    run.font.name = FONT; run.font.size = Pt(size)
    run.font.color.rgb = color; run.font.bold = bold
    rPr = run._r.get_or_add_rPr()
    # lang="ko-KR" 필수 — 이게 있어야 PowerPoint가 한글을 단어 단위로 줄바꿈(음절 분리 방지)
    rPr.set('lang', 'ko-KR'); rPr.set('altLang', 'en-US')
    ea = rPr.find(qn('a:ea'))
    if ea is None:
        ea = rPr.makeelement(qn('a:ea'), {}); rPr.append(ea)
    ea.set('typeface', FONT)


def text(s, x, y, w, h, runs, size=14, color=BODY, bold=False, align=PP_ALIGN.LEFT,
         anchor=MSO_ANCHOR.TOP, line_spacing=1.0, space_after=0):
    """runs: str | [(txt, size, color, bold)] | [[run,...] per paragraph]"""
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    if isinstance(runs, str):
        runs = [[(runs, size, color, bold)]]
    elif runs and isinstance(runs[0], tuple):
        runs = [runs]
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; p.line_spacing = line_spacing
        if space_after: p.space_after = Pt(space_after)
        # 한글 단어 잘림 금지(사출성/형 방지) — 실측: ea폰트 지정 시 eaLnBrk="0" 단독일 때만 동작.
        # latinLnBrk 속성이 존재하면(값 무관) 한글이 음절 단위로 쪼개짐 — 절대 추가하지 말 것.
        p._p.get_or_add_pPr().set('eaLnBrk', '0')
        for (t, sz, c, b) in para:
            r = p.add_run(); r.text = t.replace("What-if", "What‑if")
            _set_font(r, sz, c, b)
    return tb


def card(s, x, y, w, h, fill=CARD, line=CARD_LINE):
    c = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    c.adjustments[0] = 0.06
    c.fill.solid(); c.fill.fore_color.rgb = fill
    c.line.color.rgb = line; c.line.width = Pt(1)
    c.shadow.inherit = False
    return c


def pic(s, path, x, y, w=None, h=None, border=True):
    p = s.shapes.add_picture(os.path.join(SHOTS, path), Inches(x), Inches(y),
                             Inches(w) if w else None, Inches(h) if h else None)
    if border:
        p.line.color.rgb = CARD_LINE; p.line.width = Pt(1.25)
    return p


def accent_bar(s, x=0.55, y=0.62, w=0.45):
    b = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(0.06))
    b.fill.solid(); b.fill.fore_color.rgb = RED; b.line.fill.background(); b.shadow.inherit = False


def footer(s, idx):
    text(s, 0.55, 7.08, 9, 0.3, [( "SmartFactory XAI — 통합 스마트공장 운영 플랫폼   ·   ", 9, GRAY, False),
                                  (LIVE_URL, 9, CYAN, False)], align=PP_ALIGN.LEFT)
    text(s, 12.3, 7.08, 0.5, 0.3, str(idx), 9, GRAY, align=PP_ALIGN.RIGHT)


def stat_card(s, x, y, w, h, label, value, sub, vcolor=CYAN):
    card(s, x, y, w, h)
    text(s, x + 0.18, y + 0.13, w - 0.36, 0.3, label, 10.5, GRAY, bold=True)
    text(s, x + 0.18, y + 0.42, w - 0.36, 0.62, value, 27, vcolor, bold=True)
    text(s, x + 0.18, y + h - 0.42, w - 0.36, 0.32, sub, 9.5, GRAY)


def feature_slide(idx, tab, title, bullets, img, img_w=7.85, note=None):
    s = slide()
    accent_bar(s)
    ix = 13.333 - 0.55 - img_w          # 이미지 좌측 x (기본 4.93)
    tw = ix - 0.55 - 0.25               # 텍스트 칼럼 폭 — 이미지와 겹치지 않게
    text(s, 0.55, 0.78, tw, 0.32, tab, 12, RED, bold=True)
    text(s, 0.55, 1.1, tw, 1.35, title, 21, WHITE, bold=True, line_spacing=1.05)
    paras = []
    for head, desc in bullets:
        paras.append([("▪  ", 12, RED, True), (head, 13, WHITE, True)])
        paras.append([(desc, 10.5, BODY, False)])
    tb = text(s, 0.55, 2.5, tw, 4.35, paras, line_spacing=1.12, space_after=6)
    for i, p in enumerate(tb.text_frame.paragraphs):
        if i % 2 == 1:  # 설명 단락 — 들여쓰기(줄바꿈돼도 정렬 유지)
            p._p.get_or_add_pPr().set('marL', str(Inches(0.26).emu))
    pic(s, img, ix, 1.05, w=img_w)  # 16:9 → h = img_w*9/16
    text(s, ix, 1.05 + img_w * 9 / 16 + 0.12, img_w, 0.3,
         [("LIVE  ", 9, RED, True), (f"https://{LIVE_URL} 실제 배포 화면 캡처", 9, GRAY, False)])
    if note:
        text(s, 0.55, 6.45, tw, 0.55, note, 10, AMBER, line_spacing=1.1)
    footer(s, idx)
    return s


# ───────────────────────── 1. 표지 ─────────────────────────
s = slide()
text(s, 0.9, 0.75, 11.5, 0.4, "PORTFOLIO  ·  2026 스마트 공장 운영 시스템 MVP 개발 해커톤 본선", 13, RED, bold=True)
text(s, 0.9, 1.2, 11.8, 1.0, "SmartFactory XAI", 48, WHITE, bold=True)
text(s, 0.9, 2.18, 11.8, 0.5, [("4-AI 합의 ", 19, CYAN, True), ("기반 사출성형 이상탐지 · 진단 · 처방 — ", 19, BODY, False),
                                ("품질 · 설비 · 안전 · 생산 통합 운영 플랫폼", 19, WHITE, True)])
text(s, 0.9, 2.78, 11.8, 0.4, [("Live  ", 12, GRAY, True), ("https://" + LIVE_URL, 12, CYAN, False),
                                ("      GitHub  ", 12, GRAY, True), ("https://" + GH_URL, 12, CYAN, False),
                                ("      Data  ", 12, GRAY, True), ("KAMP 사출성형기 공개 데이터셋", 12, BODY, False)])
pic(s, "01_landing_hero.png", 3.24, 3.4, w=6.85)  # h=3.85 → 하단 7.25 (슬라이드 안)
# ──────────────────────── 2. 프로젝트 개요 ────────────────────────
s = slide(); accent_bar(s)
text(s, 0.55, 0.78, 12, 0.5, "프로젝트 개요", 26, WHITE, bold=True)
text(s, 0.55, 1.45, 12.2, 1.15, [
    [("불량률 ", 13.5, BODY, False), ("1.03%", 13.5, RED, True),
     ("의 극심한 클래스 불균형 환경에서, ", 13.5, BODY, False), ("정상 데이터만 학습한 준지도 4-AI 앙상블", 13.5, WHITE, True),
     ("이 실시간 이상탐지 → SHAP 원인분석 → 처방까지 수행하고,", 13.5, BODY, False)],
    [("같은 엔진의 출력을 ", 13.5, BODY, False), ("품질·설비·안전·생산 4축", 13.5, WHITE, True),
     ("의 운영 의사결정으로 변환하는 풀스택 웹 플랫폼입니다.  (개인 개발 · 해커톤 본선 진출작 · 실배포 운영 중)", 13.5, BODY, False)],
], line_spacing=1.25, space_after=4)
stat_card(s, 0.55, 2.95, 2.95, 1.5, "ROC-AUC (실측)", "0.9254", "검증 1,379샷 · Bootstrap 95% CI")
stat_card(s, 3.68, 2.95, 2.95, 1.5, "F1-SCORE", "0.7324", "Precision 0.8125 · Recall 0.6667")
stat_card(s, 6.81, 2.95, 2.95, 1.5, "거짓경보 (4-AI 합의)", "31 → 5건", "정밀도 0.50 → 0.84 · ≥3/4 합의", vcolor=RED)
stat_card(s, 9.94, 2.95, 2.85, 1.5, "검증 데이터", "1,379샷", "24센서 · 정상 1,340 + 불량 39")
card(s, 0.55, 4.75, 12.24, 1.95)
text(s, 0.75, 4.95, 4, 0.3, "TECH STACK", 11, RED, bold=True)
text(s, 0.75, 5.3, 11.9, 1.35, [
    [("Frontend   ", 11.5, GRAY, True), ("Next.js (App Router) · TypeScript · Vercel 배포", 11.5, BODY, False)],
    [("Backend    ", 11.5, GRAY, True), ("FastAPI · PyTorch Autoencoder (CPU 추론) · Render 배포 · 매 요청 실시간 모델 forward", 11.5, BODY, False)],
    [("ML / XAI   ", 11.5, GRAY, True), ("Autoencoder + Isolation Forest + One-Class SVM + Local Outlier Factor (4-AI 합의) · GradientSHAP · 비용가중 임계값", 11.5, BODY, False)],
    [("LLM        ", 11.5, GRAY, True), ("Claude Haiku 자연어 진단 보고서 (작업자/반장/부서장 톤 전환 · 키 부재 시 템플릿 폴백)", 11.5, BODY, False)],
], line_spacing=1.3)
footer(s, 2)
# ──────────────────────── 3. 문제 정의 ────────────────────────
s = slide(); accent_bar(s)
text(s, 0.55, 0.78, 12, 0.5, "문제 정의 — 사출성형 현장의 3대 문제", 26, WHITE, bold=True)
probs = [
    ("01", "불량은 사후에 발견된다", "불량률 1.03%의 극심한 불균형 — 라벨이 부족해 일반 지도학습이 불가능하고,\n작업자 육안 검사로는 불량을 다량 생산한 뒤에야 적발된다.", RED),
    ("02", "단일 모델은 알람 피로를 부른다", "단일 모델의 느슨한 판정은 거짓경보가 많아(합집합 기준 FP 31건 · 정밀도 0.50)\n작업자가 경보 자체를 무시하게 된다 — 현장 신뢰의 붕괴.", AMBER),
    ("03", "이상은 알아도 원인을 모른다", "\"이상입니다\"만으로는 조치 불가 — 24개 센서 중 어느 센서가, 얼마나,\n왜 문제인지 설명이 없어 대응이 지연된다.", CYAN),
]
for i, (num, title, desc, c) in enumerate(probs):
    y = 1.6 + i * 1.62
    card(s, 0.55, y, 12.24, 1.42)
    text(s, 0.85, y + 0.18, 0.9, 0.8, num, 30, c, bold=True)
    text(s, 1.95, y + 0.18, 10.5, 0.35, title, 15.5, WHITE, bold=True)
    text(s, 1.95, y + 0.58, 10.5, 0.75, desc, 11.5, BODY, line_spacing=1.18)
text(s, 0.55, 6.55, 12.2, 0.4, [("→  해법: ", 13, WHITE, True),
    ("정상만 학습하는 준지도 학습  +  성격이 다른 4개 AI의 합의  +  SHAP 설명가능성", 13, CYAN, True)])
footer(s, 3)
# ──────────────────────── 4. 아키텍처 ────────────────────────
s = slide(); accent_bar(s)
text(s, 0.55, 0.78, 12, 0.5, "시스템 아키텍처 — 단일 엔진 → 4축 분기", 26, WHITE, bold=True)
flow = ["24 센서\n실시간 수집", "4-AI 이상탐지\nAutoencoder · Isolation Forest\nOne-Class SVM · Local Outlier Factor",
        "SHAP\n원인분석", "처방 카드\nWhat-if 시뮬", "4축 운영\n의사결정"]
fx, fw, gap = 0.55, 2.25, 0.27
for i, label in enumerate(flow):
    x = fx + i * (fw + gap)
    c = card(s, x, 1.6, fw, 1.2, fill=CARD, line=RED if i == 1 else CARD_LINE)
    lines = label.split("\n")
    paras = [[(lines[0], 13, WHITE, True)]] + [[(ln, 8, GRAY, False)] for ln in lines[1:]]
    text(s, x + 0.04, 1.6, fw - 0.08, 1.2, paras,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.12)
    if i < 4:
        text(s, x + fw - 0.04, 2.0, gap + 0.1, 0.4, "→", 16, RED, bold=True)
axes = [("품질 Quality", "4-AI 합의 이상탐지\nSHAP 원인 + 자동 처방", RED),
        ("설비 Equipment", "누적 이상 외삽 → RUL\n예지정비 (3단계 임계)", CYAN),
        ("안전 Safety", "센서 이상 → 과열·과압·기계\n안전위험 자동 변환", AMBER),
        ("생산 Production", "OEE(가동률×성능×양품률)\n불량 Pareto 분석", RGBColor(0x4A, 0xDE, 0x80))]
for i, (t, d, c) in enumerate(axes):
    x = 0.55 + i * 3.16
    card(s, x, 3.05, 2.95, 1.25)
    text(s, x + 0.18, 3.2, 2.6, 0.32, t, 13.5, c, bold=True)
    text(s, x + 0.18, 3.55, 2.6, 0.65, d, 10, BODY, line_spacing=1.15)
card(s, 0.55, 4.62, 12.24, 2.1)
text(s, 0.75, 4.78, 6, 0.3, "DEPLOYMENT — 실배포 운영 구조", 11, RED, bold=True)
text(s, 0.75, 5.14, 11.9, 1.5, [
    [("Vercel (Next.js 프론트)  ──  REST  ──  Render (FastAPI + PyTorch CPU)", 12.5, WHITE, True)],
    [("· 시나리오/스트림 입력마다 백엔드가 ", 11.5, BODY, False), ("매 요청 실시간 모델 추론", 11.5, CYAN, True),
     (" — 미리 만든 화면이 아닌 라이브 시스템", 11.5, BODY, False)],
    [("· 모델·스케일러·임계값은 빌드 시 git에서 로드 — 외부 DB 의존성 없음 · /api/health 헬스체크", 11.5, BODY, False)],
    [("· KAMP 실측 스트림 재생 + 데모 시나리오(정상/경고/위험/긴급) + Claude NLG 보고서", 11.5, BODY, False)],
], line_spacing=1.3)
footer(s, 4)
# ──────────────────────── 5~12. 기능 (캡처 중심) ────────────────────────
feature_slide(5, "TAB 01 · QUALITY", "실시간 진단 — 4-AI 합의 판정", [
    ("4-AI 합의 미터", "Autoencoder · Isolation Forest · One-Class SVM · Local Outlier Factor 각 모델의 FIRE/HOLD를 공개하고 ≥3/4 합의 시에만 이상 판정 — 거짓경보 31→5건."),
    ("강도 등급 + 복원오차 게이지", "Autoencoder 복원오차/임계값(τ 0.320) 비율로 정상·경고·위험·긴급 4단계 등급화."),
    ("처방 카드 TOP 3", "|σ| 상위 센서 기반 즉시/5분 내/관찰 단계별 조치 자동 생성."),
    ("Claude 자연어 보고서", "원인·대처·리소스를 작업자 눈높이로 설명 (LLM 폴백 내장)."),
], "03_tab1_danger_top.png")

feature_slide(6, "TAB 01 · LIVE STREAM", "실측 KAMP 스트림 — 라이브 추론", [
    ("실측 데이터 재생", "KAMP 검증 1,379샷을 실시간 스트리밍 — 매 샷마다 백엔드 모델이 실제 forward 추론."),
    ("LIVE 불량률 집계", "스트리밍 중 발생 이상을 실시간 누적 집계 (캡처: 13샷 중 불량률 15.4%)."),
    ("센서 그리드 동기화", "24센서 × 5계열(시간·위치·속도·압력·온도) σ 변화를 즉시 반영."),
    ("이상 감지 이력", "세션 단위 이상 이벤트 타임라인 자동 기록."),
], "05_tab1_live_stream.png")

feature_slide(7, "TAB 02 · XAI", "불량 원인 분석 — GradientSHAP + 인과 그래프", [
    ("SHAP TOP-5 원인 센서", "GradientExplainer < 100ms — 어느 센서가 복원오차에 얼마나 기여했는지 정량화."),
    ("SHAP Waterfall", "기준 0.048 → 예측 0.639까지 센서별 기여 누적 분해."),
    ("인과 의존성 그래프", "상관 |r|≥0.4 센서를 연결해 상류 공정 원인까지 추적 (충전시간 r=0.95 등)."),
    ("PCA 클러스터", "24센서 → 2축 압축, 정상 1,340 vs 이상 39 분포 시각화."),
], "06_tab2_shap_top.png")

feature_slide(8, "TAB 03 · BATCH", "전체 이력 분석 — 임계값 시뮬레이터", [
    ("검증 1,379샷 일괄 분석", "TP 26 · FN 13 · FP 6 — 공식 혼동행렬을 화면에서 그대로 재현."),
    ("τ 민감도 라이브 시뮬레이션", "슬라이더로 임계값을 바꾸면 혼동행렬·Precision·Recall·F1이 실시간 재계산."),
    ("미탐 ↔ 거짓경보 트레이드오프", "운영자가 비용 관점에서 직접 균형점을 탐색 가능."),
    ("이상 샷 랭킹", "복원오차 상위 샷(#0455 505.92×τ 등)과 주원인 센서를 표로 제공."),
], "07_tab3_batch_top.png")

feature_slide(9, "TAB 04 · EQUIPMENT", "설비 예지정비 — 정비 우선순위 / 모델 커버리지", [
    ("계통별 불량 × 탐지 커버리지", "실측 불량 39건을 계통별 분해 — 속도·압력 100%, 시간 80%, 온도 0%."),
    ("정비 우선순위 자동 산출", "불량 발생량 내림차순으로 계통별 권장 조치 제시 (#01 온도 → 냉각수·가열대 점검)."),
    ("모델 사각지대 정직 공개", "온도 계통 탐지율 0%를 숨기지 않고 전용 룰/SPC 보강을 권장 — AI 개선안 연동."),
    ("RUL 예지정비 (외삽 데모 명시)", "시계열 가동로그 부재로 선형 외삽 데모임을 화면에 명시 — MES 연동 시 실측 전환."),
], "08_tab4_rul_top.png")

feature_slide(10, "TAB 05 · SAFETY", "안전 위험 모니터링 — ISO 12100 위험성 평가", [
    ("센서 이상 → 안전 위험 자동 변환", "과열(온도계열)·과압(압력계열)·기계(속도계열) 위험으로 σ 기반 변환."),
    ("위험성 평가 매트릭스", "발생가능성(실측 불량빈도) × 심각도(현재 σ) 2차원 평가 — ISO 12100 준용."),
    ("안전 조치 체크리스트 자동 생성", "위험 등급에 따라 E-STOP 확인 등 필수 조치를 자동 발급."),
    ("안전센서 정상률 모니터", "15/16 센서 |σ|<2 → 정상률 94% 실시간 표시."),
], "09_tab5_safety_top.png")

feature_slide(11, "TAB 06 · PRODUCTION", "생산 현황 — OEE 종합 설비효율", [
    ("OEE = A × P × Q 분해", "가동률 94.2% × 성능 91.5% × 양품률 97.17% = OEE 83.8% — 양품률만 AI 실측."),
    ("실측·가정 구분 표기", "MES 미연동 항목(가동률·성능)은 '가정' 배지로 정직하게 구분."),
    ("생산 품질 분포", "검증 1,379샷을 12구간 분해 — 양품 1,340 + 불량 39(실측)."),
    ("불량 원인 Pareto", "금형온도4(11건) → 최대 사출속도(10) → 최대 배압(8) → 충전시간(8) — 상위 4개 누적 95%."),
], "10_tab6_oee_top.png")

feature_slide(12, "TAB 07 · MODEL TRUST", "모델 신뢰도 — 학술 검증 · 정직 공개", [
    ("실측 5대 지표 + 95% CI", "ROC-AUC 0.9254 (CI 0.879~0.966) · PR-AUC 0.7080 · F1 0.7324 — Bootstrap n=1000."),
    ("ROC/PR 곡선 + 혼동행렬", "TN 1,334 · FP 6 · FN 13 · TP 26 (검증 1,379샷) 전부 화면 공개."),
    ("합의 알고리즘 9방식 비교", "Stacking(LOOCV F1 0.691)이 Hard Voting(0.761)보다 낮아 채택하지 않은 것까지 공개."),
    ("비용 민감 임계값", "미탐 50만/거짓경보 3만 가중 — 배포 τ 0.3198 vs 비용·F1 최적 τ 0.5171 비교."),
], "11_tab7_trust_top.png")
# ──────────────────────── 13. 정직한 공개 ────────────────────────
s = slide(); accent_bar(s)
text(s, 0.55, 0.78, 12, 0.5, "정직한 공개 — 한계를 숨기지 않는다", 26, WHITE, bold=True)
text(s, 0.55, 1.42, 12.2, 0.4, "포트폴리오에서 가장 강조하고 싶은 가치 — 모든 수치는 실측이며, 잘 안 되는 부분까지 화면에 명시했습니다.", 13, BODY)
lims = [
    ("온도 계통 사각지대 — 탐지율 0%", "온도 계통 불량 11건 중 탐지 0건. 4개 모델 전부 둔감 — 합의 임계를 풀어도 못 잡음.\n→ 화면에 경고 배너로 공개하고 전용 룰/SPC 관리도 보강을 권장.", RED),
    ("Recall 0.6667 — 미탐 13건", "불량 39건 중 26건 탐지. 비용가중 임계값으로 운영 최적화하고,\n합집합 기준 Recall 0.79까지 확보 가능함을 트레이드오프와 함께 제시.", AMBER),
    ("가정은 '가정'이라고 표기", "OEE 가동률·성능(MES 미연동), RUL 선형 외삽(시계열 로그 부재)은\n'가정'·'데모' 배지로 명시 — 실데이터 연동 시 즉시 실측 전환 가능한 구조.", CYAN),
    ("Stacking 미채택 사유 공개", "Meta-learner LOOCV F1 0.69 < Hard Voting 0.761 — 화려한 기법보다\n검증 성능이 높은 단순한 방법을 선택했고, 비교표를 그대로 공개.", RGBColor(0x4A, 0xDE, 0x80)),
]
for i, (t, d, c) in enumerate(lims):
    x = 0.55 + (i % 2) * 6.27; y = 2.1 + (i // 2) * 2.3
    card(s, x, y, 5.97, 2.1)
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y + 0.15), Inches(0.05), Inches(1.8))
    bar.fill.solid(); bar.fill.fore_color.rgb = c; bar.line.fill.background(); bar.shadow.inherit = False
    text(s, x + 0.25, y + 0.2, 5.5, 0.35, t, 14.5, WHITE, bold=True)
    text(s, x + 0.25, y + 0.65, 5.5, 1.3, d, 11, BODY, line_spacing=1.22)
footer(s, 13)
# ──────────────────────── 14. 마무리 ────────────────────────
s = slide(); accent_bar(s)
text(s, 0.55, 0.78, 12, 0.5, "정리 — 5가지 차별점", 26, WHITE, bold=True)
diffs = [
    ("4-AI 합의", "거짓경보 31→5건\n정밀도 0.50→0.84"),
    ("풀스택 XAI", "SHAP + What-if\n+ LLM 자연어 보고서"),
    ("단일 엔진 4축 통합", "품질·설비·안전·생산\n하나의 AI로 운영"),
    ("실측 · 정직 공개", "전 수치 실측 + 95% CI\n사각지대까지 공개"),
    ("실배포", "Vercel + Render\n누구나 접속 가능"),
]
for i, (t, d) in enumerate(diffs):
    x = 0.55 + i * 2.51
    card(s, x, 1.55, 2.31, 1.55)
    text(s, x + 0.15, 1.72, 2.0, 0.6, t, 13.5, CYAN, bold=True, line_spacing=1.05)
    text(s, x + 0.15, 2.35, 2.0, 0.65, d, 10, BODY, line_spacing=1.18)
card(s, 0.55, 3.45, 12.24, 1.55)
text(s, 0.8, 3.65, 11.7, 1.2, [
    [("Live Demo   ", 13, GRAY, True), ("https://" + LIVE_URL, 13, CYAN, True),
     ("      (Render 무료 플랜 — 첫 접속 시 콜드스타트 30~50초)", 10.5, GRAY, False)],
    [("GitHub      ", 13, GRAY, True), ("https://" + GH_URL, 13, CYAN, True)],
    [("Dataset     ", 13, GRAY, True), ("KAMP(인공지능 중소벤처 제조 플랫폼) 사출성형기 AI 데이터셋 — kamp-ai.kr", 12, BODY, False)],
], line_spacing=1.45)
text(s, 0.55, 5.45, 12.2, 0.9, [
    [("\"하나의 AI 엔진으로 공장 운영 전체를 지능화한다\"", 20, WHITE, True)],
    [("— SmartFactory XAI", 13, GRAY, False)],
], align=PP_ALIGN.CENTER, line_spacing=1.4)
footer(s, 14)

prs.save(OUT)
print("SAVED:", OUT)
print("slides:", len(prs.slides.__iter__.__self__._sldIdLst))
