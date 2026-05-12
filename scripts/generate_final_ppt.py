"""
SmartFactory XAI — 예선 기획서 최종 PPT 생성기
- 베이스: output/_template_v2.pptx (대회 공식 양식 그대로)
- 안내 박스 [작성방법] 제거
- 본문 영역에 우리 콘텐츠 추가 (스코어카드, 표, 다이어그램, 스크린샷 placeholder)
- 분량 많은 슬라이드는 양식 슬라이드 복제하여 자동 분할
- 헤더·푸터·타이틀 위치 매 슬라이드 동일 (양식 그대로)
"""
import sys, os, copy
sys.path.insert(0, os.path.dirname(os.path.dirname(os.path.abspath(__file__))))

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_SHAPE_TYPE
from pptx.oxml.ns import qn
from copy import deepcopy
from io import BytesIO

BASE = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
TEMPLATE = os.path.join(BASE, 'output', '_template_v2.pptx')
OUT      = os.path.join(BASE, 'output', '스마트공장XAI_예선기획서_최종.pptx')

# ── 디자인 토큰 (양식에서 추출) ──
FONT_MAIN  = "Pretendard"
C_BODY     = RGBColor(0x20, 0x21, 0x24)  # 본문 검정
C_ACCENT   = RGBColor(0x3C, 0x7C, 0xDE)  # 강조 파랑 (양식)
C_SUB      = RGBColor(0x5F, 0x66, 0x68)  # 서브 다크그레이 (가독성 ↑)
C_RED      = RGBColor(0xD4, 0x21, 0x21)  # 위험 빨강 (특수 강조만)
C_LIGHT    = RGBColor(0xF5, 0xF8, 0xFD)  # 카드 배경 (연파랑)
C_BORDER   = RGBColor(0xDD, 0xE3, 0xEE)  # 카드 테두리
C_GREEN    = RGBColor(0x2E, 0x7D, 0x32)  # 긍정 (정말 필요시)

# 본문 영역 (양식 기준)
BODY_X = Inches(0.62)
BODY_Y = Inches(1.85)
BODY_W = Inches(12.18)
BODY_H = Inches(5.10)  # ~6.95까지

# ─────────────────────────────────────────
# 유틸 함수
# ─────────────────────────────────────────

def remove_guide_box(slide):
    """안내 박스 [작성방법] 텍스트 박스만 제거 (헤더·푸터·타이틀 유지)"""
    shapes_to_remove = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            txt = shape.text_frame.text
            if '[작성방법]' in txt or '참고용 안내' in txt:
                shapes_to_remove.append(shape)
        # 빈 본문 placeholder 박스 (0.62, 1.11) 12.18x5.35도 제거
        elif not shape.has_text_frame:
            pass
    for shape in shapes_to_remove:
        sp = shape._element
        sp.getparent().remove(sp)
    # 빈 본문 영역 박스 제거 (0.62, 1.11) 12.18x5.35
    for shape in list(slide.shapes):
        try:
            if (shape.has_text_frame and not shape.text_frame.text.strip()
                and abs(shape.left - Inches(0.62).emu) < 100
                and abs(shape.top - Inches(1.11).emu) < 100):
                sp = shape._element
                sp.getparent().remove(sp)
        except Exception:
            pass


def duplicate_slide(target_prs, source_prs, source_idx):
    """깨끗한 source_prs[source_idx]를 target_prs의 끝에 복제.
    - 그림(로고)은 add_picture로 재생성하여 관계(rId)를 새 슬라이드에 정확히 연결.
    - 표·테이블·텍스트 등은 deepcopy로 XML 복제.
    """
    src = source_prs.slides[source_idx]
    # 동일 레이아웃 인덱스 사용
    layout = target_prs.slide_layouts[0]
    try:
        layout_idx = list(source_prs.slide_layouts).index(src.slide_layout)
        layout = target_prs.slide_layouts[layout_idx]
    except (ValueError, IndexError):
        pass
    new_slide = target_prs.slides.add_slide(layout)
    # 새 슬라이드의 기본 placeholder 제거
    for shape in list(new_slide.shapes):
        sp = shape._element
        sp.getparent().remove(sp)
    # 원본 shape 복사 (PICTURE는 BytesIO로 재삽입하여 rId 재생성)
    for shape in src.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            try:
                blob = shape.image.blob
                new_slide.shapes.add_picture(
                    BytesIO(blob), shape.left, shape.top,
                    width=shape.width, height=shape.height
                )
            except Exception:
                # 폴백: deepcopy (관계 깨질 수 있음)
                el = shape._element
                new_el = deepcopy(el)
                new_slide.shapes._spTree.insert_element_before(new_el, 'p:extLst')
        else:
            el = shape._element
            new_el = deepcopy(el)
            new_slide.shapes._spTree.insert_element_before(new_el, 'p:extLst')
    return new_slide


def reorder_slides(prs, slides_in_new_order):
    """슬라이드를 지정한 순서대로 재배열.
    slides_in_new_order: 원하는 최종 순서의 슬라이드 객체 list.
    """
    xml_lst = prs.slides._sldIdLst  # <p:sldIdLst>
    # 제거 전에 모든 slide_id를 먼저 캡처
    target_ids = [slide.slide_id for slide in slides_in_new_order]
    entries = list(xml_lst)
    by_id = {int(e.get('id')): e for e in entries}
    # 모두 제거
    for e in entries:
        xml_lst.remove(e)
    # 캡처해둔 id 순서로 다시 추가
    for sid in target_ids:
        if sid in by_id:
            xml_lst.append(by_id[sid])


def add_text(slide, text, left, top, width, height,
             size=12, bold=False, color=None, italic=False,
             align=PP_ALIGN.LEFT, anchor=None, font=FONT_MAIN):
    """텍스트 박스 추가"""
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(2); tf.margin_right = Pt(2)
    tf.margin_top = Pt(1); tf.margin_bottom = Pt(1)
    if anchor is not None:
        tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    if color is not None:
        run.font.color.rgb = color
    else:
        run.font.color.rgb = C_BODY
    return tb


def add_text_lines(slide, lines, left, top, width, height,
                   line_spacing=1.15, font=FONT_MAIN):
    """여러 줄 텍스트 (각 줄별 size/bold/color 지정 가능)"""
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(2); tf.margin_right = Pt(2)
    first = True
    for line in lines:
        if len(line) == 4:
            txt, sz, bold, col = line
            italic = False
        else:
            txt, sz, bold, col, italic = line
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.line_spacing = line_spacing
        run = p.add_run()
        run.text = txt
        run.font.name = font
        run.font.size = Pt(sz)
        run.font.bold = bold
        run.font.italic = italic
        if col is not None:
            run.font.color.rgb = col
    return tb


def add_box(slide, left, top, width, height, fill=None, line=None, line_w=0.5):
    """단순 박스"""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    if fill is not None:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line is not None:
        shape.line.color.rgb = line
        shape.line.width = Pt(line_w)
    else:
        shape.line.fill.background()
    return shape


def add_dashed_box(slide, left, top, width, height, text):
    """점선 박스 (스크린샷 placeholder)"""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0xFA, 0xFB, 0xFC)
    shape.line.color.rgb = C_SUB
    shape.line.width = Pt(1)
    # 점선 (dashStyle dash)
    line_elem = shape.line._get_or_add_ln()
    prstDash = line_elem.find(qn('a:prstDash'))
    if prstDash is None:
        from lxml import etree
        prstDash = etree.SubElement(line_elem, qn('a:prstDash'))
    prstDash.set('val', 'dash')
    # 가운데 텍스트
    add_text(slide, text, left, top + (height - Inches(0.3)) / 2,
             width, Inches(0.3),
             size=11, bold=False, color=C_SUB,
             align=PP_ALIGN.CENTER, italic=True)
    return shape


def add_card(slide, left, top, width, height, title, value, sub="", accent=C_ACCENT,
             value_size=None):
    """스코어카드 (KPI 카드) — title/value/sub 3슬롯, 카드 높이에 맞춰 자동 배치."""
    # 배경
    add_box(slide, left, top, width, height, fill=C_LIGHT, line=C_BORDER, line_w=0.5)
    # 좌측 강조 라인
    add_box(slide, left, top, Pt(3), height, fill=accent)

    # 슬롯 분할 — title(상단 고정), sub(하단 고정), value(나머지)
    inches_w = width / 914400
    inches_h = height / 914400
    title_h = 0.22
    sub_h   = 0.22 if sub else 0.0
    pad_top = 0.05
    pad_bot = 0.05
    value_h = inches_h - title_h - sub_h - pad_top - pad_bot
    if value_h < 0.22:
        value_h = 0.22

    if value_size is None:
        if inches_w < 1.6 or value_h < 0.36:
            value_size = 12
        elif inches_w < 2.2 or value_h < 0.50:
            value_size = 15
        else:
            value_size = 18

    # 라벨 (상단)
    add_text(slide, title, left + Inches(0.12), top + Inches(pad_top),
             width - Inches(0.2), Inches(title_h),
             size=9.5, bold=True, color=C_SUB,
             anchor=MSO_ANCHOR.MIDDLE)
    # 값 (중앙)
    add_text(slide, value, left + Inches(0.12), top + Inches(pad_top + title_h),
             width - Inches(0.2), Inches(value_h),
             size=value_size, bold=True, color=accent,
             anchor=MSO_ANCHOR.MIDDLE)
    # 부제 (하단)
    if sub:
        add_text(slide, sub, left + Inches(0.12),
                 top + Inches(pad_top + title_h + value_h),
                 width - Inches(0.2), Inches(sub_h),
                 size=8.5, color=C_BODY,
                 anchor=MSO_ANCHOR.MIDDLE)


def add_section(slide, left, top, width, title, accent=C_ACCENT):
    """본문 섹션 헤더 (◼ 제목)"""
    h = Inches(0.32)
    # 좌측 파랑 마커
    add_box(slide, left, top + Inches(0.08), Pt(8), Inches(0.18), fill=accent)
    add_text(slide, title, left + Inches(0.18), top,
             width - Inches(0.2), h,
             size=12, bold=True, color=C_ACCENT)
    return top + h


def add_table_simple(slide, rows, left, top, width, col_widths_ratio,
                     header=True, font_size=9.5):
    """간단한 표 (rows: list of lists)"""
    n_rows = len(rows)
    n_cols = len(rows[0])
    total_w = width
    col_widths = [int(total_w * r) for r in col_widths_ratio]
    row_h = Inches(0.32)
    for ri, row in enumerate(rows):
        x = left
        for ci, cell in enumerate(row):
            cw = col_widths[ci]
            is_header = header and ri == 0
            fill = C_LIGHT if is_header else None
            add_box(slide, x, top + ri * row_h, cw, row_h,
                    fill=fill, line=C_BORDER, line_w=0.4)
            add_text(slide, str(cell), x + Inches(0.05), top + ri * row_h + Inches(0.04),
                     cw - Inches(0.08), row_h - Inches(0.08),
                     size=font_size,
                     bold=is_header,
                     color=C_ACCENT if is_header else C_BODY,
                     align=PP_ALIGN.LEFT if ci == 0 else PP_ALIGN.CENTER,
                     anchor=MSO_ANCHOR.MIDDLE)
            x += cw
    return top + n_rows * row_h


# ═══════════════════════════════════════════════════════════════
# 슬라이드별 콘텐츠 작성
# ═══════════════════════════════════════════════════════════════

def fill_slide1_cover(prs):
    """Slide 1 — 표지 (양식 그대로 유지)"""
    slide = prs.slides[0]
    # 양식에 이미 메타정보 있음. 우리 프로젝트명만 업데이트
    for shape in slide.shapes:
        if shape.has_text_frame:
            txt = shape.text_frame.text
            if '2026 스마트 공장 운영 시스템 MVP 개발 해커톤' in txt and len(txt) < 100:
                pass  # 유지
            elif '예선 기획서' in txt and len(txt) < 30:
                pass  # 유지
        # 표 (Table) 안에 팀명/프로젝트명 들어가야 함 — 양식 테이블 활용
    # 양식 테이블 셀에 프로젝트명 입력
    for shape in slide.shapes:
        if shape.shape_type == 19:  # TABLE
            tbl = shape.table
            # 첫 row: 팀명 / 두 번째 row: 프로젝트명
            if len(tbl.rows) >= 2:
                # 프로젝트명 셀에 입력
                try:
                    cell = tbl.cell(1, 1) if len(tbl.columns) > 1 else tbl.cell(1, 0)
                    cell.text = ""
                    p = cell.text_frame.paragraphs[0]
                    r = p.add_run()
                    r.text = "SmartFactory XAI — 4-AI 합의 기반 사출성형 이상탐지·진단·처방 플랫폼"
                    r.font.name = FONT_MAIN
                    r.font.size = Pt(11)
                    r.font.bold = True
                    r.font.color.rgb = C_ACCENT
                except Exception:
                    pass


def fill_slide2_problem(prs):
    """Slide 2 — 문제 정의"""
    slide = prs.slides[1]
    remove_guide_box(slide)

    y = BODY_Y

    # ── 1. 해결하고자 하는 제조 현장 문제 (좌측) ──
    y_sec = add_section(slide, BODY_X, y, Inches(5.8), "1. 해결하고자 하는 제조 현장 문제")
    problems = [
        ("불량 발생 인지 지연", "작업자 육안 검사 한계 → 수백 개 불량 후 인지"),
        ("원인 분석 1~2시간 소요", "24개 센서 동시 모니터링 불가, 재발 방지 어려움"),
        ("사후 대응 반복 사이클", "불량 → 라인 정지 → 수동 검사 → 재가동 무한 반복"),
        ("IT 전담인력 부재", "중소 사출성형 공장 50인 규모, AI 도입 진입장벽"),
    ]
    for i, (title, desc) in enumerate(problems):
        py = y_sec + Inches(0.05) + i * Inches(0.62)
        # 번호 박스
        add_box(slide, BODY_X, py, Inches(0.32), Inches(0.32), fill=C_ACCENT)
        add_text(slide, str(i + 1), BODY_X, py, Inches(0.32), Inches(0.32),
                 size=12, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF),
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(slide, title, BODY_X + Inches(0.42), py - Inches(0.01),
                 Inches(5.4), Inches(0.22), size=11.5, bold=True, color=C_BODY)
        add_text(slide, desc, BODY_X + Inches(0.42), py + Inches(0.20),
                 Inches(5.4), Inches(0.36), size=9.5, color=C_SUB)

    # ── 2. 페르소나·적용 범위 (우측) ──
    rx = BODY_X + Inches(6.0)
    rw = Inches(6.18)
    add_section(slide, rx, y, rw, "2. 적용 대상 페르소나 · 적용 범위")
    persona_box_y = y + Inches(0.42)
    add_box(slide, rx, persona_box_y, rw, Inches(1.10), fill=C_LIGHT, line=C_BORDER, line_w=0.5)
    add_text_lines(slide, [
        ("현장 페르소나 [가정]", 9, True, C_ACCENT),
        ("• 사출성형 작업자 5년차 · 시급 12,000원 · 3교대 근무", 8.5, False, C_BODY),
        ("• 일일 생산 약 6,000샷 · 월 평균 불량 25건 인지 (지연 2~4시간)", 8.5, False, C_BODY),
    ], rx + Inches(0.15), persona_box_y + Inches(0.08), rw - Inches(0.2), Inches(0.95),
        line_spacing=1.18)

    target_box_y = persona_box_y + Inches(1.22)
    add_box(slide, rx, target_box_y, rw, Inches(1.30), fill=C_LIGHT, line=C_BORDER, line_w=0.5)
    add_text_lines(slide, [
        ("적용 대상 · 적용 영역", 9, True, C_ACCENT),
        ("• 중소 사출성형 공장 (50인 규모, IT 전담인력 부재)", 8.5, False, C_BODY),
        ("• 정부 스마트공장 보급사업 2025~2027, 3만개사 대상", 8.5, False, C_BODY),
        ("• 사출성형기 24센서 (시간 5·위치 3·속도 3·압력 5·온도 8) 실시간 모니터링", 8.5, False, C_BODY),
        ("• 이상 탐지 + SHAP 원인 진단 + 처방 자동 출력", 8.5, False, C_BODY),
    ], rx + Inches(0.15), target_box_y + Inches(0.08), rw - Inches(0.2), Inches(1.15),
        line_spacing=1.18)

    # ── 3. 문제 중요성 + KPI (하단 전폭) ──
    by = y + Inches(2.95)
    add_section(slide, BODY_X, by, Inches(12.18), "3. 문제 중요성 · 해결 후 기대 변화 (KPI)")

    # 좌측 배경
    bg_y = by + Inches(0.42)
    add_box(slide, BODY_X, bg_y, Inches(5.8), Inches(1.55), fill=C_LIGHT, line=C_BORDER, line_w=0.5)
    add_text_lines(slide, [
        ("문제 중요성 · 시장 배경", 9, True, C_ACCENT),
        ("• 불량 1건 = 약 50만원 손실 (재료+가공+폐기) [가정]", 8.5, False, C_BODY),
        ("• 공정 데이터 0.58% 극심한 불균형 [실측]", 8.5, False, C_BODY),
        ("   → 지도학습 trivial 함정 회피 필요", 8, False, C_SUB),
        ("• 사출성형 표준 벤치마크 부재", 8.5, False, C_BODY),
        ("   → KAMP 공개 데이터 활용 가치 [실측]", 8, False, C_SUB),
        ("• 정부 스마트공장 보급 2025~2027 3만개사 → 진입 적기", 8.5, True, C_ACCENT),
    ], BODY_X + Inches(0.15), bg_y + Inches(0.08), Inches(5.6), Inches(1.4), line_spacing=1.16)

    # KPI 카드 4개 (우측, 2×2)
    kpi_x = BODY_X + Inches(6.0)
    kpi_w = Inches(2.95)
    kpi_h = Inches(0.74)
    kpi_gap = Inches(0.08)
    kpis = [
        ("Recall 0.795", "[실측] 4 AI 합집합", "39 불량 중 31건 탐지"),
        ("연 1억원", "[추정치] 비용 절감", "50만 × 200건/년"),
        ("응답 수 초", "[구현 완료]", "수 시간 → 즉시"),
        ("ROI 4.2배", "[추정치] 회수율", "연 1억 ÷ 월 200만"),
    ]
    for i, (val, label, sub) in enumerate(kpis):
        kx = kpi_x + (i % 2) * (kpi_w + kpi_gap - Inches(0.05))
        ky = bg_y + (i // 2) * (kpi_h + Inches(0.08))
        add_card(slide, kx, ky, kpi_w, kpi_h, label, val, sub, accent=C_ACCENT)


def fill_slide3_solution(prs):
    """Slide 3 — 솔루션 개요"""
    slide = prs.slides[2]
    remove_guide_box(slide)

    y = BODY_Y

    # ── 슬로건 배너 (상단 전폭) ──
    add_box(slide, BODY_X, y, Inches(12.18), Inches(0.55), fill=C_ACCENT)
    add_text(slide, '"라인이 멈추기 전에 AI가 먼저 알린다"',
             BODY_X, y + Inches(0.08), Inches(12.18), Inches(0.4),
             size=16, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF),
             align=PP_ALIGN.CENTER)

    # ── 1. 솔루션 개요 (좌측) ──
    y2 = y + Inches(0.72)
    y_sec = add_section(slide, BODY_X, y2, Inches(5.8), "1. 솔루션 개요 · 핵심 아이디어")
    add_box(slide, BODY_X, y_sec + Inches(0.05), Inches(5.8), Inches(1.92),
            fill=C_LIGHT, line=C_BORDER, line_w=0.5)
    add_text_lines(slide, [
        ("SmartFactory XAI", 11, True, C_ACCENT),
        ("정상 데이터만 학습한 Autoencoder + KernelSHAP 기반", 9, False, C_BODY),
        ("4 AI(AE + IF + OCSVM + LOF) 합의로 신뢰도 강화", 9, False, C_BODY),
        ("", 5, False, C_BODY),
        ("핵심 아이디어 4가지", 9, True, C_ACCENT),
        ("• 반지도학습 → 0.58% 극불균형 trivial 함정 회피", 8.5, False, C_BODY),
        ("• SHAP → 24센서 중 어느 게 이상 원인인지 정량 분해", 8.5, False, C_BODY),
        ("• 다중 AI 합의 → Recall 0.667 → 0.795 [실측]", 8.5, True, C_ACCENT),
        ("• LIVE 디지털 트윈 → 본선 OPC-UA 동치 시연", 8.5, False, C_BODY),
    ], BODY_X + Inches(0.15), y_sec + Inches(0.12), Inches(5.6), Inches(1.8),
        line_spacing=1.18)

    # ── 2. 4단계 자동화 (우측) ──
    rx = BODY_X + Inches(6.0)
    rw = Inches(6.18)
    add_section(slide, rx, y2, rw, "2. 4단계 운영 자동화 + Feedback Loop")
    steps = [
        ("① 탐지", "Autoencoder 24→16→8→16→24 복원오차"),
        ("② 진단", "KernelSHAP K-Means 50 배경 → 24센서 기여도"),
        ("③ 처방", "센서별 24개 처방 카드 · 심각도 3단계"),
        ("④ 추적", "이력 50건 + JSON 영속 + 일일 리포트 MD"),
    ]
    for i, (num, desc) in enumerate(steps):
        sy = y2 + Inches(0.42) + i * Inches(0.40)
        add_box(slide, rx, sy, rw, Inches(0.36), fill=C_LIGHT, line=C_BORDER, line_w=0.4)
        add_text(slide, num, rx + Inches(0.08), sy + Inches(0.04),
                 Inches(0.75), Inches(0.28),
                 size=11, bold=True, color=C_ACCENT, anchor=MSO_ANCHOR.MIDDLE)
        add_text(slide, desc, rx + Inches(0.85), sy + Inches(0.04),
                 Inches(5.3), Inches(0.28),
                 size=9.5, color=C_BODY, anchor=MSO_ANCHOR.MIDDLE)
    # Feedback Loop
    fy = y2 + Inches(0.42) + 4 * Inches(0.40)
    add_box(slide, rx, fy, rw, Inches(0.32), fill=None, line=C_ACCENT, line_w=1.5)
    add_text(slide, "↺ Feedback Loop · ④ 추적 데이터 → ① 모델 재학습 (Active Learning, 30건 라벨 누적)",
             rx + Inches(0.1), fy + Inches(0.03),
             rw - Inches(0.2), Inches(0.26),
             size=9.5, bold=True, color=C_ACCENT, anchor=MSO_ANCHOR.MIDDLE)

    # ── 3. 플랫폼 가치 4 카드 + 스크린샷 (하단) ──
    by = y2 + Inches(2.45)
    add_section(slide, BODY_X, by, Inches(12.18), "3. 플랫폼 제공 가치 + 시연 화면")

    values = [
        ("실시간 탐지", "ROC-AUC 0.9254", "[실측]"),
        ("원인 자동 진단", "SHAP 24센서", "기여도 분해"),
        ("처방 자동화", "24개 처방 카드", "조작/정비 자동"),
        ("이력 표준화", "CSV / MD", "교대 인수인계"),
    ]
    for i, (t, v, s) in enumerate(values):
        cx = BODY_X + i * Inches(1.85)
        cy = by + Inches(0.42)
        add_card(slide, cx, cy, Inches(1.75), Inches(0.85), t, v, s)

    # 스크린샷 박스 (우측) — 현재 대시보드 기준 정확한 위치 명시
    ssx = BODY_X + Inches(7.6)
    ssy = by + Inches(0.42)
    add_dashed_box(slide, ssx, ssy, Inches(4.58), Inches(1.45),
                   "[스크린샷: '실시간 진단' 탭 — '판정 결과' 박스 (위험 DANGER 빨간 배너) + 다중 AI 합의 미터 + 대응 권고 카드]")


def fill_slide4_1_features(prs):
    """Slide 4-1 — 주요 기능 (1/2) MVP 핵심 F1~F5
    레이아웃: 좌측에 F1~F5 (높이 꽉 차게), 우측에 스크린샷 2개 스택 (반반).
    """
    slide = prs.slides[3]
    remove_guide_box(slide)

    y = BODY_Y

    # ── 좌측 컬럼 (BODY_X ~ +6.00) — F1~F5 ──
    LW = Inches(6.00)
    y_sec = add_section(slide, BODY_X, y, LW,
                        "MVP 핵심 5기능 (F1~F5) — 예선 단계 모두 구현 완료")

    features = [
        ("F1", "AI 자동 이상탐지",      "Autoencoder · ROC-AUC 0.9254 · F1 0.7324 [실측]"),
        ("F2", "SHAP 24센서 원인 진단", "KernelSHAP · Top-5 자동 도출 · 실시간 기여도 분해"),
        ("F3", "처방 카드 24개",         "센서별 처방 · 조작 가능 / 정비 필요 자동 분류"),
        ("F4", "심각도 3단계 + 알람",    "경고/위험/긴급 · 작업자 → 반장 → 부서장 에스컬레이션"),
        ("F5", "교대 인수인계 자동화",  "이상 이력 50건 + CSV·MD 다운로드 + JSON 영속"),
    ]
    # 가용 높이 = BODY_H - section header(0.32) - gap(0.10) = 4.68
    # 5칸 × (행 + 간격): 행 0.85, 간격 0.07 → 4.25 + 0.28 = 4.53
    row_h = Inches(0.85)
    row_gap = Inches(0.07)
    cy_start = y_sec + Inches(0.12)
    for i, (fid, name, desc) in enumerate(features):
        cy = cy_start + i * (row_h + row_gap)
        # 기능 ID 박스 (좌측, 0.70 x 0.85)
        add_box(slide, BODY_X, cy, Inches(0.70), row_h, fill=C_ACCENT)
        add_text(slide, fid, BODY_X, cy, Inches(0.70), row_h,
                 size=18, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF),
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        # 본문 박스 (우측, 5.22 x 0.85)
        bx = BODY_X + Inches(0.78)
        bw = LW - Inches(0.78)
        add_box(slide, bx, cy, bw, row_h, fill=C_LIGHT, line=C_BORDER, line_w=0.4)
        add_text(slide, name, bx + Inches(0.16), cy + Inches(0.10),
                 bw - Inches(0.2), Inches(0.30),
                 size=13, bold=True, color=C_BODY)
        add_text(slide, desc, bx + Inches(0.16), cy + Inches(0.42),
                 bw - Inches(0.2), Inches(0.36),
                 size=10, color=C_SUB)

    # ── 우측 컬럼 (BODY_X + 6.18 ~ +12.18) — 스크린샷 2개 스택 ──
    RX = BODY_X + Inches(6.18)
    RW = Inches(6.00)
    # 상단·하단 각 2.25 + 간격 0.18 → 4.68 (section header 없이 BODY_Y부터 시작)
    sh_h = Inches(2.25)
    add_dashed_box(slide, RX, y + Inches(0.05), RW, sh_h,
                   "[스크린샷: '불량 원인 분석' 탭 — 센서별 평균 SHAP 기여도 막대 차트 (Top-5 빨강 강조) + 인과 그래프]")
    add_dashed_box(slide, RX, y + Inches(0.05) + sh_h + Inches(0.18), RW, sh_h,
                   "[스크린샷: '실시간 진단' 탭 — 대응 권고 처방 카드 + 심각도 3단계 (경고/위험/긴급) 배너]")


def fill_slide4_2_features(prs, prs_clean):
    """Slide 4-2 — 주요 기능 (2/2) 차별화 F6~F10 + 본선 P1~P5"""
    # 양식 슬라이드 4를 (깨끗한 원본에서) 복제
    new_slide = duplicate_slide(prs, prs_clean, 3)
    remove_guide_box(new_slide)

    y = BODY_Y

    # ── 좌측 컬럼 (BODY_X ~ +6.00) — F6~F10 ──
    LW = Inches(6.00)
    y_sec = add_section(new_slide, BODY_X, y, LW,
                        "차별화 5기능 (F6~F10) — 다중 AI · LIVE · What-if · NL · Active Learning")

    features = [
        ("F6",  "다중 AI 합의",            "Autoencoder + IF + OCSVM + LOF · Recall 0.667→0.795 [실측]"),
        ("F7",  "LIVE 디지털 트윈",        "검증셋 1,379건 자동 스트리밍 · OPC-UA 동치 시연"),
        ("F8",  "What-if 시뮬레이터",       "Counterfactual · \"얼마나 고쳐야 정상이 되나\" 자동 계산"),
        ("F9",  "AI 자연어 진단 보고서",   "부서장·반장 보고용 Markdown 자동 생성"),
        ("F10", "Active Learning",          "라벨 30건 누적 시 모델 재학습 트리거 (Feedback Loop)"),
    ]
    row_h = Inches(0.85)
    row_gap = Inches(0.07)
    cy_start = y_sec + Inches(0.12)
    for i, (fid, name, desc) in enumerate(features):
        cy = cy_start + i * (row_h + row_gap)
        add_box(new_slide, BODY_X, cy, Inches(0.70), row_h, fill=C_ACCENT)
        add_text(new_slide, fid, BODY_X, cy, Inches(0.70), row_h,
                 size=17, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF),
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        bx = BODY_X + Inches(0.78)
        bw = LW - Inches(0.78)
        add_box(new_slide, bx, cy, bw, row_h, fill=C_LIGHT, line=C_BORDER, line_w=0.4)
        add_text(new_slide, name, bx + Inches(0.16), cy + Inches(0.10),
                 bw - Inches(0.2), Inches(0.30),
                 size=13, bold=True, color=C_BODY)
        add_text(new_slide, desc, bx + Inches(0.16), cy + Inches(0.42),
                 bw - Inches(0.2), Inches(0.36),
                 size=9.5, color=C_SUB)

    # ── 우측 컬럼 (BODY_X + 6.18 ~ +12.18) — 스크린샷 2개 스택 ──
    # P1~P5는 슬라이드 10(본선 구현 계획)으로 이전됨.
    RX = BODY_X + Inches(6.18)
    RW = Inches(6.00)
    sh_h = Inches(2.25)
    add_dashed_box(new_slide, RX, y + Inches(0.05), RW, sh_h,
                   "[스크린샷: '실시간 진단' 탭 — 다중 AI 합의 미터 (Autoencoder · Isolation Forest · One-Class SVM · LOF 4개 모델 동의도 게이지)]")
    add_dashed_box(new_slide, RX, y + Inches(0.05) + sh_h + Inches(0.18), RW, sh_h,
                   "[스크린샷: '실시간 진단' 탭 하단 — What-if 시뮬레이터 + AI 자연어 진단 보고서 + Active Learning 라벨링 표]")


def fill_slide5_1_data(prs):
    """Slide 5-1 — 데이터 및 데이터 처리"""
    slide = prs.slides[4]
    remove_guide_box(slide)

    y = BODY_Y

    # ── 1. 활용 데이터 (좌측) ──
    y_sec = add_section(slide, BODY_X, y, Inches(6.0), "1. 활용 데이터 · KAMP 공공데이터포털 [실측]")
    add_text(slide, "출처: data.go.kr/data/15089213",
             BODY_X + Inches(0.1), y_sec + Inches(0.02),
             Inches(5.5), Inches(0.22), size=9.5, italic=True, color=C_SUB)

    rows = [
        ["데이터셋", "행 수", "용도"],
        ["supervised_label_cn7", "6,736 (정상 6,697 + 불량 39)", "학습 + 검증"],
        ["moldset_labeled_cn7", "1,211 (불량 17)", "외부 검증 (CN7)"],
        ["moldset_labeled_rg3", "1,182 (불량 25)", "외부 검증 (RG3)"],
        ["labeled_data", "7,996", "Scaler raw 학습"],
    ]
    add_table_simple(slide, rows, BODY_X, y_sec + Inches(0.32),
                     Inches(6.0), [0.36, 0.34, 0.30], font_size=9)

    # 불균형 해결 전략 — 반지도학습 메커니즘 상세 설명
    imb_y = y_sec + Inches(2.18)
    add_box(slide, BODY_X, imb_y, Inches(6.0), Inches(1.55),
            fill=C_LIGHT, line=C_RED, line_w=1.2)
    add_text(slide, "⚠ 불량률 0.58% — 지도학습 trivial 함정 → 반지도학습으로 우회",
             BODY_X + Inches(0.15), imb_y + Inches(0.05),
             Inches(5.7), Inches(0.28),
             size=10.5, bold=True, color=C_RED)
    add_text_lines(slide, [
        ("문제 — 지도학습 trivial 함정:", 9, True, C_BODY),
        ("• 불량 0.58% 환경에서 \"항상 정상\"만 예측해도 정확도 99% [실측]", 8.5, False, C_BODY),
        ("• 불량 39건만으로 결정 경계 학습 불가능 → 패턴 일반화 실패", 8.5, False, C_BODY),
        ("해결 — 반지도학습 (Autoencoder):", 9, True, C_ACCENT),
        ("• 정상 5,357건만 학습 → 분포 압축·복원 (24→16→8→16→24)", 8.5, False, C_BODY),
        ("• 학습 못한 입력 = 복원 오차 큼 → 임계값 0.3198 초과 시 '이상' 판정", 8.5, False, C_BODY),
    ], BODY_X + Inches(0.18), imb_y + Inches(0.34),
        Inches(5.65), Inches(1.18), line_spacing=1.08)

    # ── 2. 검증셋 도넛 (우측 상단) ──
    rx = BODY_X + Inches(6.18)
    rw = Inches(6.0)
    add_section(slide, rx, y, rw, "2. 검증셋 구성 (1,379건)")
    # 도넛 대체 → 2 카드
    add_card(slide, rx, y + Inches(0.42), Inches(2.9), Inches(1.05),
             "정상 샘플", "1,340", "97.2% · 검증셋")
    add_card(slide, rx + Inches(3.0), y + Inches(0.42), Inches(3.0), Inches(1.05),
             "불량 샘플", "39", "2.8% · 검증셋 (전체 불량)", accent=C_RED)

    # 표 추가
    val_rows = [
        ["구분", "행 수", "라벨"],
        ["학습", "5,357", "정상 only"],
        ["검증 정상", "1,340", "0"],
        ["검증 불량", "39", "1"],
        ["합계", "6,736", "—"],
    ]
    add_table_simple(slide, val_rows, rx, y + Inches(1.58), rw, [0.34, 0.33, 0.33], font_size=9)

    # ── 3. 데이터 처리 방식 (하단 전폭) ──
    by = y + Inches(4.10)
    add_section(slide, BODY_X, by, Inches(12.18), "3. 데이터 처리 방식 · 6단계 검증 방법론 [실측]")
    methods = [
        ("①", "StandardScaler train fit only", "data leakage 방지"),
        ("②", "80/20 분할", "random_state=42, 재현성"),
        ("③", "Bootstrap 95% CI", "1,000회 · 소표본 통계 신뢰도"),
        ("④", "Pseudo Hold-out", "마지막 20건 분리"),
        ("⑤", "Cross-Machine 외부 검증", "CN7 + RG3 별도 금형"),
        ("⑥", "학습 데이터", "정상 5,357 / 검증 1,379"),
    ]
    cw = Inches(1.95)
    for i, (n, t, d) in enumerate(methods):
        mx = BODY_X + i * (cw + Inches(0.05))
        my = by + Inches(0.38)
        add_box(slide, mx, my, cw, Inches(0.62),
                fill=C_LIGHT, line=C_BORDER, line_w=0.5)
        add_text(slide, n, mx + Inches(0.05), my + Inches(0.02),
                 Inches(0.4), Inches(0.26), size=13, bold=True, color=C_ACCENT,
                 anchor=MSO_ANCHOR.MIDDLE)
        add_text(slide, t, mx + Inches(0.45), my + Inches(0.04),
                 cw - Inches(0.5), Inches(0.24), size=9, bold=True, color=C_BODY)
        add_text(slide, d, mx + Inches(0.45), my + Inches(0.28),
                 cw - Inches(0.5), Inches(0.30), size=8, color=C_SUB)


def fill_slide5_2_data(prs, prs_clean):
    """Slide 5-2 — AI 기술·기술 스택·학술 레퍼런스·차별화"""
    new_slide = duplicate_slide(prs, prs_clean, 4)
    remove_guide_box(new_slide)

    y = BODY_Y

    # ── 1. AI / 분석 기술 (좌측 상단) ──
    y_sec = add_section(new_slide, BODY_X, y, Inches(6.0), "1. 적용 AI · 분석 기술 [실측]")
    add_box(new_slide, BODY_X, y_sec + Inches(0.05), Inches(6.0), Inches(2.30),
            fill=C_LIGHT, line=C_BORDER, line_w=0.5)
    add_text_lines(new_slide, [
        ("Autoencoder (PyTorch 2.3)", 10, True, C_ACCENT),
        ("• 구조: 24→16→8→16→24 (BatchNorm + MSE + Adam, 28 epoch)", 8.5, False, C_BODY),
        ("• 반지도학습 — 정상 데이터만 학습, 불균형 trivial 함정 회피", 8.5, False, C_BODY),
        ("• 임계값: 0.3198 (99th percentile + F1-optimal)", 8.5, False, C_BODY),
        ("", 4, False, C_BODY),
        ("SHAP (shap 0.48)", 10, True, C_ACCENT),
        ("• KernelExplainer + K-Means 50 배경 압축", 8.5, False, C_BODY),
        ("• 24센서 실시간 기여도 분해 · Waterfall 시각화", 8.5, False, C_BODY),
        ("", 4, False, C_BODY),
        ("다중 AI 합의 (Isolation Forest + OCSVM + LOF)", 10, True, C_ACCENT),
        ("• 각 모델 정상 99th percentile 임계값 · 운영 모드별 합의 기준", 8.5, False, C_BODY),
    ], BODY_X + Inches(0.15), y_sec + Inches(0.12), Inches(5.7), Inches(2.18),
        line_spacing=1.18)

    # ── 2. 기술 스택 (우측 상단) ──
    rx = BODY_X + Inches(6.18)
    rw = Inches(6.0)
    add_section(new_slide, rx, y, rw, "2. 사용 기술 스택 · 시스템 구성")
    add_box(new_slide, rx, y + Inches(0.42), rw, Inches(2.0),
            fill=C_LIGHT, line=C_BORDER, line_w=0.5)
    add_text_lines(new_slide, [
        ("ML", 9, True, C_ACCENT),
        ("PyTorch 2.3 · scikit-learn 1.2 · SHAP 0.48", 8.5, False, C_BODY),
        ("UI", 9, True, C_ACCENT),
        ("Streamlit 1.36 · Plotly 5.22 · python-pptx 1.0", 8.5, False, C_BODY),
        ("Infrastructure", 9, True, C_ACCENT),
        ("Python 3.11 (Anaconda) · 온프레미스 배포 (데이터 외부 전송 0)", 8.5, False, C_BODY),
        ("Security", 9, True, C_ACCENT),
        ("모델 SHA-256 무결성 검증 · JSON 이력 영속", 8.5, False, C_BODY),
        ("Finals add-on", 9, True, C_ACCENT),
        ("OPC-UA (asyncua) · MQTT (paho-mqtt) · AWS 옵션", 8.5, False, C_BODY),
    ], rx + Inches(0.15), y + Inches(0.50), rw - Inches(0.2), Inches(1.85), line_spacing=1.16)

    # ── 3. 기술 제약 + 해결 (좌측 하단) ──
    by = y + Inches(2.55)
    add_section(new_slide, BODY_X, by, Inches(6.0), "3. 기술 제약 · 해결 전략")
    rows = [
        ["제약", "해결 [실측]"],
        ["0.58% 극심 불균형", "반지도 Autoencoder (정상만 학습)"],
        ["불량 39 소표본", "Bootstrap CI + Cross-Machine 외부 검증"],
        ["AE 단독 false negative 13건", "다중 AI 합의 → 8건 (-38%)"],
        ["Streamlit 실시간 한계", "LIVE 디지털 트윈 → 본선 OPC-UA 전환"],
    ]
    add_table_simple(new_slide, rows, BODY_X, by + Inches(0.42),
                     Inches(6.0), [0.40, 0.60], font_size=9.5)

    # ── 4. 학술 레퍼런스 + 차별화 (우측 하단) ──
    add_section(new_slide, rx, by, rw, "4. 학술 레퍼런스 7개 · 차별화 4가지")
    refs = [
        "① MDPI Processes 13(3), 912 (2025) — KAMP 직접 비교",
        "② arXiv:2511.08108 (2025) — LSTM+SHAP F1 0.92",
        "③ Brito MAKE 6(1), 16 (2024) — SHAP 베어링 98.5%",
        "④ PhysiCausalNet IEEE TII (2024) — Cross-Machine",
        "⑤ EWAD-IIoT WGAN (Sci. Rep. 2025) — Bootstrap CI",
        "⑥ Survey arXiv:2503.13195 (2025) — Deep AD",
        "⑦ Ketonen IEEE ICPS (2021) — VAE+RNN root-cause",
    ]
    add_box(new_slide, rx, by + Inches(0.42), rw, Inches(1.50),
            fill=C_LIGHT, line=C_BORDER, line_w=0.5)
    add_text_lines(new_slide, [(r, 7.8, False, C_BODY) for r in refs],
                   rx + Inches(0.12), by + Inches(0.48), rw - Inches(0.2), Inches(1.45),
                   line_spacing=1.20)
    # 차별화 배너
    db_y = by + Inches(2.00)
    add_box(new_slide, rx, db_y, rw, Inches(0.5), fill=C_ACCENT)
    add_text(new_slide, "차별화 ① 공개 KAMP  ② 비지도+KernelSHAP  ③ Cross-Machine+CI  ④ 24개 처방 통합",
             rx + Inches(0.1), db_y + Inches(0.10), rw - Inches(0.2), Inches(0.30),
             size=10, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF), align=PP_ALIGN.CENTER,
             anchor=MSO_ANCHOR.MIDDLE)


def fill_slide6_usecase(prs):
    """Slide 6 — 유즈케이스"""
    slide = prs.slides[5]
    remove_guide_box(slide)

    y = BODY_Y

    # ── 1. 3계층 사용자 (상단) ──
    y_sec = add_section(slide, BODY_X, y, Inches(12.18), "1. 주요 사용자 — 공장 운영 3계층")
    users = [
        ("작업자", "사출성형 3교대", "상시", "실시간 진단 탭"),
        ("반장", "라인 관리자", "매 교대", "실시간 진단 + 알람"),
        ("부서장", "생산기술 부서장", "일·월", "생산 이력 + 일일 리포트"),
    ]
    uw = Inches(4.0)
    for i, (role, desc, freq, screen) in enumerate(users):
        ux = BODY_X + i * (uw + Inches(0.09))
        uy = y_sec + Inches(0.08)
        add_box(slide, ux, uy, uw, Inches(0.95), fill=C_LIGHT, line=C_BORDER, line_w=0.5)
        add_box(slide, ux, uy, uw, Inches(0.30), fill=C_ACCENT)
        add_text(slide, role, ux + Inches(0.12), uy + Inches(0.03),
                 uw - Inches(0.2), Inches(0.24), size=12, bold=True,
                 color=RGBColor(0xFF, 0xFF, 0xFF), anchor=MSO_ANCHOR.MIDDLE)
        add_text(slide, freq, ux + uw - Inches(1.2), uy + Inches(0.03),
                 Inches(1.1), Inches(0.24), size=9.5, bold=True,
                 color=RGBColor(0xFF, 0xFF, 0xFF), align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)
        add_text(slide, desc, ux + Inches(0.12), uy + Inches(0.36),
                 uw - Inches(0.2), Inches(0.25), size=10, color=C_BODY)
        add_text(slide, "▸ " + screen, ux + Inches(0.12), uy + Inches(0.62),
                 uw - Inches(0.2), Inches(0.25), size=9.5, italic=True, color=C_ACCENT)

    # ── 2. 사용자 흐름 (중단) ──
    fy = y + Inches(1.50)
    add_section(slide, BODY_X, fy, Inches(12.18), "2. 사용자별 행동 흐름")
    flows = [
        ("작업자", "센서 모니터 → 이상 알람 → SHAP 원인 확인 → 처방 카드 → 조치 → 이력 기록"),
        ("반장",   "교대 시작 → 이력 50건 확인 → 심각도 분포 점검 → CSV 다운로드 → 인수인계"),
        ("부서장", "일일 리포트(MD) 수신 → 이상률 트렌드 → 장비 건강도 → 정비 권고 → 월 KPI"),
    ]
    for i, (role, flow) in enumerate(flows):
        rfy = fy + Inches(0.42) + i * Inches(0.32)
        add_box(slide, BODY_X, rfy, Inches(12.18), Inches(0.28),
                fill=C_LIGHT, line=C_BORDER, line_w=0.4)
        add_text(slide, role, BODY_X + Inches(0.1), rfy + Inches(0.03),
                 Inches(0.8), Inches(0.22), size=10, bold=True, color=C_ACCENT,
                 anchor=MSO_ANCHOR.MIDDLE)
        add_text(slide, flow, BODY_X + Inches(1.0), rfy + Inches(0.03),
                 Inches(11.1), Inches(0.22), size=9.5, color=C_BODY,
                 anchor=MSO_ANCHOR.MIDDLE)

    # ── 3. 대표 시나리오 4종 (하단) ──
    sy = y + Inches(2.92)
    add_section(slide, BODY_X, sy, Inches(12.18), "3. 대표 시나리오 4종 (사이드바 1클릭 자동 적용)")
    # 색상은 2색 원칙(파랑·빨강) — 정상/경고는 파랑, 위험/긴급은 빨강
    scenarios = [
        ("S1 정상", "모든 센서 정상", "NORMAL", "운영 유지", C_ACCENT),
        ("S2 경고", "Back_Pressure +3.5σ", "WARNING", "10분 내 재측정", C_ACCENT),
        ("S3 위험", "Mold_Temp_4 +4.2σ", "DANGER", "담당자 즉시 호출", C_RED),
        ("S4 긴급", "Filling_Time +15σ", "CRITICAL", "라인 정지 검토", C_RED),
    ]
    sw = Inches(2.95)
    for i, (sid, pat, status, action, color) in enumerate(scenarios):
        sx = BODY_X + i * (sw + Inches(0.07))
        ssy = sy + Inches(0.42)
        add_box(slide, sx, ssy, sw, Inches(1.30),
                fill=C_LIGHT, line=color, line_w=1.2)
        add_box(slide, sx, ssy, sw, Inches(0.28), fill=color)
        add_text(slide, sid, sx + Inches(0.1), ssy + Inches(0.02),
                 sw - Inches(0.2), Inches(0.24), size=11, bold=True,
                 color=RGBColor(0xFF, 0xFF, 0xFF), anchor=MSO_ANCHOR.MIDDLE)
        add_text(slide, "입력: " + pat, sx + Inches(0.1), ssy + Inches(0.35),
                 sw - Inches(0.2), Inches(0.24), size=9, bold=True, color=C_BODY)
        add_text(slide, "AI: " + status, sx + Inches(0.1), ssy + Inches(0.62),
                 sw - Inches(0.2), Inches(0.24), size=9, bold=True, color=color)
        add_text(slide, "조치: " + action, sx + Inches(0.1), ssy + Inches(0.90),
                 sw - Inches(0.2), Inches(0.30), size=8.5, color=C_SUB)


def fill_slide7_scope(prs):
    """Slide 7 — MVP 범위
    레이아웃: 좌측에 10개 체크리스트 + 재현성 보장 박스, 우측에 스크린샷 2개 스택.
    P1~P5는 슬라이드 10(본선 구현 계획)으로 이전.
    """
    slide = prs.slides[6]
    remove_guide_box(slide)

    y = BODY_Y
    LW = Inches(6.00)

    # ── 좌측 컬럼 — 10개 체크리스트 (2 × 5 그리드) ──
    y_sec = add_section(slide, BODY_X, y, LW,
                        "1. 시연 가능 핵심 기능 10개 [실측 · 모두 구현 완료]")
    checks = [
        ("AI 자동 이상탐지",       "ROC-AUC 0.9254",    "모델 신뢰도"),
        ("24센서 즉시 판정",        "이상 점수 + 심각도", "실시간 진단"),
        ("데모 시나리오 4종",       "정상/경고/위험/긴급", "사이드바"),
        ("SHAP 24센서 기여도",      "Waterfall + Bar",    "불량 원인 분석"),
        ("처방 카드 24개",           "조작 / 정비 구분",   "실시간 진단"),
        ("심각도 3단계 + 알람",      "에스컬레이션",       "실시간 진단"),
        ("이상 이력 50건 + 다운로드", "CSV / MD",          "실시간 진단"),
        ("Bootstrap 95% CI",         "Cross-Machine 검증", "모델 신뢰도"),
        ("검증셋 1,379건 스코어링",   "시계열 + 트렌드",   "전체 이력 일괄 분석"),
        ("LIVE 디지털 트윈",          "10초 스트리밍",     "사이드바"),
    ]
    cell_w = Inches(2.95)
    cell_gap = Inches(0.10)
    cell_h = Inches(0.42)
    cell_v_gap = Inches(0.07)
    cy_start = y_sec + Inches(0.10)
    for i, (name, sub, loc) in enumerate(checks):
        col = i % 2
        row = i // 2
        cx = BODY_X + col * (cell_w + cell_gap)
        cy = cy_start + row * (cell_h + cell_v_gap)
        add_box(slide, cx, cy, cell_w, cell_h, fill=C_LIGHT, line=C_BORDER, line_w=0.4)
        # ✓ 체크박스
        add_box(slide, cx + Inches(0.06), cy + Inches(0.09), Inches(0.24), Inches(0.24), fill=C_ACCENT)
        add_text(slide, "✓", cx + Inches(0.06), cy + Inches(0.06),
                 Inches(0.24), Inches(0.30),
                 size=13, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF),
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        idx_str = f"{i+1:02d}"
        add_text(slide, idx_str, cx + Inches(0.36), cy + Inches(0.05),
                 Inches(0.25), Inches(0.22), size=9, color=C_SUB)
        add_text(slide, name, cx + Inches(0.66), cy + Inches(0.04),
                 cell_w - Inches(0.75), Inches(0.22),
                 size=10, bold=True, color=C_BODY)
        # 하단 라인: sub | 위치 (한 줄)
        add_text(slide, f"{sub}  ·  {loc}",
                 cx + Inches(0.66), cy + Inches(0.23),
                 cell_w - Inches(0.75), Inches(0.18),
                 size=8.5, color=C_SUB)

    # ── 좌측 하단 — 재현성 보장 박스 ──
    rep_h = Inches(1.35)
    rep_y = cy_start + 5 * (cell_h + cell_v_gap) + Inches(0.10)
    add_box(slide, BODY_X, rep_y, LW, rep_h,
            fill=C_LIGHT, line=C_ACCENT, line_w=1)
    add_text(slide, "재현성 보장 — 평가위원 검증 가능",
             BODY_X + Inches(0.15), rep_y + Inches(0.06),
             LW - Inches(0.2), Inches(0.26),
             size=11, bold=True, color=C_ACCENT)
    add_text_lines(slide, [
        ("• install.bat 원클릭 conda 환경 (Python 3.11 자동 설치)", 9, False, C_BODY),
        ("• REPRODUCE.md 7단계 절차 (Data → Train → Eval → Run)", 9, False, C_BODY),
        ("• verify_model.py SHA-256 모델 무결성 자동 검증", 9, False, C_BODY),
        ("• 학술 레퍼런스 7개 (KAMP·LSTM+SHAP·Cross-Machine 등)", 9, False, C_BODY),
    ], BODY_X + Inches(0.18), rep_y + Inches(0.36),
        LW - Inches(0.2), rep_h - Inches(0.40), line_spacing=1.20)

    # ── 우측 컬럼 — 스크린샷 2개 스택 ──
    RX = BODY_X + Inches(6.18)
    RW = Inches(6.00)
    sh_h = Inches(2.40)
    add_dashed_box(slide, RX, y + Inches(0.05), RW, sh_h,
                   "[스크린샷: '모델 신뢰도 확인' 탭 — ROC Curve + Bootstrap 95% CI + 탐지 결과 요약 (TP/FP/FN)]")
    add_dashed_box(slide, RX, y + Inches(0.05) + sh_h + Inches(0.18), RW, sh_h,
                   "[스크린샷: 사이드바 — LIVE 디지털 트윈 모드 ON 펄스 배너 + 데모 시나리오 4종 (정상/경고/위험/긴급) 1클릭 적용]")


def fill_slide8_1_impact(prs):
    """Slide 8-1 — 기대 효과 (1/2): 정량 + 정성 + 현장 적용"""
    slide = prs.slides[7]
    remove_guide_box(slide)

    y = BODY_Y

    # ── 1. 정량 효과 4 카드 (상단) ──
    y_sec = add_section(slide, BODY_X, y, Inches(12.18), "1. 정량 기대 효과 — 도입 첫해 기준")
    kpis = [
        ("Recall 0.795", "[실측] 4 AI 합집합", "39 불량 중 31건 탐지"),
        ("연 1억원 절감", "[추정치] 비용 회피", "50만원 × 200건/년"),
        ("응답 수 초", "[구현 완료]", "수 시간 → 즉시"),
        ("ROI 4.2배", "[추정치] 투자 회수", "연 1억 ÷ 월 200만 SaaS"),
    ]
    cw = Inches(2.95)
    for i, (val, label, sub) in enumerate(kpis):
        cx = BODY_X + i * (cw + Inches(0.07))
        cy = y_sec + Inches(0.08)
        add_card(slide, cx, cy, cw, Inches(1.10), label, val, sub)

    # ── 2. 정성 효과 (좌측 하단) ──
    qy = y + Inches(1.78)
    add_section(slide, BODY_X, qy, Inches(6.0), "2. 정성적 효과")
    add_box(slide, BODY_X, qy + Inches(0.42), Inches(6.0), Inches(1.85),
            fill=C_LIGHT, line=C_BORDER, line_w=0.5)
    add_text_lines(slide, [
        ("노하우 디지털화", 9.5, True, C_ACCENT),
        ("• 담당자 경험 → SHAP 데이터로 영속 (퇴사·교대 시도 누적)", 8.5, False, C_BODY),
        ("교대 인수인계 표준화", 9.5, True, C_ACCENT),
        ("• 수기 일지 → CSV / Markdown 자동 생성", 8.5, False, C_BODY),
        ("ISO 9001 이력 보존", 9.5, True, C_ACCENT),
        ("• JSON 영속 + SHA-256 모델 무결성 검증", 8.5, False, C_BODY),
        ("부서장 의사결정 지원", 9.5, True, C_ACCENT),
        ("• 일일 자연어 보고서 자동 생성 (Markdown)", 8.5, False, C_BODY),
    ], BODY_X + Inches(0.15), qy + Inches(0.50), Inches(5.7), Inches(1.7), line_spacing=1.16)

    # ── 3. 현장 적용 가능성 (우측 하단) ──
    rx = BODY_X + Inches(6.18)
    rw = Inches(6.0)
    add_section(slide, rx, qy, rw, "3. 현장 적용 가능성")
    add_box(slide, rx, qy + Inches(0.42), rw, Inches(1.85),
            fill=C_LIGHT, line=C_BORDER, line_w=0.5)
    add_text_lines(slide, [
        ("온프레미스 배포", 9.5, True, C_ACCENT),
        ("• 공정 데이터 외부 전송 0 (보안 우려 해소)", 8.5, False, C_BODY),
        ("최소 사양", 9.5, True, C_ACCENT),
        ("• 8GB RAM 보급형 PC · CPU 전용 가능 (GPU 불필요)", 8.5, False, C_BODY),
        ("원클릭 설치", 9.5, True, C_ACCENT),
        ("• install.bat 더블클릭 · IT 인력 부담 최소화", 8.5, False, C_BODY),
        ("다설비 확장 준비", 9.5, True, C_ACCENT),
        ("• 사이드바 설비 선택기 UI 이미 구현", 8.5, False, C_BODY),
    ], rx + Inches(0.15), qy + Inches(0.50), rw - Inches(0.2), Inches(1.7), line_spacing=1.16)

    # 하단 스크린샷 (크게)
    ssy = qy + Inches(2.42)
    add_dashed_box(slide, BODY_X, ssy, Inches(12.18), Inches(0.95),
                   "[스크린샷: '생산 이력' 탭 — 장비 건강도 스코어카드 7개 + 정비 권고 + 검증셋 1,379건 시계열]")


def fill_slide8_2_impact(prs, prs_clean):
    """Slide 8-2 — 기대 효과 (2/2): 수평·수직 확장 + 본선 시연 가치 + 차별화 배너
    (사업화 로드맵은 사용자 요청으로 제거)
    """
    new_slide = duplicate_slide(prs, prs_clean, 7)
    remove_guide_box(new_slide)

    y = BODY_Y

    # ── 1. 향후 고도화 — 수평·수직 확장 (상단) ──
    y_sec = add_section(new_slide, BODY_X, y, Inches(12.18),
                        "1. 향후 고도화 — 수평·수직 확장")

    # 수평 (좌측)
    add_box(new_slide, BODY_X, y_sec + Inches(0.10), Inches(6.0), Inches(1.55),
            fill=C_LIGHT, line=C_BORDER, line_w=0.5)
    add_text(new_slide, "수평 확장 — 도메인 / 설비",
             BODY_X + Inches(0.15), y_sec + Inches(0.16),
             Inches(5.7), Inches(0.28), size=13, bold=True, color=C_ACCENT)
    add_text_lines(new_slide, [
        ("• 사출성형 → 다이캐스팅 / CNC / 프레스 (24센서 표준 인터페이스)", 10, False, C_BODY),
        ("• KAMP 12종 공공 데이터셋 → 공정별 모델 카탈로그 SaaS", 10, False, C_BODY),
        ("• 다지점 공장 통합 모니터링 (사이드바 설비 선택기 → 클러스터)", 10, False, C_BODY),
    ], BODY_X + Inches(0.18), y_sec + Inches(0.50), Inches(5.7), Inches(1.05),
        line_spacing=1.28)

    # 수직 (우측)
    rx = BODY_X + Inches(6.18)
    rw = Inches(6.0)
    add_box(new_slide, rx, y_sec + Inches(0.10), rw, Inches(1.55),
            fill=C_LIGHT, line=C_BORDER, line_w=0.5)
    add_text(new_slide, "수직 확장 — AI · 분석 기술",
             rx + Inches(0.15), y_sec + Inches(0.16),
             rw - Inches(0.2), Inches(0.28), size=13, bold=True, color=C_ACCENT)
    add_text_lines(new_slide, [
        ("• Autoencoder → Transformer · TimeSeries Foundation Model", 10, False, C_BODY),
        ("• 예측 정비 모듈 (잔여수명 RUL 추정 + 정비 권고)", 10, False, C_BODY),
        ("• 강화학습 기반 자동 처방 (Closed-Loop Control · 장기)", 10, False, C_BODY),
    ], rx + Inches(0.18), y_sec + Inches(0.50), rw - Inches(0.2), Inches(1.05),
        line_spacing=1.28)

    # ── 2. 본선 1일 시연 핵심 가치 (중단) ──
    cy = y + Inches(2.30)
    add_section(new_slide, BODY_X, cy, Inches(12.18),
                "2. 본선 1일 구현 — 실측 가능한 시연 가치")
    demo_values = [
        ("탐지",  "ROC-AUC",  "0.9254 [실측]"),
        ("진단",  "SHAP",     "24센서 기여도 [실측]"),
        ("처방",  "처방 카드", "24개 자동 [실측]"),
        ("추적",  "JSON / CSV", "이력 영속 [실측]"),
        ("재학습", "Active Learning", "30건 누적 트리거"),
    ]
    dw = Inches(2.36)
    for i, (cat, t, d) in enumerate(demo_values):
        dx = BODY_X + i * (dw + Inches(0.05))
        dy = cy + Inches(0.42)
        add_box(new_slide, dx, dy, dw, Inches(1.10),
                fill=C_LIGHT, line=C_ACCENT, line_w=1)
        add_box(new_slide, dx, dy, dw, Inches(0.30), fill=C_ACCENT)
        add_text(new_slide, cat, dx, dy, dw, Inches(0.30),
                 size=13, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF),
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(new_slide, t, dx + Inches(0.1), dy + Inches(0.38),
                 dw - Inches(0.2), Inches(0.30), size=11, bold=True, color=C_BODY,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(new_slide, d, dx + Inches(0.1), dy + Inches(0.74),
                 dw - Inches(0.2), Inches(0.32), size=9.5, color=C_SUB,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # ── 3. 차별화 한 줄 배너 (하단) ──
    bby = cy + Inches(1.95)
    add_box(new_slide, BODY_X, bby, Inches(12.18), Inches(0.65), fill=C_ACCENT)
    add_text(new_slide,
             '"상용 0건 · 논문 4편 따로 — 그걸 한국 사출성형 24센서 위에서 한 화면에 묶은 첫 MVP"',
             BODY_X, bby + Inches(0.12), Inches(12.18), Inches(0.45),
             size=16, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF),
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


# ═══════════════════════════════════════════════════════════════
# 메인
# ═══════════════════════════════════════════════════════════════

def fill_slide7_5_gallery(prs, prs_clean):
    """Slide 10 — 본선 1일 추가 구현 계획
    오프라인 본선(5/22, 09:30~17:30, 하루)에 예선 MVP 위에 얹을 핵심 5 기능 + 시연 흐름.
    """
    new_slide = duplicate_slide(prs, prs_clean, 6)
    remove_guide_box(new_slide)

    y = BODY_Y

    # ── 1. 본선 추가 구현 항목 P1~P5 (상단, 5 카드 가로 배치) ──
    y_sec = add_section(new_slide, BODY_X, y, Inches(12.18),
                        "1. 본선 1일 추가 구현 항목 — 예선 MVP 위에 얹는 5 기능")
    priorities = [
        ("P1", "실시간 PLC 연동",       "OPC-UA / MQTT",
         "사출성형기 → 대시보드 무중단 스트리밍 · 슬라이더 입력을 실제 센서로 치환"),
        ("P2", "알람 자동화",            "SMS / 이메일",
         "심각도 3단계(경고/위험/긴급) 자동 발송 · 작업자→반장→부서장 에스컬레이션"),
        ("P3", "멀티 설비 관리",         "사이드바 설비 선택",
         "사출 / 다이캐스팅 / CNC 모델 카탈로그 자동 로딩 · 다지점 통합 모니터링"),
        ("P4", "예측 정비",              "Predictive Maintenance",
         "잔여수명(RUL) 추정 + 정비 권고 · 사후 대응 → 사전 정비 전환"),
        ("P5", "클라우드 옵션",          "AWS 배포",
         "원격 모니터링 · 다지점 통합 대시보드 · 권한 분리 (작업자/관리자)"),
    ]
    pw = Inches(2.36)
    p_gap = Inches(0.05)
    p_h = Inches(2.10)
    py_start = y_sec + Inches(0.15)
    for i, (pid, title, sub, desc) in enumerate(priorities):
        px = BODY_X + i * (pw + p_gap)
        # 카드 배경
        add_box(new_slide, px, py_start, pw, p_h,
                fill=C_LIGHT, line=C_ACCENT, line_w=1)
        # 상단 ID 띠
        add_box(new_slide, px, py_start, pw, Inches(0.36), fill=C_ACCENT)
        add_text(new_slide, pid, px, py_start, pw, Inches(0.36),
                 size=14, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF),
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        # 타이틀
        add_text(new_slide, title,
                 px + Inches(0.10), py_start + Inches(0.44),
                 pw - Inches(0.2), Inches(0.30),
                 size=12, bold=True, color=C_BODY, align=PP_ALIGN.CENTER)
        # 서브 (기술 키워드)
        add_text(new_slide, sub,
                 px + Inches(0.10), py_start + Inches(0.76),
                 pw - Inches(0.2), Inches(0.26),
                 size=9.5, bold=True, color=C_ACCENT, align=PP_ALIGN.CENTER)
        # 설명
        add_text(new_slide, desc,
                 px + Inches(0.12), py_start + Inches(1.08),
                 pw - Inches(0.24), Inches(0.95),
                 size=8.5, color=C_SUB)

    # ── 2. 본선 1일 시연 흐름 (하단, 5단계 타임라인) ──
    fy = y + Inches(2.85)
    add_section(new_slide, BODY_X, fy, Inches(12.18),
                "2. 본선 1일 시연 흐름 — 09:30~17:30 (8시간) 운영 계획")
    timeline = [
        ("09:30 ~ 11:00", "환경 셋업 + 예선 재현",
         "install.bat · verify_model.py · 예선 MVP 결과 재현 (ROC-AUC 0.9254 확인)"),
        ("11:00 ~ 13:00", "P1·P2 우선 구현",
         "OPC-UA 모의 스트림 연결 + SMS 알람 통합 (P3-P5 보다 평가 가중치 큼)"),
        ("13:00 ~ 15:00", "P3 멀티 설비 통합",
         "사이드바 설비 선택기 → 모델 카탈로그 연결 · 다지점 동시 모니터링 시연"),
        ("15:00 ~ 16:30", "P4 예측 정비 프로토타입",
         "RUL 모듈 · 정비 권고 카드 추가 · '사후 대응 → 사전 정비' 가치 어필"),
        ("16:30 ~ 17:30", "최종 시연 리허설",
         "8 시나리오 End-to-End 데모 (정상→경고→위험→긴급) · Q&A 30개 대비"),
    ]
    tw = Inches(2.36)
    t_gap = Inches(0.05)
    t_h = Inches(1.40)
    ty_start = fy + Inches(0.42)
    for i, (t, name, desc) in enumerate(timeline):
        tx = BODY_X + i * (tw + t_gap)
        add_box(new_slide, tx, ty_start, tw, t_h,
                fill=C_LIGHT, line=C_BORDER, line_w=0.5)
        # 시간 띠
        add_box(new_slide, tx, ty_start, tw, Inches(0.32), fill=C_ACCENT)
        add_text(new_slide, t, tx, ty_start, tw, Inches(0.32),
                 size=10.5, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF),
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        # 단계명
        add_text(new_slide, name,
                 tx + Inches(0.10), ty_start + Inches(0.38),
                 tw - Inches(0.2), Inches(0.30),
                 size=10.5, bold=True, color=C_BODY, align=PP_ALIGN.CENTER)
        # 작업 내용
        add_text(new_slide, desc,
                 tx + Inches(0.12), ty_start + Inches(0.72),
                 tw - Inches(0.24), Inches(0.65),
                 size=8.5, color=C_SUB)


def main():
    print(f"[로드] {TEMPLATE}")
    prs = Presentation(TEMPLATE)
    prs_clean = Presentation(TEMPLATE)  # 복제 전용 깨끗한 원본
    print(f"   원본 슬라이드 수: {len(prs.slides)}")

    print("[빌드] Slide 1 표지")
    fill_slide1_cover(prs)
    print("[빌드] Slide 2 문제 정의")
    fill_slide2_problem(prs)
    print("[빌드] Slide 3 솔루션 개요")
    fill_slide3_solution(prs)
    print("[빌드] Slide 4-1 주요 기능 MVP 핵심")
    fill_slide4_1_features(prs)
    print("[빌드] Slide 4-2 주요 기능 차별화 + 본선 (복제)")
    fill_slide4_2_features(prs, prs_clean)
    print("[빌드] Slide 5-1 데이터·처리")
    fill_slide5_1_data(prs)
    print("[빌드] Slide 5-2 AI 기술·차별화 (복제)")
    fill_slide5_2_data(prs, prs_clean)
    print("[빌드] Slide 6 유즈케이스")
    fill_slide6_usecase(prs)
    print("[빌드] Slide 7 MVP 범위")
    fill_slide7_scope(prs)
    print("[빌드] Slide 7.5 주요 화면 시연 갤러리 (신규)")
    fill_slide7_5_gallery(prs, prs_clean)
    print("[빌드] Slide 8-1 기대 효과")
    fill_slide8_1_impact(prs)
    print("[빌드] Slide 8-2 확장 + 마일스톤 (복제)")
    fill_slide8_2_impact(prs, prs_clean)

    # ── 슬라이드 순서 재배열: 분할 슬라이드가 짝 슬라이드 직후에 오도록 ──
    # 현재 순서(0-idx): 0 Cover, 1 문제, 2 솔루션, 3 주요기능1, 4 데이터1,
    #                   5 시나리오, 6 MVP범위, 7 기대1, 8 주요기능2(dup),
    #                   9 데이터2(dup), 10 갤러리(dup), 11 기대2(dup)
    # 목표:            0 Cover, 1 문제, 2 솔루션, 3 주요기능1, 4 주요기능2,
    #                   5 데이터1, 6 데이터2, 7 시나리오, 8 MVP범위,
    #                   9 갤러리, 10 기대1, 11 기대2
    slides = list(prs.slides)
    if len(slides) == 12:
        new_order = [
            slides[0],   # Cover
            slides[1],   # 1. 문제 정의
            slides[2],   # 2. 솔루션
            slides[3],   # 3. 주요 기능 1
            slides[8],   # 3. 주요 기능 2 (dup of 3)
            slides[4],   # 4. 데이터·기술 1
            slides[9],   # 4. 데이터·기술 2 (dup of 4)
            slides[5],   # 5. 사용자 시나리오
            slides[6],   # 6. MVP 범위
            slides[10],  # 주요 화면 시연 갤러리 (dup of 6)
            slides[7],   # 7. 기대 효과 1
            slides[11],  # 7. 기대 효과 2 (dup of 7)
        ]
        reorder_slides(prs, new_order)
        print("[정렬] 슬라이드 순서 재배열 완료 (자연스러운 흐름)")

    prs.save(OUT)
    print(f"\n[완료] {OUT}")
    print(f"   최종 슬라이드 수: {len(prs.slides)}")


if __name__ == '__main__':
    main()
