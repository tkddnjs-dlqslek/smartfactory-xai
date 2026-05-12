"""
스마트 공장 XAI 이상탐지 플랫폼 — 예선 기획서 PPT (공식 양식)
hackathon_template_full.md 양식 100% 준수: 8슬라이드 (표지 1 + 본문 7)
디자인: 흰 배경 + 진한 네이비 본문 타이틀 + 회색 헤더 + 검정 본문
"""
import sys, os, json
sys.path.insert(0, os.path.dirname(os.path.dirname(os.path.abspath(__file__))))

import numpy as np
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# ── 경로 ──
BASE_DIR   = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
RESULT_DIR = os.path.join(BASE_DIR, 'results')
MODEL_DIR  = os.path.join(BASE_DIR, 'models')
OUT_PATH   = os.path.join(BASE_DIR, 'output', '스마트공장XAI_예선기획서_공식양식.pptx')
os.makedirs(os.path.join(BASE_DIR, 'output'), exist_ok=True)

# ── 실제 수치 로드 ──
with open(os.path.join(RESULT_DIR, 'metrics.json'), encoding='utf-8') as f:
    M = json.load(f)
ROC, PR, F1 = M['roc_auc'], M['pr_auc'], M['f1']
REC, PREC, THR = M['recall'], M['precision'], M['threshold']

# SHAP 상위 센서
SHAP_TOP = []
shap_path = os.path.join(RESULT_DIR, 'shap_values.npy')
if os.path.exists(shap_path):
    from src.config import SENSOR_COLS
    sv  = np.load(shap_path)
    ma  = np.abs(sv).mean(axis=0)
    idx = np.argsort(ma)[::-1][:5]
    SHAP_TOP = [(SENSOR_COLS[i], float(ma[i])) for i in idx]

# 스코어링 통계
SC = None
sc_path = os.path.join(RESULT_DIR, 'scored_unlabeled.parquet')
if os.path.exists(sc_path):
    import pandas as pd
    df  = pd.read_parquet(sc_path)
    anom = (df['recon_error'] >= THR).sum()
    SC  = {'total': len(df), 'anomaly': int(anom), 'rate': anom / len(df) * 100}

# ══════════════════════════════════════════════════════════════
# 디자인 토큰 (공식 양식)
# ══════════════════════════════════════════════════════════════
C_BG       = RGBColor(0xFF, 0xFF, 0xFF)  # 흰 배경
C_NAVY     = RGBColor(0x1A, 0x2B, 0x5C)  # 진한 네이비 (본문 타이틀)
C_NAVY_BD  = RGBColor(0x0E, 0x1B, 0x40)  # 더 진한 네이비 (밑줄)
C_BLACK    = RGBColor(0x1A, 0x1A, 0x1A)  # 본문 검정
C_GRAY     = RGBColor(0x88, 0x88, 0x88)  # 헤더 우측 회색
C_GRAY_BG  = RGBColor(0xF5, 0xF5, 0xF5)  # 카드 배경 연회색
C_GRAY_BD  = RGBColor(0xDD, 0xDD, 0xDD)  # 카드 테두리
C_RED      = RGBColor(0xC0, 0x00, 0x00)  # 안내박스/강조 빨강
C_RED_LT   = RGBColor(0xFF, 0xEE, 0xEE)  # 강조 박스 연빨강 배경
C_BLUE_LT  = RGBColor(0xEE, 0xF2, 0xFA)  # 정보 박스 연파랑 배경
C_GREEN    = RGBColor(0x2E, 0x7D, 0x32)  # 완료/긍정 녹색
C_DIM      = RGBColor(0x55, 0x55, 0x55)  # 본문 부가 회색
C_FOOTER   = RGBColor(0xAA, 0xAA, 0xAA)  # 푸터

W = Inches(13.33)
H = Inches(7.5)

FONT = "Noto Sans KR"

# ══════════════════════════════════════════════════════════════
# 헬퍼
# ══════════════════════════════════════════════════════════════
def new_prs():
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H
    return prs

def blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])

def fill_bg(slide, color=C_BG):
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = color

def add_box(slide, left, top, width, height, fill=None, border=None, border_w=Pt(0.5)):
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.line.width = Pt(0)
    f = shape.fill
    if fill is not None:
        f.solid()
        f.fore_color.rgb = fill
    else:
        f.background()
    if border is not None:
        shape.line.color.rgb = border
        shape.line.width = border_w
    return shape

def add_text(slide, text, left, top, width, height,
             size=Pt(11), bold=False, color=C_BLACK,
             align=PP_ALIGN.LEFT, italic=False, font_name=FONT,
             v_anchor=None):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(2); tf.margin_right = Pt(2)
    tf.margin_top = Pt(1); tf.margin_bottom = Pt(1)
    if v_anchor is not None:
        tf.vertical_anchor = v_anchor
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = size
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = font_name
    return txBox

def add_text_lines(slide, lines, left, top, width, height,
                   line_spacing=1.15, font_name=FONT):
    """lines: [(txt, size, bold, color), ...]"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(2); tf.margin_right = Pt(2)
    first = True
    for line in lines:
        if len(line) == 4:
            txt, sz, bold, col = line
            italic = False
        else:
            txt, sz, bold, col, italic = line
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = line_spacing
        run = p.add_run()
        run.text = txt
        run.font.size = sz
        run.font.bold = bold
        run.font.italic = italic
        run.font.color.rgb = col
        run.font.name = font_name
    return txBox

def navy_underline(slide, left, top, width, h=Pt(2.5)):
    shape = slide.shapes.add_shape(1, left, top, width, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = C_NAVY_BD
    shape.line.width = Pt(0)
    return shape

# ══════════════════════════════════════════════════════════════
# 공식 양식 레이아웃 (모든 본문 슬라이드 공통)
# ══════════════════════════════════════════════════════════════
def official_header(slide, num, title_kor, title_eng):
    """■ 예선 산출물 – N. 섹션명 (English) | 우측: 회색 브랜딩"""
    # 좌측 ■ 마커 + 텍스트
    add_text(slide,
             f"■ 예선 산출물 – {num}. {title_kor} ({title_eng})",
             Inches(0.45), Inches(0.25), Inches(8.0), Inches(0.32),
             size=Pt(10.5), bold=True, color=C_BLACK)
    # 우측 회색 브랜딩
    add_text(slide,
             "2026 스마트 공장 운영 시스템 MVP 개발 해커톤",
             Inches(8.5), Inches(0.27), Inches(4.5), Inches(0.3),
             size=Pt(8.5), color=C_GRAY, align=PP_ALIGN.RIGHT)

def official_title(slide, num, title_kor, title_eng):
    """본문 타이틀: N. 섹션명 (English) — 진한 네이비 + 가로 밑줄"""
    add_text(slide,
             f"{num}. {title_kor} ({title_eng})",
             Inches(0.45), Inches(0.7), Inches(12.5), Inches(0.5),
             size=Pt(20), bold=True, color=C_NAVY)
    navy_underline(slide, Inches(0.45), Inches(1.2), Inches(12.45))

def official_footer(slide):
    """푸터: 좌측 빈공간 + 우측 차세대융합기술연구원 + DACON"""
    add_text(slide,
             "주최/주관: 차세대융합기술연구원 (AICT)  ·  DACON",
             Inches(7.5), Inches(7.05), Inches(5.5), Inches(0.3),
             size=Pt(8), color=C_FOOTER, align=PP_ALIGN.RIGHT)

def add_bullet_card(slide, x, y, w, h, title, lines, accent=C_NAVY):
    """제목 카드 + 불릿 라인"""
    add_box(slide, x, y, w, h, fill=C_GRAY_BG, border=C_GRAY_BD)
    # 좌측 accent 막대
    add_box(slide, x, y, Pt(3), h, fill=accent)
    add_text(slide, title, x + Inches(0.12), y + Inches(0.05),
             w - Inches(0.2), Inches(0.3),
             size=Pt(10.5), bold=True, color=accent)
    body_y = y + Inches(0.36)
    body_lines = [(f"• {ln}", Pt(8.5), False, C_BLACK) for ln in lines]
    add_text_lines(slide, body_lines, x + Inches(0.15), body_y,
                   w - Inches(0.25), h - Inches(0.4))

# ══════════════════════════════════════════════════════════════
# Slide 1 — 표지
# ══════════════════════════════════════════════════════════════
def slide01_cover(prs):
    sl = blank_slide(prs); fill_bg(sl)

    # 메인 제목 (중앙 정렬)
    add_text(sl, "2026 스마트 공장 운영 시스템",
             Inches(0.5), Inches(2.0), Inches(12.33), Inches(0.7),
             size=Pt(34), bold=True, color=C_NAVY, align=PP_ALIGN.CENTER)
    add_text(sl, "MVP 개발 해커톤",
             Inches(0.5), Inches(2.65), Inches(12.33), Inches(0.7),
             size=Pt(34), bold=True, color=C_NAVY, align=PP_ALIGN.CENTER)

    # 가로 밑줄
    navy_underline(sl, Inches(3.0), Inches(3.42), Inches(7.33), h=Pt(3.5))

    # 서브타이틀
    add_text(sl, "예선 기획서",
             Inches(0.5), Inches(3.6), Inches(12.33), Inches(0.5),
             size=Pt(20), bold=True, color=C_BLACK, align=PP_ALIGN.CENTER)

    # 표 (2x2) - 중앙 배치
    tbl_x = Inches(4.0)
    tbl_y = Inches(4.5)
    col_lbl_w = Inches(2.0)
    col_val_w = Inches(3.5)
    row_h = Inches(0.55)

    # Row 1 - 팀명
    add_box(sl, tbl_x, tbl_y, col_lbl_w, row_h, fill=C_GRAY_BG, border=C_GRAY_BD)
    add_text(sl, "팀명", tbl_x, tbl_y, col_lbl_w, row_h,
             size=Pt(12), bold=True, color=C_NAVY, align=PP_ALIGN.CENTER,
             v_anchor=None)
    add_box(sl, tbl_x + col_lbl_w, tbl_y, col_val_w, row_h, fill=C_BG, border=C_GRAY_BD)
    add_text(sl, "[ 팀명 작성란 ]", tbl_x + col_lbl_w, tbl_y, col_val_w, row_h,
             size=Pt(12), color=C_DIM, align=PP_ALIGN.CENTER, italic=True)

    # Row 2 - 프로젝트명
    add_box(sl, tbl_x, tbl_y + row_h, col_lbl_w, row_h, fill=C_GRAY_BG, border=C_GRAY_BD)
    add_text(sl, "프로젝트명", tbl_x, tbl_y + row_h, col_lbl_w, row_h,
             size=Pt(12), bold=True, color=C_NAVY, align=PP_ALIGN.CENTER)
    add_box(sl, tbl_x + col_lbl_w, tbl_y + row_h, col_val_w, row_h, fill=C_BG, border=C_GRAY_BD)
    add_text(sl, "SmartFactory XAI — 사출성형 이상탐지·진단 플랫폼",
             tbl_x + col_lbl_w, tbl_y + row_h, col_val_w, row_h,
             size=Pt(11), bold=True, color=C_BLACK, align=PP_ALIGN.CENTER)

    # 푸터 로고 영역 (텍스트 대체)
    add_text(sl, "주최 · 주관",
             Inches(8.5), Inches(6.7), Inches(2.5), Inches(0.25),
             size=Pt(8), color=C_GRAY, align=PP_ALIGN.RIGHT)
    add_text(sl, "차세대융합기술연구원 (AICT)  ·  DACON",
             Inches(7.5), Inches(6.95), Inches(5.5), Inches(0.3),
             size=Pt(10), bold=True, color=C_NAVY, align=PP_ALIGN.RIGHT)

# ══════════════════════════════════════════════════════════════
# Slide 2 — 1. 문제 정의 (Problem Definition)
# ══════════════════════════════════════════════════════════════
def slide02_problem(prs):
    sl = blank_slide(prs); fill_bg(sl)
    official_header(sl, 1, "문제 정의", "Problem Definition")
    official_title(sl, 1, "문제 정의", "Problem Definition")

    # ── 좌측 상단: 현장 문제 4가지 ──
    add_text(sl, "[ 해결하고자 하는 제조 현장 문제 ]",
             Inches(0.45), Inches(1.42), Inches(6.4), Inches(0.32),
             size=Pt(11), bold=True, color=C_NAVY)
    pains = [
        ("불량 발생 인지 지연", "작업자 육안 검사 한계 → 수백 개 불량 후 인지"),
        ("원인 분석 1~2시간", "24개 센서 동시 모니터링 불가 → 재발 방지 어려움"),
        ("사후 대응 반복 사이클", "불량 → 라인 정지 → 수동 검사 → 재가동 무한 반복"),
        ("IT 전담인력 부재", "중소 사출성형 공장 50인 규모 — AI 도입 진입장벽"),
    ]
    for i, (t, b) in enumerate(pains):
        y = Inches(1.78) + i * Inches(0.7)
        add_box(sl, Inches(0.45), y, Inches(6.4), Inches(0.62), fill=C_GRAY_BG, border=C_GRAY_BD)
        add_box(sl, Inches(0.45), y, Pt(3), Inches(0.62), fill=C_RED)
        add_text(sl, f"{i+1}.  {t}", Inches(0.6), y + Inches(0.06),
                 Inches(6.0), Inches(0.28), size=Pt(10.5), bold=True, color=C_BLACK)
        add_text(sl, b, Inches(0.6), y + Inches(0.32),
                 Inches(6.0), Inches(0.26), size=Pt(8.5), color=C_DIM)

    # ── 우측 상단: 현장 페르소나 + 적용 영역 + 활용 목적 (통합) ──
    add_text(sl, "[ 적용 대상 페르소나 · 활용 목적 · 시장 배경 ]",
             Inches(7.0), Inches(1.42), Inches(6.0), Inches(0.32),
             size=Pt(11), bold=True, color=C_NAVY)
    add_box(sl, Inches(7.0), Inches(1.78), Inches(6.0), Inches(3.82),
            fill=C_BLUE_LT, border=C_GRAY_BD)
    add_text_lines(sl, [
        ("현장 페르소나 (구체 수치 기반)", Pt(10.5), True, C_NAVY),
        ("• 작업자 김OO (사출성형 5년차, 시급 12,000원, 3교대 근무)", Pt(8.5), False, C_BLACK),
        ("• 일일 생산: 약 6,000샷 · 점검 가능량 35,000샷 중 일부만", Pt(8.5), False, C_BLACK),
        ("• 월 평균 불량 25건 · 인지 지연 평균 2~4시간 · 수기 일지", Pt(8.5), False, C_BLACK),
        ("적용 대상 & 영역", Pt(10.5), True, C_NAVY),
        ("• 중소 사출성형 공장 (50인 규모, IT 전담 인력 부재)", Pt(8.5), False, C_BLACK),
        ("• 정부 스마트공장 보급사업 (2025~2027, 3만개사 대상)", Pt(8.5), False, C_BLACK),
        ("• 24센서 실시간 모니터링 + SHAP 원인 진단 + 처방 자동", Pt(8.5), False, C_BLACK),
        ("활용 목적 (도입 후 변화)", Pt(10.5), True, C_RED),
        ("• 점검 시간 2시간 → 수 초 (자동화) · 작업자 시간 -90%", Pt(8.5), True, C_BLACK),
        ("• 수동 점검 인력 1명 대체 = 연 5,000만원 인건비 절감", Pt(8.5), True, C_BLACK),
        ("• 불량 1건당 50만원 손실 × 200건/년 = 1억원 직접 회피", Pt(8.5), True, C_RED),
    ], Inches(7.15), Inches(1.85), Inches(5.7), Inches(3.72), line_spacing=1.16)

    # ── 하단 전폭: 해결 후 KPI ──
    add_text(sl, "[ 해결 후 기대 KPI ]",
             Inches(0.45), Inches(5.78), Inches(12.5), Inches(0.32),
             size=Pt(11), bold=True, color=C_NAVY)
    kpis = [
        ("Recall 0.667+", "불량 10건 중\n6.7건 자동 탐지", C_RED),
        ("응답 즉시", "원인 분석\n2시간 → 수 초", C_NAVY),
        ("연 1억원 절감", "월 25건 × 50만원\n× 12개월 ROI", C_GREEN),
        ("ROI 4.2배", "월 200만 SaaS 대비\n절감액 회수율", C_NAVY),
    ]
    for i, (big, sub, clr) in enumerate(kpis):
        x = Inches(0.45) + i * Inches(3.13)
        add_box(sl, x, Inches(6.13), Inches(3.0), Inches(0.85),
                fill=C_GRAY_BG, border=clr, border_w=Pt(1.2))
        add_text(sl, big, x + Inches(0.1), Inches(6.18), Inches(2.8), Inches(0.3),
                 size=Pt(13), bold=True, color=clr)
        add_text(sl, sub, x + Inches(0.1), Inches(6.5), Inches(2.8), Inches(0.45),
                 size=Pt(8.5), color=C_DIM)

    official_footer(sl)

# ══════════════════════════════════════════════════════════════
# Slide 3 — 2. 솔루션 개요 (Solution Overview)
# ══════════════════════════════════════════════════════════════
def slide03_solution(prs):
    sl = blank_slide(prs); fill_bg(sl)
    official_header(sl, 2, "제안 솔루션 개요", "Solution Overview")
    official_title(sl, 2, "제안 솔루션 개요", "Solution Overview")

    # ── 핵심 메시지 (강조) ──
    add_box(sl, Inches(0.45), Inches(1.42), Inches(12.45), Inches(0.7),
            fill=C_NAVY, border=None)
    add_text(sl, "\"라인이 멈추기 전에 AI가 먼저 알린다\"",
             Inches(0.45), Inches(1.45), Inches(12.45), Inches(0.4),
             size=Pt(15), bold=True, color=C_BG, align=PP_ALIGN.CENTER)
    add_text(sl, "정상 데이터로만 학습한 Autoencoder + SHAP XAI = 라벨 없이 이상탐지 + 원인 즉시 설명",
             Inches(0.45), Inches(1.82), Inches(12.45), Inches(0.28),
             size=Pt(10), color=C_BG, align=PP_ALIGN.CENTER)

    # ── 좌측: 4단계 운영 자동화 (탐지→진단→처방→추적) ──
    add_text(sl, "[ 핵심 아이디어: 4단계 운영 자동화 ]",
             Inches(0.45), Inches(2.32), Inches(6.4), Inches(0.3),
             size=Pt(11), bold=True, color=C_NAVY)
    steps = [
        ("①", "탐지 (Detect)",  "Autoencoder 복원오차 → ROC-AUC 0.9254"),
        ("②", "진단 (Diagnose)","SHAP KernelSHAP → 24센서 실시간 기여도"),
        ("③", "처방 (Prescribe)","센서별 처방 카드 24개 (조작가능 / 정비필요 분류)"),
        ("④", "추적 (Track)",  "이력 50건 + JSON 영속 + 일일 부서장 리포트 MD"),
    ]
    for i, (num, t, b) in enumerate(steps):
        y = Inches(2.62) + i * Inches(0.55)
        add_box(sl, Inches(0.45), y, Inches(6.4), Inches(0.48),
                fill=C_GRAY_BG, border=C_GRAY_BD)
        add_text(sl, num, Inches(0.5), y + Inches(0.03), Inches(0.5), Inches(0.42),
                 size=Pt(18), bold=True, color=C_RED, align=PP_ALIGN.CENTER)
        add_text(sl, t, Inches(1.0), y + Inches(0.03), Inches(2.0), Inches(0.22),
                 size=Pt(10), bold=True, color=C_NAVY)
        add_text(sl, b, Inches(1.0), y + Inches(0.25), Inches(5.4), Inches(0.22),
                 size=Pt(8), color=C_BLACK)

    # ── 피드백 루프 (플랫폼 관점 강조) ──
    _fy = Inches(4.85)
    add_box(sl, Inches(0.45), _fy, Inches(6.4), Inches(0.4),
            fill=C_RED_LT, border=C_RED, border_w=Pt(1))
    add_text(sl, "↺ Feedback Loop",
             Inches(0.55), _fy + Inches(0.0), Inches(2.0), Inches(0.2),
             size=Pt(9), bold=True, color=C_RED)
    add_text(sl, "④ 추적 데이터 → ① 모델 재학습 (조치 결과·라벨 누적 → 정확도 지속 개선)",
             Inches(0.55), _fy + Inches(0.2), Inches(6.2), Inches(0.2),
             size=Pt(7.5), color=C_BLACK)

    # ── 우측: 시스템 아키텍처 ──
    add_text(sl, "[ 시스템 전체 구성 ]",
             Inches(7.0), Inches(2.32), Inches(6.0), Inches(0.3),
             size=Pt(11), bold=True, color=C_NAVY)
    add_box(sl, Inches(7.0), Inches(2.68), Inches(6.0), Inches(2.35),
            fill=C_GRAY_BG, border=C_GRAY_BD)

    arch = [
        ("입력 (Input)",  "24센서 실시간 / CSV 1행 / 슬라이더 시뮬레이터", C_NAVY),
        ("↓", "", C_DIM),
        ("처리 (Pipeline)", "StandardScaler → Autoencoder (24→16→8→16→24)", C_NAVY),
        ("↓", "", C_DIM),
        ("XAI (SHAP)", "KernelExplainer (배경 K-Means 50개) → 센서 기여도", C_NAVY),
        ("↓", "", C_DIM),
        ("출력 (Output)", "이상점수 + 심각도 3단계 + 처방 + 이력 누적", C_RED),
    ]
    for i, (t, b, clr) in enumerate(arch):
        y = Inches(2.78) + i * Inches(0.32)
        if t == "↓":
            add_text(sl, "↓", Inches(7.0), y, Inches(6.0), Inches(0.25),
                     size=Pt(11), bold=True, color=C_DIM, align=PP_ALIGN.CENTER)
        else:
            add_text(sl, t, Inches(7.15), y, Inches(2.0), Inches(0.28),
                     size=Pt(9.5), bold=True, color=clr)
            add_text(sl, b, Inches(9.15), y, Inches(3.8), Inches(0.28),
                     size=Pt(8.5), color=C_BLACK)

    # ── 하단: 제공 가치 4가지 (전체폭) ──
    add_text(sl, "[ 플랫폼 제공 가치 ]",
             Inches(0.45), Inches(5.18), Inches(12.5), Inches(0.3),
             size=Pt(11), bold=True, color=C_NAVY)
    values = [
        ("실시간 탐지", "ROC-AUC 0.9254\n불량 인지 즉시"),
        ("원인 자동 진단", "SHAP 24센서 기여도\n경험 의존 탈피"),
        ("처방 자동화", "센서별 24개 처방\n조치 시간 단축"),
        ("이력 누적·표준화", "교대 인수인계 CSV\n일일 리포트 MD"),
    ]
    for i, (t, b) in enumerate(values):
        x = Inches(0.45) + i * Inches(3.13)
        add_box(sl, x, Inches(5.5), Inches(3.0), Inches(1.45),
                fill=C_BLUE_LT, border=C_NAVY, border_w=Pt(1))
        add_text(sl, t, x + Inches(0.1), Inches(5.6), Inches(2.8), Inches(0.32),
                 size=Pt(11.5), bold=True, color=C_NAVY)
        add_text(sl, b, x + Inches(0.1), Inches(5.95), Inches(2.8), Inches(0.85),
                 size=Pt(9), color=C_BLACK)

    official_footer(sl)

# ══════════════════════════════════════════════════════════════
# Slide 4 — 3. 주요 기능 정의 (Key Features)
# ══════════════════════════════════════════════════════════════
def slide04_features(prs):
    sl = blank_slide(prs); fill_bg(sl)
    official_header(sl, 3, "주요 기능 정의", "Key Features")
    official_title(sl, 3, "주요 기능 정의", "Key Features")

    # ── 7가지 핵심 기능 (3+4 레이아웃) ──
    features = [
        ("F1", "AI 자동 이상탐지", "Autoencoder 복원오차 → ROC-AUC 0.9254\nF1 0.7324, 95% CI [0.597, 0.835]", "MVP"),
        ("F2", "SHAP 원인 진단",   "KernelExplainer 24센서 실시간 기여도\nTop-3 자동 도출 + Waterfall 시각화", "MVP"),
        ("F3", "처방 카드 자동",   "센서별 24개 처방 (조작가능/정비필요)\n현장 작업자 즉시 조치 가능", "MVP"),
        ("F4", "심각도 3단계",     "경고/위험/긴급 자동 분류\n알람 에스컬레이션 (작업자→반장→부서장)", "MVP"),
        ("F5", "교대 인수인계",     "이력 50건 누적 + JSON 영속\nCSV 다운로드 + 인수인계 메모", "MVP"),
        ("F6", "Cross-Machine 검증", "별도 금형 세트 (CN7+RG3) 외부 검증\nCircular Evaluation 방지", "본선"),
        ("F7", "일일 부서장 리포트", "Markdown 자동 생성\n심각도 분포 + 조치 현황 집계", "본선"),
    ]
    # 3행 x 4열 (7개 + 1개 empty 또는 마지막은 강조)
    cols = 4
    cell_w = Inches(3.07)
    cell_h = Inches(1.55)
    start_x = Inches(0.45)
    start_y = Inches(1.4)

    for i, (fid, title, desc, scope) in enumerate(features):
        row, col = i // cols, i % cols
        x = start_x + col * (cell_w + Inches(0.05))
        y = start_y + row * (cell_h + Inches(0.1))
        is_mvp = (scope == "MVP")
        add_box(sl, x, y, cell_w, cell_h, fill=C_GRAY_BG, border=C_NAVY if is_mvp else C_GRAY_BD)
        add_box(sl, x, y, Pt(3), cell_h, fill=C_RED if is_mvp else C_GRAY)
        # 기능 ID + 범위
        add_text(sl, fid, x + Inches(0.12), y + Inches(0.05),
                 Inches(0.4), Inches(0.25), size=Pt(11), bold=True, color=C_RED if is_mvp else C_GRAY)
        add_box(sl, x + cell_w - Inches(0.55), y + Inches(0.08), Inches(0.45), Inches(0.22),
                fill=C_RED if is_mvp else C_GRAY)
        add_text(sl, scope, x + cell_w - Inches(0.55), y + Inches(0.08),
                 Inches(0.45), Inches(0.22), size=Pt(7), bold=True, color=C_BG, align=PP_ALIGN.CENTER)
        # 제목
        add_text(sl, title, x + Inches(0.12), y + Inches(0.32),
                 cell_w - Inches(0.2), Inches(0.32), size=Pt(11), bold=True, color=C_NAVY)
        # 설명
        add_text(sl, desc, x + Inches(0.12), y + Inches(0.68),
                 cell_w - Inches(0.2), cell_h - Inches(0.7), size=Pt(8.5), color=C_BLACK)

    # 마지막 빈자리에 "통합 흐름도" 카드 — 데이터 + 사용자 + Feedback
    x_last = start_x + 3 * (cell_w + Inches(0.05))
    y_last = start_y + 1 * (cell_h + Inches(0.1))
    add_box(sl, x_last, y_last, cell_w, cell_h, fill=C_NAVY, border=None)
    add_text(sl, "통합 플랫폼 흐름 (데이터+사용자)",
             x_last + Inches(0.1), y_last + Inches(0.05),
             cell_w - Inches(0.2), Inches(0.3), size=Pt(10), bold=True, color=C_BG)
    add_text_lines(sl, [
        ("[ 데이터 흐름 ]", Pt(8), True, C_RED_LT),
        ("F1탐지 → F2진단(SHAP)", Pt(7.5), True, C_BG),
        ("→ F3처방 + F4심각도", Pt(7.5), True, C_BG),
        ("→ F5이력 + F7리포트", Pt(7.5), True, C_BG),
        ("↺ Feedback → 재학습", Pt(7.5), True, C_RED_LT),
        ("[ 사용자 흐름 ]", Pt(8), True, C_RED_LT),
        ("작업자 → 반장 → 부서장", Pt(7.5), True, C_BG),
    ], x_last + Inches(0.12), y_last + Inches(0.36),
       cell_w - Inches(0.22), cell_h - Inches(0.4), line_spacing=1.05)

    # ── 하단 전폭: 차별화 포지셔닝 ──
    add_box(sl, Inches(0.45), Inches(5.4), Inches(12.45), Inches(1.55),
            fill=C_BLUE_LT, border=C_NAVY, border_w=Pt(1))
    add_text(sl, "[ 본선 MVP 핵심 기능 명확화 — 경쟁 솔루션 대비 차별화 ]",
             Inches(0.6), Inches(5.5), Inches(12.2), Inches(0.3),
             size=Pt(11), bold=True, color=C_NAVY)
    headers = ["", "수동 점검", "범용 MES", "외부 컨설팅", "SmartFactory XAI"]
    rows = [
        ("실시간 탐지",      "✗",     "△ 임계값",   "✗",      "✓ AI 자동"),
        ("원인 진단",        "△ 경험", "✗",         "△ 사후",  "✓ SHAP"),
        ("처방 자동화",      "✗",     "✗",         "△ 보고서","✓ 24개"),
        ("월 비용 / TCO",    "인건비",  "수천만원/년", "회당 수백",  "월 200만"),
    ]
    col_w = [2.4, 2.0, 2.0, 2.0, 2.45]
    # 헤더
    x_cur = 0.6
    for ci, (h, w) in enumerate(zip(headers, col_w)):
        add_box(sl, Inches(x_cur), Inches(5.85), Inches(w), Inches(0.28),
                fill=C_NAVY if ci == 4 else C_GRAY_BG, border=C_GRAY_BD)
        add_text(sl, h, Inches(x_cur), Inches(5.85), Inches(w), Inches(0.28),
                 size=Pt(8.5), bold=True,
                 color=C_BG if ci == 4 else C_NAVY, align=PP_ALIGN.CENTER)
        x_cur += w
    # 데이터
    for ri, row in enumerate(rows):
        x_cur = 0.6
        for ci, (val, w) in enumerate(zip(row, col_w)):
            y = Inches(6.13) + ri * Inches(0.21)
            add_box(sl, Inches(x_cur), y, Inches(w), Inches(0.21),
                    fill=C_RED_LT if ci == 4 else C_BG, border=C_GRAY_BD)
            add_text(sl, val, Inches(x_cur), y, Inches(w), Inches(0.21),
                     size=Pt(7.5),
                     bold=(ci == 0 or ci == 4),
                     color=C_RED if ci == 4 else (C_NAVY if ci == 0 else C_BLACK),
                     align=PP_ALIGN.CENTER)
            x_cur += w

    official_footer(sl)

# ══════════════════════════════════════════════════════════════
# Slide 5 — 4. 데이터 및 기술 활용 계획 (Data & Tech Plan)
# ══════════════════════════════════════════════════════════════
def slide05_data_tech(prs):
    sl = blank_slide(prs); fill_bg(sl)
    official_header(sl, 4, "데이터 및 기술 활용 계획", "Data & Tech Plan")
    official_title(sl, 4, "데이터 및 기술 활용 계획", "Data & Tech Plan")

    # ── 좌측 상단: 데이터 ──
    add_text(sl, "[ 활용 데이터 — KAMP 공공데이터포털 ]",
             Inches(0.45), Inches(1.4), Inches(6.4), Inches(0.3),
             size=Pt(11), bold=True, color=C_NAVY)
    add_box(sl, Inches(0.45), Inches(1.72), Inches(6.4), Inches(2.0),
            fill=C_GRAY_BG, border=C_GRAY_BD)
    add_text_lines(sl, [
        ("출처: data.go.kr/data/15089213 (KAMP)", Pt(9), True, C_NAVY),
        ("• supervised_label_cn7.csv  6,736행 (학습 + Pseudo Hold-out)", Pt(8.5), False, C_BLACK),
        ("• moldset_labeled_cn7/rg3.csv  외부 검증 (Cross-Machine)", Pt(8.5), False, C_BLACK),
        ("• unlabeled_data.csv  795,315행 (대규모 스코어링)", Pt(8.5), False, C_BLACK),
        ("", Pt(4), False, C_BLACK),
        ("24개 센서 분류", Pt(9), True, C_NAVY),
        ("• 시간 5  /  위치 3  /  속도·RPM 3  /  압력 5  /  온도 8", Pt(8.5), False, C_BLACK),
        ("불량 클래스 0.58% (39 / 6,736) — 극심한 불균형", Pt(8.5), True, C_RED),
    ], Inches(0.6), Inches(1.78), Inches(6.1), Inches(1.9), line_spacing=1.2)

    # ── 좌측 하단: 데이터 처리 ──
    add_text(sl, "[ 데이터 처리 방식 ]",
             Inches(0.45), Inches(3.8), Inches(6.4), Inches(0.3),
             size=Pt(11), bold=True, color=C_NAVY)
    add_box(sl, Inches(0.45), Inches(4.12), Inches(6.4), Inches(1.45),
            fill=C_GRAY_BG, border=C_GRAY_BD)
    add_text_lines(sl, [
        ("• StandardScaler — train fit only (data leakage 방지)", Pt(8.5), False, C_BLACK),
        ("• 80/20 분할 (random_state=42, 재현성 고정)", Pt(8.5), False, C_BLACK),
        ("• 학습: 정상 7,238 / 검증: 정상 1,810 + 불량 81", Pt(8.5), False, C_BLACK),
        ("• Bootstrap 95% CI 1,000회 (소표본 통계 신뢰성 확보)", Pt(8.5), False, C_BLACK),
        ("• Pseudo Hold-out (마지막 20건 분리) + Cross-Machine 외부 검증", Pt(8.5), True, C_NAVY),
    ], Inches(0.6), Inches(4.18), Inches(6.1), Inches(1.35), line_spacing=1.2)

    # ── 우측 상단: AI/분석 기술 ──
    add_text(sl, "[ 적용 AI / 분석 기술 ]",
             Inches(7.0), Inches(1.4), Inches(6.0), Inches(0.3),
             size=Pt(11), bold=True, color=C_NAVY)
    add_box(sl, Inches(7.0), Inches(1.72), Inches(6.0), Inches(2.0),
            fill=C_GRAY_BG, border=C_GRAY_BD)
    add_text_lines(sl, [
        ("Autoencoder (PyTorch 2.3.1)", Pt(9), True, C_NAVY),
        ("• 구조: 24→16→8→16→24 (BatchNorm + MSE + Adam)", Pt(8.5), False, C_BLACK),
        ("• 비지도 학습 (정상 데이터만) — 0.58% 불균형 trivial 함정 회피", Pt(8.5), False, C_BLACK),
        ("• 임계값: 99th percentile + F1-optimal 균형", Pt(8.5), False, C_BLACK),
        ("", Pt(4), False, C_BLACK),
        ("SHAP XAI (shap 0.48)", Pt(9), True, C_NAVY),
        ("• KernelExplainer + K-Means 50개 배경 압축 (안정성)", Pt(8.5), False, C_BLACK),
        ("• 24센서 실시간 기여도 — Waterfall + Bar Summary", Pt(8.5), False, C_BLACK),
    ], Inches(7.15), Inches(1.78), Inches(5.7), Inches(1.9), line_spacing=1.2)

    # ── 우측 중단: 기술 스택 ──
    add_text(sl, "[ 사용 예정 기술 스택 ]",
             Inches(7.0), Inches(3.8), Inches(6.0), Inches(0.3),
             size=Pt(11), bold=True, color=C_NAVY)
    add_box(sl, Inches(7.0), Inches(4.12), Inches(6.0), Inches(1.45),
            fill=C_GRAY_BG, border=C_GRAY_BD)
    add_text_lines(sl, [
        ("• ML: PyTorch 2.3 + scikit-learn 1.2 + SHAP 0.48", Pt(8.5), False, C_BLACK),
        ("• UI: Streamlit 1.36 + Plotly 5.22 + python-pptx", Pt(8.5), False, C_BLACK),
        ("• 인프라: Python 3.11 (Anaconda) + 온프레미스 배포", Pt(8.5), False, C_BLACK),
        ("• 본선: OPC-UA / MQTT 실시간 연동, AWS 옵션", Pt(8.5), False, C_BLACK),
        ("• 보안: 모델 SHA-256 무결성 검증 + JSON 이력 영속", Pt(8.5), True, C_NAVY),
    ], Inches(7.15), Inches(4.18), Inches(5.7), Inches(1.35), line_spacing=1.2)

    # ── 하단 전폭: 학술 레퍼런스 + 차별화 + 제약 ──
    add_text(sl, "[ 학술 레퍼런스 + 차별화 + 기술적 제약 / 해결 전략 ]",
             Inches(0.45), Inches(5.65), Inches(12.5), Inches(0.3),
             size=Pt(11), bold=True, color=C_NAVY)
    # 좌측 레퍼런스
    add_box(sl, Inches(0.45), Inches(5.97), Inches(6.4), Inches(0.95),
            fill=C_BG, border=C_NAVY, border_w=Pt(1))
    add_text_lines(sl, [
        ("핵심 인용 4개 (전체 7개 인용)", Pt(8.5), True, C_NAVY),
        ("• MDPI Processes 13(3),912 (2025) — KAMP 직접 비교", Pt(7.5), False, C_BLACK),
        ("• arXiv:2511.08108 (2025) — LSTM+SHAP F1 0.92", Pt(7.5), False, C_BLACK),
        ("• Brito MAKE 6(1),16 (2024) — SHAP 베어링 98.5%", Pt(7.5), False, C_BLACK),
        ("• PhysiCausalNet IEEE TII (2024) — Cross-Machine FD", Pt(7.5), False, C_BLACK),
    ], Inches(0.6), Inches(6.02), Inches(6.1), Inches(0.85), line_spacing=1.1)

    # 우측 — AI 의사결정 자동화 수준 + 차별화
    add_box(sl, Inches(7.0), Inches(5.97), Inches(6.0), Inches(0.95),
            fill=C_BG, border=C_RED, border_w=Pt(1))
    add_text_lines(sl, [
        ("AI 실질 기여 — 의사결정 자동화 4단계", Pt(8.5), True, C_RED),
        ("①탐지 AI100%  ②원인진단 AI100%  ③처방추천 AI100%  ④최종조치 작업자", Pt(7.5), True, C_BLACK),
        ("차별화 4가지 (선행연구 대비)", Pt(8.5), True, C_RED),
        ("①공개KAMP데이터  ②비지도+KernelSHAP  ③Cross-Machine+CI  ④24개처방통합", Pt(7.5), False, C_BLACK),
    ], Inches(7.15), Inches(6.02), Inches(5.7), Inches(0.85), line_spacing=1.18)

    official_footer(sl)

# ══════════════════════════════════════════════════════════════
# Slide 6 — 5. 사용자 시나리오 / 유즈케이스 (User Case)
# ══════════════════════════════════════════════════════════════
def slide06_usecase(prs):
    sl = blank_slide(prs); fill_bg(sl)
    official_header(sl, 5, "사용자 시나리오 / 유즈케이스", "User Case")
    official_title(sl, 5, "사용자 시나리오 / 유즈케이스", "User Case")

    # ── 상단: 주요 사용자 정의 (3계층) ──
    add_text(sl, "[ 주요 사용자 정의 — 공장 운영 3계층 ]",
             Inches(0.45), Inches(1.4), Inches(12.5), Inches(0.3),
             size=Pt(11), bold=True, color=C_NAVY)
    users = [
        ("🛠 작업자", "Operator", [
            "사이드바 시나리오 선택 (정상/경고/위험/긴급)",
            "즉시 판정 + SHAP 처방 카드 확인",
            "조치 완료 / 모니터링 / 정비 요청 기록",
            "교대 인수인계 메모 작성",
        ], C_NAVY),
        ("👤 반장 / 관리자", "Shift Lead", [
            "이상 이력 50건 + 심각도 집계 박스 확인",
            "교대별 조치 현황 모니터링",
            "이상 이력 CSV 다운로드",
            "긴급 알람 시 작업자 지원·에스컬레이션",
        ], C_RED),
        ("🏢 부서장", "Department Head", [
            "일일 부서장 리포트 (Markdown) 자동 수신",
            "구간별 이상률 트렌드 차트",
            "외부 기계 검증 (CN7+RG3) 점검",
            "월 KPI 보고 (Recall, 비용 절감)",
        ], C_GREEN),
    ]
    for i, (icon_t, role, tasks, clr) in enumerate(users):
        x = Inches(0.45) + i * Inches(4.18)
        add_box(sl, x, Inches(1.72), Inches(4.05), Inches(2.18),
                fill=C_GRAY_BG, border=clr, border_w=Pt(1.2))
        add_box(sl, x, Inches(1.72), Inches(4.05), Inches(0.4), fill=clr)
        add_text(sl, icon_t, x + Inches(0.12), Inches(1.76),
                 Inches(2.5), Inches(0.32), size=Pt(11), bold=True, color=C_BG)
        add_text(sl, role, x + Inches(2.6), Inches(1.78),
                 Inches(1.4), Inches(0.3), size=Pt(8.5), color=C_BG, align=PP_ALIGN.RIGHT)
        for ti, t in enumerate(tasks):
            ty = Inches(2.22) + ti * Inches(0.38)
            add_text(sl, f"• {t}", x + Inches(0.15), ty,
                     Inches(3.85), Inches(0.36), size=Pt(8.5), color=C_BLACK)

    # ── 하단: 대표 시나리오 4개 ──
    add_text(sl, "[ 대표 시나리오 — 데모 시나리오 4종 (사이드바 즉시 적용) ]",
             Inches(0.45), Inches(4.0), Inches(12.5), Inches(0.3),
             size=Pt(11), bold=True, color=C_NAVY)
    scenarios = [
        ("S1", "정상 운영 (NORMAL)",  "모든 센서 정상 분포",     "정상 운영 유지", C_GREEN),
        ("S2", "경고 (WARNING)",     "Max_Back_Pressure +3.5σ",  "10분 내 재측정 권고", RGBColor(0xFF, 0xA5, 0x00)),
        ("S3", "위험 (DANGER)",      "Mold_Temperature_4 +4.2σ", "담당자 즉시 호출",   C_RED),
        ("S4", "긴급 (CRITICAL)",    "Filling_Time +15σ (게이트)", "라인 정지 검토",  RGBColor(0x8B, 0x00, 0x00)),
    ]
    for i, (sid, t, sym, action, clr) in enumerate(scenarios):
        x = Inches(0.45) + i * Inches(3.13)
        add_box(sl, x, Inches(4.32), Inches(3.0), Inches(2.6),
                fill=C_GRAY_BG, border=clr, border_w=Pt(1.5))
        # 헤더
        add_box(sl, x, Inches(4.32), Inches(3.0), Inches(0.45), fill=clr)
        add_text(sl, sid, x + Inches(0.12), Inches(4.36),
                 Inches(0.55), Inches(0.36), size=Pt(15), bold=True, color=C_BG)
        add_text(sl, t, x + Inches(0.65), Inches(4.4),
                 Inches(2.3), Inches(0.32), size=Pt(10), bold=True, color=C_BG)
        # 본문
        add_text(sl, "[ 입력 패턴 ]", x + Inches(0.15), Inches(4.92),
                 Inches(2.7), Inches(0.25), size=Pt(8), bold=True, color=C_NAVY)
        add_text(sl, sym, x + Inches(0.15), Inches(5.18),
                 Inches(2.7), Inches(0.5), size=Pt(8.5), color=C_BLACK)
        add_text(sl, "[ AI 응답 ]", x + Inches(0.15), Inches(5.7),
                 Inches(2.7), Inches(0.25), size=Pt(8), bold=True, color=C_NAVY)
        add_text(sl, action, x + Inches(0.15), Inches(5.95),
                 Inches(2.7), Inches(0.5), size=Pt(8.5), bold=True, color=clr)
        # 하단 사용자 흐름
        add_text(sl, "→ 작업자 즉시 처방 확인 → 조치 → 이력 누적",
                 x + Inches(0.15), Inches(6.5),
                 Inches(2.7), Inches(0.35), size=Pt(7.5), italic=True, color=C_DIM)

    official_footer(sl)

# ══════════════════════════════════════════════════════════════
# Slide 7 — 6. MVP 구현 범위 (MVP Scope)
# ══════════════════════════════════════════════════════════════
def slide07_mvp(prs):
    sl = blank_slide(prs); fill_bg(sl)
    official_header(sl, 6, "MVP 구현 범위", "MVP Scope")
    official_title(sl, 6, "MVP 구현 범위", "MVP Scope")

    # ── 상단: 핵심 메시지 + 데모 접속 ──
    add_box(sl, Inches(0.45), Inches(1.4), Inches(12.45), Inches(0.55),
            fill=C_NAVY, border=None)
    add_text(sl, "심사위원 즉시 시연 가능 — Streamlit 5탭 / 데모 시나리오 4종 / 시연 체크리스트 10개",
             Inches(0.6), Inches(1.45), Inches(8.5), Inches(0.45),
             size=Pt(11.5), bold=True, color=C_BG)
    add_text(sl, "데모 접속: streamlit run app.py → http://localhost:8501",
             Inches(9.1), Inches(1.5), Inches(3.7), Inches(0.4),
             size=Pt(8.5), color=C_BG, align=PP_ALIGN.RIGHT)

    # ══════════════════════════════════════════════════════════════
    # ── 좌측 (대형): 시연 가능 체크리스트 10개 ──
    # ══════════════════════════════════════════════════════════════
    add_box(sl, Inches(0.45), Inches(2.05), Inches(7.5), Inches(4.9),
            fill=C_GRAY_BG, border=C_GREEN, border_w=Pt(1.5))
    add_box(sl, Inches(0.45), Inches(2.05), Inches(7.5), Inches(0.42), fill=C_GREEN)
    add_text(sl, "✅  예선 완료 — 심사위원 시연 체크리스트 10개",
             Inches(0.6), Inches(2.1), Inches(7.2), Inches(0.32),
             size=Pt(12), bold=True, color=C_BG)

    demo_checklist = [
        ("1", "AI 자동 이상탐지",  f"ROC-AUC {ROC:.3f}, F1 {F1:.3f}", "Tab 1"),
        ("2", "24센서 슬라이더 즉시 판정", "복원 오차 실시간 계산", "Tab 2"),
        ("3", "데모 시나리오 4종 1클릭", "정상/경고/위험/긴급", "사이드바"),
        ("4", "SHAP 24센서 기여도", "KernelSHAP Waterfall", "Tab 2 / Tab 4"),
        ("5", "처방 카드 자동 출력", "24개 센서별 조작/정비 구분", "Tab 2"),
        ("6", "심각도 3단계 알람",  "경고·위험·긴급 + 에스컬레이션", "Tab 2"),
        ("7", "이상 이력 50건 + CSV", "JSON 영속 + 인수인계 메모", "Tab 2 하단"),
        ("8", "일일 부서장 리포트", "Markdown 자동 생성", "Tab 2 하단"),
        ("9", "Bootstrap 95% CI 평가", "1,000회 + Cross-Machine 검증", "Tab 1"),
        ("10", "795K 일괄 스코어링", "비라벨 데이터 트렌드 분석", "Tab 3"),
    ]
    # 2열 5행 그리드
    for i, (num, t, b, loc) in enumerate(demo_checklist):
        col, row = i % 2, i // 2
        x = Inches(0.55) + col * Inches(3.7)
        y = Inches(2.58) + row * Inches(0.87)
        add_box(sl, x, y, Inches(3.55), Inches(0.78), fill=C_BG, border=C_GRAY_BD)
        # 체크 번호
        add_box(sl, x, y, Inches(0.4), Inches(0.78), fill=C_GREEN)
        add_text(sl, num, x, y + Inches(0.23), Inches(0.4), Inches(0.32),
                 size=Pt(14), bold=True, color=C_BG, align=PP_ALIGN.CENTER)
        # 제목
        add_text(sl, t, x + Inches(0.48), y + Inches(0.04),
                 Inches(2.4), Inches(0.28), size=Pt(9.5), bold=True, color=C_NAVY)
        # 설명
        add_text(sl, b, x + Inches(0.48), y + Inches(0.3),
                 Inches(2.4), Inches(0.26), size=Pt(7.5), color=C_BLACK)
        # 위치 (탭)
        add_box(sl, x + Inches(2.9), y + Inches(0.5), Inches(0.6), Inches(0.22),
                fill=C_BLUE_LT, border=C_NAVY, border_w=Pt(0.5))
        add_text(sl, loc, x + Inches(2.9), y + Inches(0.5),
                 Inches(0.6), Inches(0.22),
                 size=Pt(7), bold=True, color=C_NAVY, align=PP_ALIGN.CENTER)

    # ══════════════════════════════════════════════════════════════
    # ── 우측 (소): 본선 구현 P1~P5 ──
    # ══════════════════════════════════════════════════════════════
    add_box(sl, Inches(8.1), Inches(2.05), Inches(4.85), Inches(4.9),
            fill=C_GRAY_BG, border=C_RED, border_w=Pt(1.5))
    add_box(sl, Inches(8.1), Inches(2.05), Inches(4.85), Inches(0.42), fill=C_RED)
    add_text(sl, "🎯  본선 구현 우선순위 P1~P5",
             Inches(8.25), Inches(2.1), Inches(4.6), Inches(0.32),
             size=Pt(12), bold=True, color=C_BG)

    todo_items = [
        ("P1", "OPC-UA/MQTT 실시간 연동", "PLC 직접 데이터 수신", "1.5개월"),
        ("P2", "알람 자동화 (SMS/이메일)", "긴급 시 자동 발송",   "2주"),
        ("P3", "멀티 설비 모델 관리",     "설비 ID별 임계값",     "1개월"),
        ("P4", "예측 정비 모듈",          "고장 사전 예측",       "1.5개월"),
        ("P5", "AWS 클라우드 옵션",       "온프레미스+SaaS",      "1개월"),
    ]
    y = Inches(2.6)
    for pid, title, desc, dur in todo_items:
        add_box(sl, Inches(8.25), y, Inches(4.55), Inches(0.82),
                fill=C_BG, border=C_GRAY_BD)
        # 우선순위
        add_box(sl, Inches(8.25), y, Inches(0.45), Inches(0.82), fill=C_RED)
        add_text(sl, pid, Inches(8.25), y + Inches(0.26),
                 Inches(0.45), Inches(0.32),
                 size=Pt(12), bold=True, color=C_BG, align=PP_ALIGN.CENTER)
        # 본문
        add_text(sl, title, Inches(8.8), y + Inches(0.04),
                 Inches(3.0), Inches(0.3), size=Pt(9.5), bold=True, color=C_NAVY)
        add_text(sl, desc, Inches(8.8), y + Inches(0.3),
                 Inches(3.0), Inches(0.28), size=Pt(8), color=C_BLACK)
        # 기간
        add_box(sl, Inches(11.85), y + Inches(0.25), Inches(0.9), Inches(0.3),
                fill=C_BLUE_LT, border=C_NAVY)
        add_text(sl, dur, Inches(11.85), y + Inches(0.25),
                 Inches(0.9), Inches(0.3),
                 size=Pt(8), bold=True, color=C_NAVY, align=PP_ALIGN.CENTER)
        y += Inches(0.88)

    # ══════════════════════════════════════════════════════════════
    # ── 하단: 재현 가이드 + 데모 접속 안내 ──
    # ══════════════════════════════════════════════════════════════
    add_box(sl, Inches(0.45), Inches(7.0), Inches(12.45), Inches(0.0))  # 자리만
    add_text(sl,
             "재현 가이드: install.bat 원클릭  ·  REPRODUCE.md 7단계  ·  모델 SHA-256 무결성 검증 (verify_model.py)  ·  학술 레퍼런스 7개 인용",
             Inches(0.45), Inches(6.99), Inches(12.45), Inches(0.18),
             size=Pt(7.5), italic=True, color=C_DIM, align=PP_ALIGN.CENTER)

    official_footer(sl)

# ══════════════════════════════════════════════════════════════
# Slide 8 — 7. 기대 효과 및 향후 확장성 (Expected Impact)
# ══════════════════════════════════════════════════════════════
def slide08_impact(prs):
    sl = blank_slide(prs); fill_bg(sl)
    official_header(sl, 7, "기대 효과 및 향후 확장성", "Expected Impact")
    official_title(sl, 7, "기대 효과 및 향후 확장성", "Expected Impact")

    # ── 상단: 정량 효과 4개 KPI 카드 ──
    add_text(sl, "[ 정량 기대 효과 — 도입 첫해 기준 ]",
             Inches(0.45), Inches(1.4), Inches(12.5), Inches(0.3),
             size=Pt(11), bold=True, color=C_NAVY)
    quant = [
        ("Recall 0.667", "월 25건 중\n약 17건 자동 탐지", "불량 조기 인지", C_RED),
        ("연 1억원 절감", "200건/년 ×\n50만원/건",        "불량 비용 회피", C_GREEN),
        ("응답 즉시", "수 시간 →\n수 초",                 "원인 분석 단축",  C_NAVY),
        ("ROI 4.2배", "연 절감 1억 ÷\n월 200만원 SaaS", "투자 회수율",   C_RED),
    ]
    for i, (big, val, sub, clr) in enumerate(quant):
        x = Inches(0.45) + i * Inches(3.13)
        add_box(sl, x, Inches(1.72), Inches(3.0), Inches(1.25),
                fill=C_GRAY_BG, border=clr, border_w=Pt(1.5))
        add_text(sl, big, x + Inches(0.1), Inches(1.78),
                 Inches(2.8), Inches(0.4), size=Pt(15), bold=True, color=clr)
        add_text(sl, val, x + Inches(0.1), Inches(2.18),
                 Inches(2.8), Inches(0.5), size=Pt(9), color=C_BLACK)
        add_text(sl, sub, x + Inches(0.1), Inches(2.7),
                 Inches(2.8), Inches(0.25), size=Pt(8.5), italic=True, color=C_DIM)

    # ── 중단 좌측: 정성 효과 + 현장 적용성 ──
    add_text(sl, "[ 정성 효과 + 현장 적용 가능성 ]",
             Inches(0.45), Inches(3.1), Inches(6.4), Inches(0.3),
             size=Pt(11), bold=True, color=C_NAVY)
    add_box(sl, Inches(0.45), Inches(3.42), Inches(6.4), Inches(2.0),
            fill=C_BLUE_LT, border=C_NAVY, border_w=Pt(1))
    add_text_lines(sl, [
        ("정성 효과", Pt(10), True, C_NAVY),
        ("• 노하우 디지털화 — 담당자 경험 → SHAP 데이터 누적", Pt(8.5), False, C_BLACK),
        ("• 교대 인수인계 표준화 — 수기 일지 → CSV/MD 자동", Pt(8.5), False, C_BLACK),
        ("• ISO 9001 이력 보존 — JSON 영속 + 무결성 해시", Pt(8.5), False, C_BLACK),
        ("", Pt(4), False, C_BLACK),
        ("현장 적용 가능성", Pt(10), True, C_NAVY),
        ("• 온프레미스 배포 — 공정 데이터 외부 전송 없음", Pt(8.5), False, C_BLACK),
        ("• 8GB RAM 보급형 PC 동작 (CPU 전용 가능)", Pt(8.5), False, C_BLACK),
        ("• install.bat 원클릭 설치 — IT 인력 부담 최소화", Pt(8.5), False, C_BLACK),
    ], Inches(0.6), Inches(3.5), Inches(6.1), Inches(1.85), line_spacing=1.18)

    # ── 중단 우측: 확장성 (수직/수평) ──
    add_text(sl, "[ 향후 고도화 / 확장성 ]",
             Inches(7.0), Inches(3.1), Inches(6.0), Inches(0.3),
             size=Pt(11), bold=True, color=C_NAVY)
    add_box(sl, Inches(7.0), Inches(3.42), Inches(6.0), Inches(2.0),
            fill=C_RED_LT, border=C_RED, border_w=Pt(1))
    add_text_lines(sl, [
        ("도메인 수평 확장 (사출 → 타 공정)", Pt(10), True, C_RED),
        ("• 다이캐스팅, CNC, 프레스 — 24센서 인터페이스 표준화", Pt(8.5), False, C_BLACK),
        ("• KAMP 12종 데이터셋 → 공정별 모델 카탈로그", Pt(8.5), False, C_BLACK),
        ("", Pt(4), False, C_BLACK),
        ("수직 확장 (기능 고도화)", Pt(10), True, C_RED),
        ("• Autoencoder → Transformer 마이그레이션 로드맵", Pt(8.5), False, C_BLACK),
        ("• 예측 정비 (Predictive Maintenance) 모듈", Pt(8.5), False, C_BLACK),
        ("• 강화학습 기반 자동 처방 (장기)", Pt(8.5), False, C_BLACK),
    ], Inches(7.15), Inches(3.5), Inches(5.7), Inches(1.85), line_spacing=1.18)

    # ── 하단: 사업화 로드맵 (6개월 마일스톤) ──
    add_text(sl, "[ 사업화 로드맵 — 6개월 마일스톤 (KAMP + 정부 정책 연계) ]",
             Inches(0.45), Inches(5.55), Inches(12.5), Inches(0.3),
             size=Pt(11), bold=True, color=C_NAVY)
    miles = [
        ("M1~M2", "예선 통과 + 본선 구현", "OPC-UA 연동, 알람 자동화 완료", C_NAVY),
        ("M3", "PoC 3건 시작", "KAMP 파트너 + 정부 보급사업 연계", C_NAVY),
        ("M4", "유료 전환 2건", "MRR 400만원 / Recall 0.75+", C_RED),
        ("M5", "Series A 준비", "ARR 6,000만 / 유료 5개사 / NPS 40+", C_RED),
        ("M6", "Series A 클로징", "정부 정책 만기 (3만개사) 진입 적기", C_GREEN),
    ]
    for i, (m, t, b, clr) in enumerate(miles):
        x = Inches(0.45) + i * Inches(2.52)
        add_box(sl, x, Inches(5.87), Inches(2.42), Inches(0.92),
                fill=C_GRAY_BG, border=clr, border_w=Pt(1))
        add_box(sl, x, Inches(5.87), Inches(2.42), Inches(0.28), fill=clr)
        add_text(sl, m, x, Inches(5.87), Inches(2.42), Inches(0.28),
                 size=Pt(10), bold=True, color=C_BG, align=PP_ALIGN.CENTER)
        add_text(sl, t, x + Inches(0.1), Inches(6.22),
                 Inches(2.22), Inches(0.3), size=Pt(9.5), bold=True, color=C_BLACK)
        add_text(sl, b, x + Inches(0.1), Inches(6.5),
                 Inches(2.22), Inches(0.4), size=Pt(7.5), color=C_DIM)
        # 화살표
        if i < len(miles) - 1:
            add_text(sl, "→", x + Inches(2.3), Inches(6.12),
                     Inches(0.3), Inches(0.4),
                     size=Pt(14), bold=True, color=C_DIM)

    # ── 발표 시간 분배 + 예상 Q&A 키워드 (발표/전달력 평가 반영) ──
    add_box(sl, Inches(0.45), Inches(6.86), Inches(12.45), Inches(0.18),
            fill=C_GRAY_BG, border=C_NAVY, border_w=Pt(0.5))
    add_text(sl,
             "발표 6분 — 표지15s·문제45s·솔루션45s·기능60s·데이터60s·유즈45s·MVP45s·효과60s  |  Q&A: ①0.58%불균형→반지도 ②외부검증→CN7+RG3 ③사업화→KAMP+정부정책",
             Inches(0.55), Inches(6.88), Inches(12.3), Inches(0.16),
             size=Pt(7), bold=True, color=C_NAVY, align=PP_ALIGN.CENTER)

    official_footer(sl)

# ══════════════════════════════════════════════════════════════
# 메인
# ══════════════════════════════════════════════════════════════
def main():
    print("[공식 양식] PPT 생성 시작...")
    prs = new_prs()

    print("  Slide 1/8: 표지")
    slide01_cover(prs)
    print("  Slide 2/8: 1. 문제 정의")
    slide02_problem(prs)
    print("  Slide 3/8: 2. 솔루션 개요")
    slide03_solution(prs)
    print("  Slide 4/8: 3. 주요 기능")
    slide04_features(prs)
    print("  Slide 5/8: 4. 데이터 및 기술")
    slide05_data_tech(prs)
    print("  Slide 6/8: 5. 사용자 시나리오")
    slide06_usecase(prs)
    print("  Slide 7/8: 6. MVP 구현 범위")
    slide07_mvp(prs)
    print("  Slide 8/8: 7. 기대 효과 및 향후 확장성")
    slide08_impact(prs)

    prs.save(OUT_PATH)
    print(f"\n[완료] 저장: {OUT_PATH}")
    print(f"   슬라이드 수: {len(prs.slides)}")
    print(f"\n   모델 성능 반영:")
    print(f"   ROC-AUC  = {ROC:.4f}")
    print(f"   F1-Score = {F1:.4f}")
    print(f"   Recall   = {REC:.4f}")
    if SHAP_TOP:
        print(f"\n   SHAP Top-3: {', '.join(f'{c}({v:.2f})' for c,v in SHAP_TOP[:3])}")

if __name__ == '__main__':
    main()
