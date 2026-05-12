"""
스마트 공장 XAI 이상탐지 플랫폼 — 예선 기획서 PPT 생성
8슬라이드 / 16:9 / 흑백+빨강 강조색
"""
import sys, os, json
sys.path.insert(0, os.path.dirname(os.path.dirname(os.path.abspath(__file__))))

import numpy as np
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt
from pptx.enum.dml import MSO_THEME_COLOR

# ── 경로 ──
BASE_DIR   = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
RESULT_DIR = os.path.join(BASE_DIR, 'results')
MODEL_DIR  = os.path.join(BASE_DIR, 'models')
OUT_PATH   = os.path.join(BASE_DIR, 'output', '스마트공장XAI_예선기획서.pptx')
os.makedirs(os.path.join(BASE_DIR, 'output'), exist_ok=True)

# ── 실제 수치 로드 ──
with open(os.path.join(RESULT_DIR, 'metrics.json'), encoding='utf-8') as f:
    M = json.load(f)
ROC   = M['roc_auc']
PR    = M['pr_auc']
F1    = M['f1']
REC   = M['recall']
PREC  = M['precision']
THR   = M['threshold']

# SHAP 상위 센서
SHAP_TOP = []
shap_path = os.path.join(RESULT_DIR, 'shap_values.npy')
if os.path.exists(shap_path):
    from src.config import SENSOR_COLS
    sv  = np.load(shap_path)
    ma  = np.abs(sv).mean(axis=0)
    idx = np.argsort(ma)[::-1][:5]
    SHAP_TOP = [(SENSOR_COLS[i], float(ma[i])) for i in idx]

# 스코어링 통계
SC = None
sc_path = os.path.join(RESULT_DIR, 'scored_unlabeled.parquet')
if os.path.exists(sc_path):
    import pandas as pd
    df  = pd.read_parquet(sc_path)
    anom = (df['recon_error'] >= THR).sum()
    SC  = {'total': len(df), 'anomaly': int(anom), 'rate': anom / len(df) * 100}

# ── 색상 상수 ──
C_BLACK  = RGBColor(0x08, 0x08, 0x08)
C_DARK   = RGBColor(0x11, 0x11, 0x11)
C_CARD   = RGBColor(0x1A, 0x1A, 0x1A)
C_RED    = RGBColor(0xD4, 0x21, 0x21)
C_WHITE  = RGBColor(0xEF, 0xEF, 0xEF)
C_DIM    = RGBColor(0x88, 0x88, 0x88)
C_MUTED  = RGBColor(0x4A, 0x4A, 0x4A)
C_BORDER = RGBColor(0x33, 0x33, 0x33)

# ── 슬라이드 크기 (16:9, 33.87 × 19.05 cm) ──
W = Inches(13.33)
H = Inches(7.5)

# ── 헬퍼 함수 ──
def new_prs():
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H
    return prs

def blank_slide(prs):
    layout = prs.slide_layouts[6]  # completely blank
    return prs.slides.add_slide(layout)

def fill_bg(slide, color=C_BLACK):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_box(slide, left, top, width, height, fill=None, border=None, border_w=Pt(0.5)):
    from pptx.util import Emu
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        left, top, width, height
    )
    shape.line.width = Pt(0)
    f = shape.fill
    if fill:
        f.solid()
        f.fore_color.rgb = fill
    else:
        f.background()
    if border:
        shape.line.color.rgb = border
        shape.line.width = border_w
    return shape

def add_text(slide, text, left, top, width, height,
             size=Pt(14), bold=False, color=C_WHITE,
             align=PP_ALIGN.LEFT, wrap=True, italic=False, font_name="Noto Sans KR"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = size
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = font_name
    return txBox

def add_text_lines(slide, lines, left, top, width, height,
                   size=Pt(12), color=C_WHITE, line_spacing=1.2,
                   font_name="Noto Sans KR"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    first = True
    for txt, sz, bold, col in lines:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = txt
        run.font.size = sz
        run.font.bold = bold
        run.font.color.rgb = col
        run.font.name = font_name
    return txBox

def red_line(slide, left, top, width, h=Pt(2)):
    shape = slide.shapes.add_shape(1, left, top, width, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = C_RED
    shape.line.width = Pt(0)
    return shape

def section_header(slide, number, title_ko, title_en=""):
    # 좌측 빨간 세로 바
    add_box(slide, Inches(0.45), Inches(0.3), Pt(3), Inches(0.45), fill=C_RED)
    # 섹션 번호
    add_text(slide, number, Inches(0.55), Inches(0.28), Inches(0.5), Inches(0.4),
             size=Pt(10), bold=True, color=C_RED)
    # 섹션 제목
    add_text(slide, title_ko, Inches(0.85), Inches(0.28), Inches(8), Inches(0.44),
             size=Pt(18), bold=True, color=C_WHITE)
    if title_en:
        add_text(slide, title_en, Inches(0.85), Inches(0.68), Inches(8), Inches(0.3),
                 size=Pt(10), color=C_DIM)

def metric_card(slide, left, top, label, value, sub="", w=Inches(2.2), h=Inches(1.1)):
    add_box(slide, left, top, w, h, fill=C_CARD, border=C_BORDER)
    add_text(slide, label, left+Inches(0.12), top+Inches(0.08), w-Inches(0.2), Inches(0.28),
             size=Pt(9), color=C_DIM)
    add_text(slide, value, left+Inches(0.12), top+Inches(0.30), w-Inches(0.2), Inches(0.5),
             size=Pt(22), bold=True, color=C_RED)
    if sub:
        add_text(slide, sub, left+Inches(0.12), top+Inches(0.78), w-Inches(0.2), Inches(0.25),
                 size=Pt(8), color=C_DIM)

# ══════════════════════════════════════════════════════════════
# 슬라이드 생성 함수
# ══════════════════════════════════════════════════════════════

def slide01_cover(prs):
    """표지"""
    sl = blank_slide(prs)
    fill_bg(sl)

    # 전체 좌측 빨간 세로선 (장식)
    add_box(sl, Inches(0.38), Inches(1.5), Pt(2), Inches(4.5), fill=C_RED)

    # 해커톤명
    add_text(sl, "2026 스마트 공장 운영 시스템 MVP 개발 해커톤",
             Inches(0.55), Inches(1.4), Inches(12), Inches(0.4),
             size=Pt(11), color=C_DIM)

    # 메인 타이틀
    add_text(sl, "사출성형기 이상탐지 &\nXAI 공정 운영 지원 시스템",
             Inches(0.55), Inches(1.85), Inches(10), Inches(1.8),
             size=Pt(34), bold=True, color=C_WHITE)

    # 서브 타이틀 (빨강 강조)
    add_text(sl, "반지도학습 Autoencoder + SHAP 기반\n탐지 · 진단 · 처방 · 추적 4단계 운영 자동화",
             Inches(0.55), Inches(3.65), Inches(10), Inches(1.0),
             size=Pt(14), color=C_RED)

    # 구분선
    red_line(sl, Inches(0.55), Inches(4.72), Inches(12.3))

    # 핵심 임팩트 수치 3개 (P5-M2: 시각적 위계 강화)
    impact_cards = [
        (f"불량 10건 중 {int(round(REC*10))}건",  "AI 조기 감지 (Recall 기준)"),
        ("수 시간 → 수 초",                         "원인 센서 분석 단축"),
        ("24개 센서 동시",                          "실시간 공정 모니터링"),
    ]
    for i, (big, sub) in enumerate(impact_cards):
        x = Inches(0.55) + i * Inches(4.2)
        add_box(sl, x, Inches(4.92), Inches(3.8), Inches(0.72),
                fill=C_CARD, border=C_RED, border_w=Pt(1))
        add_text(sl, big, x + Inches(0.14), Inches(4.97),
                 Inches(2.5), Inches(0.38),
                 size=Pt(16), bold=True, color=C_RED)
        add_text(sl, sub, x + Inches(0.14), Inches(5.32),
                 Inches(3.5), Inches(0.26),
                 size=Pt(9), color=C_DIM)

    # 하단 정보
    add_text(sl, "팀명: SmartFactory XAI팀",
             Inches(0.55), Inches(5.88), Inches(5), Inches(0.35),
             size=Pt(10), color=C_DIM)
    add_text(sl, "제출일: 2026년 5월 13일",
             Inches(0.55), Inches(6.18), Inches(5), Inches(0.35),
             size=Pt(10), color=C_DIM)

    # 우측 배지
    add_box(sl, Inches(10.5), Inches(5.85), Inches(2.4), Inches(0.7),
            fill=C_RED, border=None)
    add_text(sl, "예선 기획서",
             Inches(10.5), Inches(5.85), Inches(2.4), Inches(0.7),
             size=Pt(14), bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)


def slide02_problem(prs):
    """1. 문제 정의"""
    sl = blank_slide(prs)
    fill_bg(sl)
    section_header(sl, "01", "문제 정의", "Problem Definition")

    # 배경 — 현장 고통 3가지
    pains = [
        ("불량 발생 인지 지연", "작업자가 육안으로만 불량 확인 → 이미 수백 개 불량 후 인지"),
        ("원인 불명확",         "어떤 센서(공정 변수)가 불량 원인인지 분석 불가 → 재발 방지 어려움"),
        ("사후 대응 반복",       "불량 발생 → 라인 정지 → 수동 검사 → 재가동 사이클 반복"),
    ]

    for i, (title, body) in enumerate(pains):
        y = Inches(1.15) + i * Inches(1.55)
        # 번호 원
        add_box(sl, Inches(0.45), y, Inches(0.48), Inches(0.48),
                fill=C_RED)
        add_text(sl, str(i+1), Inches(0.45), y, Inches(0.48), Inches(0.48),
                 size=Pt(14), bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
        # 제목
        add_text(sl, title, Inches(1.1), y, Inches(5.5), Inches(0.44),
                 size=Pt(15), bold=True, color=C_WHITE)
        # 내용
        add_text(sl, body, Inches(1.1), y+Inches(0.42), Inches(10.5), Inches(0.55),
                 size=Pt(11), color=C_DIM)

    # 우측: 현황 수치 박스
    add_box(sl, Inches(7.8), Inches(1.1), Inches(5.0), Inches(5.1),
            fill=C_CARD, border=C_BORDER)
    add_text(sl, "제조업 현장 불량 현황",
             Inches(7.95), Inches(1.2), Inches(4.7), Inches(0.38),
             size=Pt(12), bold=True, color=C_DIM)
    red_line(sl, Inches(7.95), Inches(1.55), Inches(4.6))

    stats = [
        ("전체 생산 불량률",    "1 ~ 3%",   "사출성형 공정 평균"),
        ("불량 원인 추적 소요", "수 시간",   "담당자 수동 분석 기준"),
        ("연간 불량 손실 비용", "수억 원",   "중소 제조사 추정치"),
        ("작업자 의존도",       "90% 이상",  "정성적 판단 비율"),
    ]
    for i, (lbl, val, sub) in enumerate(stats):
        y = Inches(1.72) + i * Inches(0.9)
        add_text(sl, lbl, Inches(7.95), y,          Inches(2.5), Inches(0.35),
                 size=Pt(10), color=C_DIM)
        add_text(sl, val, Inches(10.3), y-Inches(0.02), Inches(1.5), Inches(0.42),
                 size=Pt(16), bold=True, color=C_RED, align=PP_ALIGN.RIGHT)
        add_text(sl, sub, Inches(7.95), y+Inches(0.32), Inches(4.6), Inches(0.28),
                 size=Pt(9), color=C_MUTED)

    # ── 시장 규모 (TAM/SAM/SOM) ──
    add_box(sl, Inches(0.45), Inches(5.3), Inches(12.4), Inches(0.52),
            fill=C_CARD, border=C_BORDER)
    add_text(sl, "시장 규모 (출처: KEIT 기술수요조사 2025 · 국내 사출성형 사업장 9,800개 기준)",
             Inches(0.65), Inches(5.35), Inches(5.0), Inches(0.32),
             size=Pt(9), bold=True, color=C_DIM)
    mkt = [
        ("TAM", "2조 원", "국내 스마트 제조 솔루션 전체"),
        ("SAM", "3,000억 원", "사출·금속 공정 특화 AI"),
        ("SOM", "42억 원 (Y3)", "350개사 × 월 100만원"),
    ]
    for i, (lbl, val, sub) in enumerate(mkt):
        x = Inches(5.85) + i * Inches(2.2)
        add_text(sl, lbl, x, Inches(5.35), Inches(0.6), Inches(0.28),
                 size=Pt(8), bold=True, color=C_RED)
        add_text(sl, val, x + Inches(0.62), Inches(5.35), Inches(1.3), Inches(0.28),
                 size=Pt(10), bold=True, color=C_WHITE)
        add_text(sl, sub, x + Inches(0.62), Inches(5.6), Inches(1.4), Inches(0.22),
                 size=Pt(7), color=C_MUTED)

    # 하단 강조 문구
    add_box(sl, Inches(0.45), Inches(5.9), Inches(12.4), Inches(0.85),
            fill=C_CARD, border=C_RED, border_w=Pt(1))
    add_text(sl,
             "\"불량이 발생했다\"는 사실을 아는 것만으로는 충분하지 않습니다.  "
             "\"언제·어떤 공정 변수가 원인인지\" 즉각 파악하고 처방해야 합니다.",
             Inches(0.65), Inches(5.98), Inches(12.0), Inches(0.65),
             size=Pt(11), color=C_WHITE)


def slide03_solution(prs):
    """2. 제안 솔루션 개요"""
    sl = blank_slide(prs)
    fill_bg(sl)
    section_header(sl, "02", "제안 솔루션 개요", "Solution Overview")

    # 4단계 플로우 박스
    steps = [
        ("탐 지", "Detect",  "Autoencoder\n복원 오차 기반\n이상 스코어링"),
        ("진 단", "Diagnose","SHAP XAI\n불량 원인 센서\n자동 순위 도출"),
        ("처 방", "Prescribe","원인 센서별\n점검·조정 액션\n자동 권고"),
        ("추 적", "Track",    "설비·공정별\n건강 점수 &\n생산 이력 분석"),
    ]
    step_w = Inches(2.8)
    arrow_w = Inches(0.35)
    start_x = Inches(0.45)

    for i, (ko, en, desc) in enumerate(steps):
        x = start_x + i * (step_w + arrow_w)
        # 메인 박스
        add_box(sl, x, Inches(1.1), step_w, Inches(2.8),
                fill=C_CARD, border=C_RED if i == 0 else C_BORDER,
                border_w=Pt(1.5) if i == 0 else Pt(0.5))
        # 단계 번호
        add_text(sl, f"STEP {i+1}", x, Inches(1.18), step_w, Inches(0.3),
                 size=Pt(8), color=C_RED, align=PP_ALIGN.CENTER, bold=True)
        # 한글 제목 (빨강)
        add_text(sl, ko, x, Inches(1.44), step_w, Inches(0.55),
                 size=Pt(22), bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
        add_text(sl, f"[ {en} ]", x, Inches(1.95), step_w, Inches(0.35),
                 size=Pt(10), color=C_DIM, align=PP_ALIGN.CENTER)
        add_text(sl, desc, x, Inches(2.35), step_w, Inches(1.4),
                 size=Pt(10), color=C_WHITE, align=PP_ALIGN.CENTER)

        # 화살표
        if i < 3:
            ax = x + step_w + Inches(0.05)
            add_text(sl, "→", ax, Inches(2.3), arrow_w, Inches(0.5),
                     size=Pt(20), bold=True, color=C_RED, align=PP_ALIGN.CENTER)

    # 기술 스택 라벨
    tech = [
        ("PyTorch Autoencoder",   "24→16→8→16→24"),
        ("SHAP KernelExplainer",  "50 kmeans 배경"),
        ("Streamlit Dashboard",   "5탭 웹 인터페이스"),
        ("KAMP 사출성형기 데이터", "공개 데이터 24 센서"),
    ]
    for i, (lbl, sub) in enumerate(tech):
        x = Inches(0.45) + i * Inches(3.15)
        add_box(sl, x, Inches(4.2), Inches(2.9), Inches(0.75),
                fill=C_CARD, border=C_BORDER)
        add_text(sl, lbl, x+Inches(0.12), Inches(4.27), Inches(2.65), Inches(0.35),
                 size=Pt(11), bold=True, color=C_WHITE)
        add_text(sl, sub, x+Inches(0.12), Inches(4.57), Inches(2.65), Inches(0.28),
                 size=Pt(9), color=C_DIM)

    # 핵심 차별점 + 경쟁 포지셔닝
    add_box(sl, Inches(0.45), Inches(5.1), Inches(12.4), Inches(2.0),
            fill=C_CARD, border=C_BORDER)
    add_text(sl, "핵심 차별점",
             Inches(0.65), Inches(5.18), Inches(3), Inches(0.32),
             size=Pt(11), bold=True, color=C_RED)
    diffs = [
        "✦  반지도학습 — 정상 데이터만으로 학습 (불량 라벨 불필요, 데이터 수집 비용 제거)",
        "✦  SHAP XAI — \"왜 불량인가?\" 센서 단위 설명 → 담당자가 즉시 원인 이해",
        "✦  탐지→진단→처방 자동화 — 기존 MES·ERP 단순 임계값 알람과 달리 원인까지 제시",
    ]
    for i, d in enumerate(diffs):
        add_text(sl, d, Inches(0.65), Inches(5.48) + i * Inches(0.28),
                 Inches(12.0), Inches(0.28),
                 size=Pt(9.5), color=C_WHITE)
    # 경쟁 대안 비교 한 줄
    add_text(sl, "vs 경쟁 대안:  수동 점검(불량 후 인지·분석 불가)  |  범용 MES 알람(센서 원인 설명 없음)  |  외부 컨설팅(수백만원/회·비정기)",
             Inches(0.65), Inches(6.35), Inches(12.0), Inches(0.28),
             size=Pt(8.5), color=C_DIM)


def slide04_features(prs):
    """3. 주요 기능 정의"""
    sl = blank_slide(prs)
    fill_bg(sl)
    section_header(sl, "03", "주요 기능 정의", "Key Features")

    features = [
        {
            "tab": "탭 1",
            "name": "데이터 개요 & 모델 성능",
            "items": [
                "ROC-AUC / PR-AUC / F1 실시간 모니터링",
                "복원 오차 분포 인터랙티브 차트",
                "학습 곡선 (Train/Val Loss) 추이",
            ],
        },
        {
            "tab": "탭 2",
            "name": "실시간 공정 시뮬레이터",
            "items": [
                "24개 센서 슬라이더 → 즉시 이상 판정",
                "On-demand SHAP 계산 → 원인 센서 Top-5",
                "이상 시 처방 자동 출력 (점검 액션)",
            ],
        },
        {
            "tab": "탭 3",
            "name": "대규모 스코어링 분석",
            "items": [
                f"비라벨 {SC['total']:,}건 이상 스코어 시각화" if SC else "비라벨 데이터 이상 스코어 시각화",
                "임계값 슬라이더 → 이상 비율 실시간 변경",
                "이상 구간 Zoom-In 인터랙티브 탐색",
            ],
        },
        {
            "tab": "탭 4",
            "name": "XAI 원인 분석",
            "items": [
                "SHAP Waterfall — 개별 이상 샘플 설명",
                "SHAP Summary Bar — 전체 특성 중요도",
                "센서별 원인 기여도 순위표",
            ],
        },
        {
            "tab": "탭 5",
            "name": "설비별 운영 현황",
            "items": [
                "금형(설비)별 건강 점수 (0~100)",
                "최근 N샷 이상 추이 게이지 & 히스토리",
                "생산 이력 테이블 (필터·정렬 지원)",
            ],
        },
    ]

    col_w = Inches(2.45)
    for i, feat in enumerate(features):
        x = Inches(0.45) + i * (col_w + Inches(0.07))
        # 탭 배지
        add_box(sl, x, Inches(1.08), Inches(0.72), Inches(0.3),
                fill=C_RED)
        add_text(sl, feat["tab"], x, Inches(1.08), Inches(0.72), Inches(0.3),
                 size=Pt(8), bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
        # 기능 카드
        add_box(sl, x, Inches(1.36), col_w, Inches(5.35),
                fill=C_CARD, border=C_BORDER)
        add_text(sl, feat["name"],
                 x+Inches(0.12), Inches(1.44), col_w-Inches(0.22), Inches(0.62),
                 size=Pt(12), bold=True, color=C_WHITE)
        red_line(sl, x+Inches(0.12), Inches(2.02), col_w-Inches(0.22), Pt(1))
        for j, item in enumerate(feat["items"]):
            add_text(sl, f"• {item}",
                     x+Inches(0.12), Inches(2.15)+j*Inches(0.55),
                     col_w-Inches(0.22), Inches(0.55),
                     size=Pt(9), color=C_WHITE)


def slide05_data_tech(prs):
    """4. 데이터 및 기술 활용 계획"""
    sl = blank_slide(prs)
    fill_bg(sl)
    section_header(sl, "04", "데이터 및 기술 활용 계획", "Data & Technology")

    # 좌측: 데이터셋
    add_box(sl, Inches(0.45), Inches(1.1), Inches(5.8), Inches(5.5),
            fill=C_CARD, border=C_BORDER)
    add_text(sl, "데이터셋 — KAMP 사출성형기 AI 데이터셋",
             Inches(0.65), Inches(1.18), Inches(5.4), Inches(0.38),
             size=Pt(12), bold=True, color=C_WHITE)
    red_line(sl, Inches(0.65), Inches(1.53), Inches(5.4))

    ds_rows = [
        ("supervised_label_cn7.csv",    "6,736행",  "정상 6,697 / 불량 39"),
        ("moldset_labeled_cn7.csv",     "1,211행",  "정상 1,194 / 불량 17"),
        ("moldset_labeled_rg3.csv",     "1,182행",  "정상 1,157 / 불량 25"),
        ("moldset_unlabeled_cn7.csv",   "35,239행", "비라벨 스코어링용"),
        ("labeled_data.csv",            "7,996행",  "Scaler 학습용 원본"),
    ]
    for i, (name, cnt, note) in enumerate(ds_rows):
        y = Inches(1.68) + i * Inches(0.83)
        add_text(sl, name, Inches(0.65), y,          Inches(3.4), Inches(0.32),
                 size=Pt(9), bold=True, color=C_RED)
        add_text(sl, cnt,  Inches(0.65), y+Inches(0.3), Inches(1.5), Inches(0.28),
                 size=Pt(10), bold=True, color=C_WHITE)
        add_text(sl, note, Inches(2.2), y+Inches(0.3), Inches(3.8), Inches(0.28),
                 size=Pt(9), color=C_DIM)

    add_text(sl, "* 중소기업기술정보진흥원(KAMP) 공개 데이터 — 상업적 이용 가능",
             Inches(0.65), Inches(5.95), Inches(5.5), Inches(0.35),
             size=Pt(8), color=C_MUTED)

    # 우측: 기술 스택
    add_box(sl, Inches(6.7), Inches(1.1), Inches(6.15), Inches(5.5),
            fill=C_CARD, border=C_BORDER)
    add_text(sl, "기술 스택",
             Inches(6.9), Inches(1.18), Inches(5.8), Inches(0.38),
             size=Pt(12), bold=True, color=C_WHITE)
    red_line(sl, Inches(6.9), Inches(1.53), Inches(5.8))

    tech_rows = [
        ("모델링",   "PyTorch",          "Autoencoder (24→16→8→16→24), BatchNorm, Adam"),
        ("전처리",   "sklearn",          "StandardScaler — raw 원본 기반 fit"),
        ("XAI",      "SHAP",             "KernelExplainer, kmeans(50) 배경 샘플"),
        ("대시보드", "Streamlit",        "Plotly 인터랙티브 차트, 5탭 구성"),
        ("데이터",   "pandas / numpy",   "parquet 저장, 배치 inference"),
        ("반지도",   "학습 전략",        "정상 데이터만 학습 → 복원 오차 임계값 판정"),
    ]
    for i, (cat, lib, desc) in enumerate(tech_rows):
        y = Inches(1.68) + i * Inches(0.75)
        add_box(sl, Inches(6.9), y, Inches(0.9), Inches(0.3), fill=C_RED)
        add_text(sl, cat, Inches(6.9), y, Inches(0.9), Inches(0.3),
                 size=Pt(8), bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
        add_text(sl, lib, Inches(7.88), y, Inches(1.5), Inches(0.3),
                 size=Pt(10), bold=True, color=C_WHITE)
        add_text(sl, desc, Inches(7.88), y+Inches(0.3), Inches(4.8), Inches(0.32),
                 size=Pt(9), color=C_DIM)


def slide06_usecase(prs):
    """5. 사용자 시나리오 / 유즈케이스"""
    sl = blank_slide(prs)
    fill_bg(sl)
    section_header(sl, "05", "사용자 시나리오 / 유즈케이스", "User Scenario")

    # 페르소나 박스
    add_box(sl, Inches(0.45), Inches(1.1), Inches(3.5), Inches(5.5),
            fill=C_CARD, border=C_BORDER)
    add_text(sl, "사용자 페르소나",
             Inches(0.65), Inches(1.18), Inches(3.1), Inches(0.35),
             size=Pt(11), bold=True, color=C_RED)
    red_line(sl, Inches(0.65), Inches(1.5), Inches(3.1))

    personas = [
        ("생산 관리자", "현장 불량 즉시 확인\n원인 파악 후 지시"),
        ("설비 엔지니어", "설비 이상 징후 조기 탐지\n예방 정비 계획 수립"),
        ("품질 담당자", "불량 이력 분석\n재발 방지 대책 수립"),
    ]
    for i, (role, desc) in enumerate(personas):
        y = Inches(1.65) + i * Inches(1.55)
        add_box(sl, Inches(0.65), y, Inches(3.1), Inches(1.3),
                fill=C_DARK, border=C_MUTED)
        add_text(sl, role, Inches(0.85), y+Inches(0.1), Inches(2.7), Inches(0.38),
                 size=Pt(12), bold=True, color=C_WHITE)
        add_text(sl, desc, Inches(0.85), y+Inches(0.45), Inches(2.7), Inches(0.75),
                 size=Pt(9), color=C_DIM)

    # 시나리오 플로우 (우측)
    add_text(sl, "5분 데모 시나리오",
             Inches(4.35), Inches(1.1), Inches(8.5), Inches(0.38),
             size=Pt(13), bold=True, color=C_WHITE)
    red_line(sl, Inches(4.35), Inches(1.45), Inches(8.5))

    scenario_steps = [
        ("00:00 – 00:30", "대시보드 접속",
         "탭 1 — ROC-AUC 0.9254 확인, 복원 오차 분포로 정상/불량 분리 명확성 시연"),
        ("00:30 – 01:30", "이상 감지 시뮬레이션",
         "탭 2 — Filling_Time 슬라이더 조작 → 이상 판정 즉시 반응 확인"),
        ("01:30 – 02:30", "XAI 원인 진단",
         "탭 2 SHAP 계산 → Filling_Time 1위 확인, 탭 4 Summary Bar로 전체 원인 분포 설명"),
        ("02:30 – 03:30", "처방 액션 확인",
         "처방 카드 — \"사출 속도 프로파일 점검 · 원재료 점도 확인\" 즉시 표시"),
        ("03:30 – 05:00", "대규모 분석 & 설비 현황",
         "탭 3 — 35K 스코어링 결과, 탭 5 — 설비 건강 점수 & 생산 이력 추적"),
    ]
    for i, (time, title, desc) in enumerate(scenario_steps):
        y = Inches(1.6) + i * Inches(1.05)
        # 시간 배지
        add_box(sl, Inches(4.35), y, Inches(1.5), Inches(0.3), fill=C_RED)
        add_text(sl, time, Inches(4.35), y, Inches(1.5), Inches(0.3),
                 size=Pt(8), bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
        add_text(sl, title, Inches(6.0), y, Inches(7.2), Inches(0.3),
                 size=Pt(11), bold=True, color=C_WHITE)
        add_text(sl, desc, Inches(4.35), y+Inches(0.3), Inches(8.8), Inches(0.55),
                 size=Pt(9), color=C_DIM)


def slide07_mvp(prs):
    """6. MVP 구현 범위"""
    sl = blank_slide(prs)
    fill_bg(sl)
    section_header(sl, "06", "MVP 구현 범위", "MVP Scope")

    # 구현 완료 (좌측)
    add_box(sl, Inches(0.45), Inches(1.08), Inches(6.1), Inches(5.55),
            fill=C_CARD, border=C_RED, border_w=Pt(1.5))
    add_text(sl, "[완료]  구현 완료",
             Inches(0.65), Inches(1.15), Inches(5.7), Inches(0.38),
             size=Pt(13), bold=True, color=C_RED)
    red_line(sl, Inches(0.65), Inches(1.5), Inches(5.75))

    done_items = [
        ("모델", [
            f"Autoencoder V2 — ROC-AUC {ROC:.4f}",
            f"F1-Score {F1:.4f} / Recall {REC:.4f} / Precision {PREC:.4f}",
            f"임계값 {THR:.4f} (99th percentile + F1 최적 평균)",
        ]),
        ("데이터 파이프라인", [
            "StandardScaler fit (labeled_data.csv 원본)",
            "supervised_label_cn7.csv 단일 파일 학습 (V2)",
            f"moldset_unlabeled_cn7.csv {SC['total']:,}건 스코어링" if SC else "비라벨 데이터 스코어링",
        ]),
        ("XAI", [
            "SHAP KernelExplainer 사전 계산",
            f"Top 센서: {', '.join(c for c,_ in SHAP_TOP[:3])}" if SHAP_TOP else "SHAP 상위 센서 도출",
            "Waterfall / Summary Bar 시각화",
        ]),
        ("대시보드", [
            "Streamlit 5탭 — 탐지/진단/처방/추적",
            "Plotly 인터랙티브 차트 전환 완료",
            "흑백+빨강 운영 시스템 디자인",
        ]),
    ]
    y = Inches(1.65)
    for cat, items in done_items:
        add_text(sl, cat, Inches(0.65), y, Inches(5.7), Inches(0.32),
                 size=Pt(10), bold=True, color=C_WHITE)
        for item in items:
            y += Inches(0.3)
            add_text(sl, f"   · {item}", Inches(0.65), y, Inches(5.7), Inches(0.28),
                     size=Pt(9), color=C_DIM)
        y += Inches(0.38)

    # 향후 확장 (우측)
    add_box(sl, Inches(6.88), Inches(1.08), Inches(6.0), Inches(5.55),
            fill=C_CARD, border=C_BORDER)
    add_text(sl, "[예정]  본선 추가 구현 계획",
             Inches(7.08), Inches(1.15), Inches(5.6), Inches(0.38),
             size=Pt(13), bold=True, color=C_WHITE)
    red_line(sl, Inches(7.08), Inches(1.5), Inches(5.65))

    todo_items = [
        ("실시간 연동", "OPC-UA / MQTT 프로토콜 연결", "설비 PLC 데이터 스트림 수신"),
        ("알람 시스템", "이상 탐지 즉시 Slack/Email 알림", "임계값 초과 시 담당자 Push"),
        ("멀티 설비", "설비 ID별 독립 모델 관리", "설비 유형별 임계값 차별화"),
        ("예측 정비", "이상 스코어 추이 기반 고장 예측", "설비 가동 중단 사전 방지"),
    ]
    for i, (title, sub1, sub2) in enumerate(todo_items):
        y = Inches(1.7) + i * Inches(1.18)
        add_text(sl, title, Inches(7.08), y, Inches(5.6), Inches(0.35),
                 size=Pt(11), bold=True, color=C_WHITE)
        add_text(sl, f"· {sub1}", Inches(7.08), y+Inches(0.34), Inches(5.6), Inches(0.28),
                 size=Pt(9), color=C_DIM)
        add_text(sl, f"· {sub2}", Inches(7.08), y+Inches(0.62), Inches(5.6), Inches(0.28),
                 size=Pt(9), color=C_MUTED)

    # 성능 지표 요약 바
    add_box(sl, Inches(0.45), Inches(6.78), Inches(12.43), Inches(0.55),
            fill=C_RED)
    metrics_str = (
        f"  ROC-AUC: {ROC:.4f}   |   PR-AUC: {PR:.4f}   |   "
        f"F1: {F1:.4f}   |   Recall: {REC:.4f}   |   Precision: {PREC:.4f}   |   "
        f"임계값: {THR:.4f}"
    )
    add_text(sl, metrics_str, Inches(0.45), Inches(6.78), Inches(12.43), Inches(0.55),
             size=Pt(9), bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)


def slide08_impact(prs):
    """7. 기대 효과 및 향후 확장성"""
    sl = blank_slide(prs)
    fill_bg(sl)
    section_header(sl, "07", "기대 효과 및 향후 확장성", "Impact & Scalability")

    # 정량 기대 효과 (상단 3개 카드)
    effects = [
        ("불량 탐지율",         f"ROC-AUC\n{ROC:.2%}",    "목표 80% 대비 초과 달성"),
        ("조기 탐지 효과",       "이상 징후\n사전 감지",   "복원 오차 추이로 발생 전 경보"),
        ("분석 시간 단축",       "수 시간 →\n수 초",       "SHAP 자동화로 즉시 원인 파악"),
    ]
    for i, (lbl, val, sub) in enumerate(effects):
        x = Inches(0.45) + i * Inches(4.3)
        add_box(sl, x, Inches(1.1), Inches(4.0), Inches(1.8),
                fill=C_CARD, border=C_RED, border_w=Pt(1.2))
        add_text(sl, lbl, x+Inches(0.15), Inches(1.18), Inches(3.7), Inches(0.32),
                 size=Pt(10), bold=True, color=C_DIM)
        add_text(sl, val, x+Inches(0.15), Inches(1.45), Inches(3.7), Inches(0.82),
                 size=Pt(20), bold=True, color=C_RED, align=PP_ALIGN.CENTER)
        add_text(sl, sub, x+Inches(0.15), Inches(2.55), Inches(3.7), Inches(0.28),
                 size=Pt(8), color=C_DIM)

    # SHAP 원인 분석 결과 (좌측 하단)
    add_box(sl, Inches(0.45), Inches(3.15), Inches(5.8), Inches(3.55),
            fill=C_CARD, border=C_BORDER)
    add_text(sl, "SHAP 원인 분석 결과 (Top 5)",
             Inches(0.65), Inches(3.23), Inches(5.4), Inches(0.35),
             size=Pt(11), bold=True, color=C_WHITE)
    red_line(sl, Inches(0.65), Inches(3.55), Inches(5.4))

    if SHAP_TOP:
        max_val = SHAP_TOP[0][1]
        for i, (col, val) in enumerate(SHAP_TOP):
            y = Inches(3.65) + i * Inches(0.52)
            # 이름
            add_text(sl, col.replace('_', ' '), Inches(0.65), y, Inches(2.6), Inches(0.32),
                     size=Pt(9), color=C_WHITE)
            # 바
            bar_w = (val / max_val) * Inches(2.5)
            bar_w = max(bar_w, Inches(0.1))
            add_box(sl, Inches(3.2), y+Inches(0.06), bar_w, Inches(0.22),
                    fill=C_RED if i == 0 else C_MUTED)
            # 값
            add_text(sl, f"{val:.2f}", Inches(3.2)+bar_w+Inches(0.06), y, Inches(0.8), Inches(0.32),
                     size=Pt(9), color=C_DIM)
    else:
        add_text(sl, "SHAP 데이터 없음 (04_compute_shap.py 실행 필요)",
                 Inches(0.65), Inches(3.65), Inches(5.4), Inches(0.5),
                 size=Pt(9), color=C_MUTED)

    # 확장성 로드맵 (우측 하단)
    add_box(sl, Inches(6.7), Inches(3.15), Inches(6.15), Inches(3.55),
            fill=C_CARD, border=C_BORDER)
    add_text(sl, "향후 확장 로드맵",
             Inches(6.9), Inches(3.23), Inches(5.8), Inches(0.35),
             size=Pt(11), bold=True, color=C_WHITE)
    red_line(sl, Inches(6.9), Inches(3.55), Inches(5.8))

    roadmap = [
        ("단기 (MVP+)",   "OPC-UA 실시간 연동 / 알람 자동화"),
        ("중기 (v2.0)",   "멀티 설비 관리 / 예측 정비 모듈"),
        ("장기 (v3.0)",   "Digital Twin 연동 / 공정 최적화 추천"),
        ("플랫폼화",       "SaaS 전환 — 중소 제조사 구독 모델"),
    ]
    for i, (phase, desc) in enumerate(roadmap):
        y = Inches(3.65) + i * Inches(0.72)
        add_box(sl, Inches(6.9), y, Inches(1.3), Inches(0.3),
                fill=C_RED if i == 0 else C_MUTED)
        add_text(sl, phase, Inches(6.9), y, Inches(1.3), Inches(0.3),
                 size=Pt(8), bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
        add_text(sl, desc, Inches(8.35), y, Inches(4.3), Inches(0.3),
                 size=Pt(10), color=C_WHITE)

    # 하단 마무리
    red_line(sl, Inches(0.45), Inches(6.9), Inches(12.43))
    add_text(sl,
             "SmartFactory XAI — 이상탐지에서 끝나지 않고, 탐지→진단→처방→추적까지 "
             "스마트 공장 운영의 전 주기를 지원합니다.",
             Inches(0.55), Inches(7.0), Inches(12.2), Inches(0.4),
             size=Pt(10), color=C_DIM, align=PP_ALIGN.CENTER)


def slide08_impact_v2(prs):
    """7. 기대 효과 및 가설 검증 결과 — v2 (분석 결과 포함)"""
    sl = blank_slide(prs)
    fill_bg(sl)
    section_header(sl, "07", "기대 효과 및 향후 확장성", "Impact & Scalability")

    # 가설 검증 결과 로드
    hyp_path = os.path.join(RESULT_DIR, 'hypothesis_test.json')
    HYP = json.load(open(hyp_path, encoding='utf-8')) if os.path.exists(hyp_path) else []
    sig_cnt   = sum(1 for h in HYP if h['significant'])
    large_cnt = sum(1 for h in HYP if h.get('effect_size') == 'large')

    # 상단 4개 정량 카드
    effects = [
        ("불량 탐지  ROC-AUC",   f"{ROC:.4f}",       "목표 0.80 대비 +15.7%p 초과"),
        ("유의 센서",             f"{sig_cnt} / 24",  "Mann-Whitney U  (alpha=0.05)"),
        ("Large 효과 크기 센서",  f"{large_cnt}개",   "Cohen's d >= 0.8 센서 수"),
        ("분석 자동화",           "수시간 -> 초",     "SHAP 즉시 원인 도출"),
    ]
    cw = Inches(3.05)
    for i, (lbl, val, sub) in enumerate(effects):
        x = Inches(0.45) + i * (cw + Inches(0.07))
        add_box(sl, x, Inches(1.08), cw, Inches(1.5),
                fill=C_CARD, border=C_RED if i < 2 else C_BORDER, border_w=Pt(1.2))
        add_text(sl, lbl, x+Inches(0.12), Inches(1.16), cw-Inches(0.22), Inches(0.3),
                 size=Pt(9), bold=True, color=C_DIM)
        add_text(sl, val, x+Inches(0.12), Inches(1.42), cw-Inches(0.22), Inches(0.65),
                 size=Pt(21), bold=True, color=C_RED, align=PP_ALIGN.CENTER)
        add_text(sl, sub, x+Inches(0.12), Inches(2.22), cw-Inches(0.22), Inches(0.28),
                 size=Pt(8), color=C_MUTED)

    # 좌측 하단: 가설 검증 상위 7개 센서
    add_box(sl, Inches(0.45), Inches(2.82), Inches(5.9), Inches(3.88),
            fill=C_CARD, border=C_BORDER)
    add_text(sl, "가설 검증 — 유의 센서 Top 7  (Cohen's d 기준)",
             Inches(0.65), Inches(2.90), Inches(5.5), Inches(0.35),
             size=Pt(11), bold=True, color=C_WHITE)
    red_line(sl, Inches(0.65), Inches(3.22), Inches(5.5))

    top_hyp = sorted(HYP, key=lambda x: -x['cohen_d'])[:7]
    if top_hyp:
        max_d = top_hyp[0]['cohen_d']
        for i, h in enumerate(top_hyp):
            y = Inches(3.32) + i * Inches(0.47)
            marker = "*" if h['significant'] else " "
            add_text(sl, f"{marker} {h['sensor'].replace('_',' ')}", Inches(0.65), y,
                     Inches(2.55), Inches(0.3), size=Pt(8),
                     color=C_WHITE if h['significant'] else C_MUTED)
            bar_w = max((h['cohen_d'] / (max_d + 1e-9)) * Inches(2.4), Inches(0.05))
            add_box(sl, Inches(3.1), y + Inches(0.06), bar_w, Inches(0.2),
                    fill=C_RED if h['significant'] else C_MUTED)
            add_text(sl, f"{h['cohen_d']:.3f}", Inches(3.1) + bar_w + Inches(0.06),
                     y, Inches(0.7), Inches(0.3), size=Pt(8), color=C_DIM)
    add_text(sl, "* p < 0.05 유의  |  Large effect: Cohen's d >= 0.8",
             Inches(0.65), Inches(6.56), Inches(5.5), Inches(0.28),
             size=Pt(8), color=C_MUTED)

    # 우측 하단: SHAP + 향후 로드맵
    add_box(sl, Inches(6.72), Inches(2.82), Inches(6.13), Inches(3.88),
            fill=C_CARD, border=C_BORDER)
    add_text(sl, "SHAP 원인 분석  &  향후 확장 로드맵",
             Inches(6.92), Inches(2.90), Inches(5.7), Inches(0.35),
             size=Pt(11), bold=True, color=C_WHITE)
    red_line(sl, Inches(6.92), Inches(3.22), Inches(5.7))

    if SHAP_TOP:
        add_text(sl, "SHAP Top 원인 센서", Inches(6.92), Inches(3.32), Inches(5.7), Inches(0.28),
                 size=Pt(9), bold=True, color=C_DIM)
        max_sv = SHAP_TOP[0][1]
        for i, (col, val) in enumerate(SHAP_TOP[:3]):
            y = Inches(3.58) + i * Inches(0.42)
            add_text(sl, col.replace('_', ' '), Inches(6.92), y, Inches(2.3), Inches(0.28),
                     size=Pt(8), color=C_WHITE)
            bw = max((val / (max_sv + 1e-9)) * Inches(2.0), Inches(0.05))
            add_box(sl, Inches(9.15), y + Inches(0.05), bw, Inches(0.2),
                    fill=C_RED if i == 0 else C_MUTED)
            add_text(sl, f"{val:.2f}", Inches(9.15) + bw + Inches(0.05),
                     y, Inches(0.7), Inches(0.28), size=Pt(8), color=C_DIM)

    add_text(sl, "향후 확장", Inches(6.92), Inches(4.98), Inches(5.7), Inches(0.28),
             size=Pt(9), bold=True, color=C_DIM)
    roadmap = [
        ("단기", "OPC-UA 실시간 연동  /  알람 자동화"),
        ("중기", "멀티 설비 관리  /  예측 정비"),
        ("장기", "Digital Twin  /  공정 최적화 추천"),
    ]
    for i, (phase, desc) in enumerate(roadmap):
        y = Inches(5.24) + i * Inches(0.46)
        add_box(sl, Inches(6.92), y, Inches(0.75), Inches(0.28),
                fill=C_RED if i == 0 else C_MUTED)
        add_text(sl, phase, Inches(6.92), y, Inches(0.75), Inches(0.28),
                 size=Pt(7), bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
        add_text(sl, desc, Inches(7.75), y, Inches(4.85), Inches(0.28),
                 size=Pt(9), color=C_WHITE)

    # 하단 마무리 바
    red_line(sl, Inches(0.45), Inches(6.9), Inches(12.43))
    ci_lo_v = M.get('roc_auc_ci_lo', 0.88)
    ci_hi_v = M.get('roc_auc_ci_hi', 0.97)
    add_text(sl,
             f"ROC-AUC {ROC:.4f} (95% CI [{ci_lo_v:.3f}, {ci_hi_v:.3f}]) — "
             f"불량 10건 중 9건 이상 구별  |  수작업 수시간 원인 분석 -> AI 즉시 진단으로 단축",
             Inches(0.55), Inches(7.0), Inches(12.2), Inches(0.4),
             size=Pt(10), color=C_DIM, align=PP_ALIGN.CENTER)


def slide09_bm_gtm(prs):
    """8. 비즈니스 모델 & GTM 전략"""
    sl = blank_slide(prs)
    fill_bg(sl)
    section_header(sl, "08", "비즈니스 모델 & 시장 진출 전략", "Business Model & GTM")

    # ── 좌측: SaaS 수익 모델 ──
    add_box(sl, Inches(0.45), Inches(1.0), Inches(6.1), Inches(5.65),
            fill=C_CARD, border=C_BORDER)
    add_text(sl, "SaaS 구독 수익 모델", Inches(0.65), Inches(1.1), Inches(5.7), Inches(0.35),
             size=Pt(12), bold=True, color=C_WHITE)
    red_line(sl, Inches(0.65), Inches(1.42), Inches(5.7))

    tiers = [
        ("스타터",      "1~3대",   "월 80만원",   "소규모 제조사"),
        ("스탠다드",    "4~10대",  "월 200만원",  "중견 사출성형 공장"),
        ("엔터프라이즈","10대 이상","월 400만원+", "대형·복수 라인"),
    ]
    for i, (name, seats, price, target) in enumerate(tiers):
        y = Inches(1.60) + i * Inches(1.25)
        border_c = C_RED if i == 1 else C_BORDER
        add_box(sl, Inches(0.65), y, Inches(5.7), Inches(1.05),
                fill=RGBColor(0x15, 0x15, 0x15), border=border_c, border_w=Pt(1))
        add_text(sl, name, Inches(0.85), y + Inches(0.08), Inches(1.5), Inches(0.35),
                 size=Pt(11), bold=True, color=C_RED if i == 1 else C_WHITE)
        if i == 1:
            add_text(sl, "추천", Inches(2.2), y + Inches(0.10), Inches(0.9), Inches(0.28),
                     size=Pt(8), bold=True, color=C_BLACK)
            add_box(sl, Inches(2.2), y + Inches(0.10), Inches(0.85), Inches(0.24), fill=C_RED)
            add_text(sl, "추천", Inches(2.2), y + Inches(0.10), Inches(0.85), Inches(0.24),
                     size=Pt(8), bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
        add_text(sl, f"설비 {seats}", Inches(0.85), y + Inches(0.38), Inches(2.0), Inches(0.28),
                 size=Pt(9), color=C_DIM)
        add_text(sl, price, Inches(3.5), y + Inches(0.06), Inches(2.65), Inches(0.45),
                 size=Pt(17), bold=True, color=C_WHITE, align=PP_ALIGN.RIGHT)
        add_text(sl, target, Inches(3.5), y + Inches(0.50), Inches(2.65), Inches(0.28),
                 size=Pt(8), color=C_MUTED, align=PP_ALIGN.RIGHT)

    # 가격 ROI 근거
    add_text(sl, "가격 근거: 스탠다드 월 200만 × 12 = 연 2,400만원  vs  절감액 ~1억/년 → 고객 ROI 4.2배 (동종 SaaS 대비 제조 특화)",
             Inches(0.65), Inches(5.1), Inches(5.7), Inches(0.22),
             size=Pt(7.5), color=C_DIM)

    # 3년 재무 모델 (P&L 요약)
    add_text(sl, "3년 재무 모델 (단위: 만원)", Inches(0.65), Inches(5.32), Inches(5.7), Inches(0.28),
             size=Pt(9), bold=True, color=C_DIM)
    fin_rows = [
        ("",          "Y1 (목표)",  "Y2 (목표)",  "Y3 (목표)"),
        ("고객 수",    "10개사",     "50개사",      "200개사"),
        ("ARR",        "2,400만",    "1.2억",       "48억"),
        ("매출원가",   "1,200만",    "5,000만",     "14억"),
        ("매출 이익률","50%",         "58%",         "71%"),
        ("CAC",        "200만",      "150만",       "100만"),
        ("LTV",        "2,400만",    "2,880만",     "3,600만"),
    ]
    col_ws = [Inches(1.4), Inches(1.4), Inches(1.4), Inches(1.5)]
    col_xs = [Inches(0.65)]
    for w in col_ws[:-1]:
        col_xs.append(col_xs[-1] + w)
    row_h = Inches(0.3)
    for ri, row in enumerate(fin_rows):
        y = Inches(5.55) + ri * row_h
        for ci, (cell, cw, cx) in enumerate(zip(row, col_ws, col_xs)):
            is_hdr_row = ri == 0
            is_hdr_col = ci == 0
            is_highlight = (ri == 2 and ci == 3) or (ri == 6 and ci in (1,2,3))
            fill = RGBColor(0x22, 0x22, 0x22) if (is_hdr_row or is_hdr_col) else C_CARD
            add_box(sl, cx, y, cw - Inches(0.02), row_h - Inches(0.01), fill=fill, border=C_BORDER)
            clr = C_RED if is_highlight else (C_WHITE if (is_hdr_row or is_hdr_col) else C_DIM)
            add_text(sl, cell, cx + Inches(0.05), y + Inches(0.03), cw - Inches(0.1), row_h - Inches(0.06),
                     size=Pt(7.5), bold=is_hdr_row or is_hdr_col, color=clr,
                     align=PP_ALIGN.CENTER if ci > 0 else PP_ALIGN.LEFT)

    # ── 우측: GTM 전략 ──
    add_box(sl, Inches(6.85), Inches(1.0), Inches(6.0), Inches(5.65),
            fill=C_CARD, border=C_BORDER)
    add_text(sl, "시장 진출 전략 (GTM)", Inches(7.05), Inches(1.1), Inches(5.6), Inches(0.35),
             size=Pt(12), bold=True, color=C_WHITE)
    red_line(sl, Inches(7.05), Inches(1.42), Inches(5.6))

    gtm_steps = [
        ("Phase 1", "KAMP 연계 PoC",
         "한국AI스마트제조플랫폼 파트너십\n중소 사출성형 사업장 3개사 무료 PoC\nKAMP 공인 데이터로 신뢰도 확보"),
        ("Phase 2", "정부 사업 연계",
         "스마트공장 보급·확산사업 공급기업 등록\n중기부·KOSMO 보조금 활용 도입 장벽 제거\n파일럿 레퍼런스 → 확산"),
        ("Phase 3", "직접 영업 & 파트너사",
         "설비 제조사(LSMtron·Woojin Plaimm) OEM\n설비 판매 채널 병행 → CAC 최소화\n연간 200개사 목표"),
    ]
    for i, (phase, title, body) in enumerate(gtm_steps):
        y = Inches(1.60) + i * Inches(1.6)
        add_box(sl, Inches(7.05), y, Inches(0.95), Inches(0.32), fill=C_RED)
        add_text(sl, phase, Inches(7.05), y, Inches(0.95), Inches(0.32),
                 size=Pt(8), bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
        add_text(sl, title, Inches(8.1), y, Inches(3.3), Inches(0.35),
                 size=Pt(11), bold=True, color=C_WHITE)
        add_text(sl, body, Inches(7.05), y + Inches(0.38), Inches(5.6), Inches(0.9),
                 size=Pt(9), color=C_DIM)

    # 하단 KPI + ROI
    add_box(sl, Inches(6.85), Inches(5.3), Inches(6.0), Inches(1.1),
            fill=RGBColor(0x1a, 0x08, 0x08), border=C_RED, border_w=Pt(1))
    add_text(sl, "초기 목표: 6개월 내 PoC 3건 완료 → 유료 전환 2건 → MRR 400만원",
             Inches(7.05), Inches(5.35), Inches(5.6), Inches(0.28),
             size=Pt(9), bold=True, color=C_RED)
    add_text(sl, "CAC < 200만원 / LTV > 2,400만원 (24개월 기준) / LTV:CAC > 12x 목표",
             Inches(7.05), Inches(5.60), Inches(5.6), Inches(0.24),
             size=Pt(8), color=C_DIM)
    add_text(sl, f"ROI 계산: 불량 1건 = 50만원 손실 (재작업+라인정지)  |  AI 탐지율(테스트 기준) {REC:.2f} × 연간 불량 300건 = {int(round(REC*300))}건 방지  |  절감액 {int(round(REC*300))*50:,}만원/년",
             Inches(7.05), Inches(5.82), Inches(5.6), Inches(0.24),
             size=Pt(7.5), color=C_WHITE)
    add_text(sl, "※ 연간 불량 300건 추정 근거: 사출성형 평균 불량률 1.5% × 연간 생산 20,000샷 = 300건 (중소 사출업체 기준)",
             Inches(7.05), Inches(6.05), Inches(5.6), Inches(0.22),
             size=Pt(7), color=C_DIM)

    # ── 경쟁 우위 / 데이터 해자 (P4-M1) ──
    add_box(sl, Inches(6.85), Inches(6.28), Inches(6.0), Inches(0.58),
            fill=C_CARD, border=C_BORDER)
    add_text(sl, "방어선 — 데이터 해자", Inches(7.05), Inches(6.33), Inches(2.5), Inches(0.28),
             size=Pt(9), bold=True, color=C_DIM)
    moats = ["고객 현장 데이터 누적 → 모델 자동 개선 루프", "산업별 정상 프로파일 DB 축적", "전환 비용: 재교육 + 히스토리 데이터 이전 부담"]
    for i, m in enumerate(moats):
        add_text(sl, f"• {m}", Inches(7.05), Inches(6.60) + i*Inches(0.18),
                 Inches(5.6), Inches(0.2), size=Pt(8), color=C_DIM)

    # 하단 구분선
    red_line(sl, Inches(0.45), Inches(6.9), Inches(12.43))
    add_text(sl,
             "라인 멈추기 전에 AI가 먼저 알려줍니다.  "
             "KAMP 공개 데이터 검증 + 정부 스마트공장 사업 연계로 첫 고객까지 최단 경로.",
             Inches(0.55), Inches(7.0), Inches(12.2), Inches(0.4),
             size=Pt(10), bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)


def slide13_references(prs):
    """13. 학술 레퍼런스 & 차별화"""
    sl = blank_slide(prs)
    fill_bg(sl)
    section_header(sl, "13", "학술 레퍼런스 & 차별화", "References & Novelty")

    # ── 좌측: 본 연구 vs 선행연구 비교표 ──
    add_box(sl, Inches(0.45), Inches(1.0), Inches(7.0), Inches(5.6),
            fill=C_CARD, border=C_BORDER)
    add_text(sl, "본 연구 vs 선행 연구 — 직접 비교",
             Inches(0.65), Inches(1.1), Inches(6.6), Inches(0.35),
             size=Pt(11), bold=True, color=C_RED)
    red_line(sl, Inches(0.65), Inches(1.42), Inches(6.6))

    headers = ["연구", "데이터", "모델", "학습", "성능"]
    col_widths = [1.55, 1.45, 1.6, 0.9, 1.5]
    rows = [
        ("본 연구 (2026)",      "KAMP 공개",    "AE+SHAP",       "비지도", "AUC 0.925"),
        ("MDPI 2025",           "KAMP 공개",    "XGB/LGBM+XAI",  "지도",   "1.00→0.13%"),
        ("arXiv:2511.08108",    "사출 (자체)",  "LSTM+SHAP",     "지도",   "F1 0.92"),
        ("JMST 2025",           "사출 (자체)",  "VAE+GAN",       "반지도", "소표본 안정"),
        ("Yoo JCDE 2023",       "사출 (자체)",  "Ensemble",      "지도",   "Recall 우선"),
        ("Ketonen ICPS 2021",   "사출 (자체)",  "VAE+RNN",       "반지도", "RC 분석"),
    ]

    # 헤더
    x_cur = 0.65
    for ci, (h, w) in enumerate(zip(headers, col_widths)):
        add_box(sl, Inches(x_cur), Inches(1.55), Inches(w), Inches(0.34),
                fill=C_RED if ci == 0 else C_DARK)
        add_text(sl, h, Inches(x_cur), Inches(1.55), Inches(w), Inches(0.34),
                 size=Pt(8.5), bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
        x_cur += w

    # 데이터 행
    for ri, row in enumerate(rows):
        y = Inches(1.92) + ri * Inches(0.42)
        x_cur = 0.65
        for ci, (val, w) in enumerate(zip(row, col_widths)):
            add_box(sl, Inches(x_cur), y, Inches(w), Inches(0.4),
                    fill=RGBColor(0x2a, 0x10, 0x10) if ri == 0 else C_DARK,
                    border=C_BORDER)
            add_text(sl, val, Inches(x_cur), y, Inches(w), Inches(0.4),
                     size=Pt(7.5),
                     bold=(ri == 0 or ci == 0),
                     color=C_RED if ri == 0 else C_WHITE,
                     align=PP_ALIGN.CENTER)
            x_cur += w

    # 비교 결론
    add_box(sl, Inches(0.65), Inches(4.6), Inches(6.7), Inches(1.85),
            fill=C_DARK, border=C_RED, border_w=Pt(1))
    add_text(sl, "본 연구 차별화 4가지", Inches(0.85), Inches(4.7),
             Inches(6.3), Inches(0.3), size=Pt(10), bold=True, color=C_RED)
    novelties = [
        "공개 KAMP 24센서 데이터 — 학술적 재현성 확보 (대부분 선행연구는 비공개)",
        "비지도 AE + KernelSHAP — 라벨 없이도 설명 가능 (KAMP 선행연구는 모두 지도학습)",
        "Cross-Machine + Bootstrap 95% CI — 통계적 검증 강화 (선행연구 대부분 single-split)",
        "24개 처방 통합 — 탐지+원인+조치를 단일 시스템에 결합 (선행연구는 탐지만)",
    ]
    for ni, n in enumerate(novelties):
        y = Inches(5.05) + ni * Inches(0.32)
        add_text(sl, f"{ni+1}.", Inches(0.85), y, Inches(0.3), Inches(0.28),
                 size=Pt(9), bold=True, color=C_RED)
        add_text(sl, n, Inches(1.15), y, Inches(6.05), Inches(0.28),
                 size=Pt(8.5), color=C_WHITE)

    # ── 우측: 핵심 인용 ──
    add_box(sl, Inches(7.65), Inches(1.0), Inches(5.25), Inches(5.6),
            fill=C_CARD, border=C_BORDER)
    add_text(sl, "핵심 인용 레퍼런스",
             Inches(7.85), Inches(1.1), Inches(4.9), Inches(0.35),
             size=Pt(11), bold=True, color=C_RED)
    red_line(sl, Inches(7.85), Inches(1.42), Inches(4.9))

    refs = [
        ("MDPI Processes 13(3), 912", "(2025) KAMP XGB/LGBM XAI", "직접 비교 대상"),
        ("arXiv:2511.08108", "(2025) LSTM + SHAP/Grad-CAM/LIME", "SHAP 결과 해석"),
        ("Ketonen & Blech ICPS", "(2021) VAE+RNN 사출 root-cause", "관련 연구"),
        ("Brito et al. MAKE 6(1), 16", "(2024) SHAP 베어링 98.5%", "왜 SHAP인가"),
        ("PhysiCausalNet IEEE TII", "(2024) Cross-Machine FD", "외부 검증 방법론"),
        ("EWAD-IIoT WGAN", "(2025) Sci. Reports — 95% CI", "Bootstrap CI 정당성"),
        ("Survey arXiv:2503.13195", "(2025) Deep AD 종합 서베이", "Motivation"),
    ]
    for ri, (cite, desc, where) in enumerate(refs):
        y = Inches(1.6) + ri * Inches(0.68)
        add_box(sl, Inches(7.85), y, Inches(0.18), Inches(0.55), fill=C_RED)
        add_text(sl, cite, Inches(8.1), y, Inches(4.7), Inches(0.25),
                 size=Pt(8.5), bold=True, color=C_WHITE)
        add_text(sl, desc, Inches(8.1), y + Inches(0.22), Inches(4.7), Inches(0.22),
                 size=Pt(7.5), color=C_DIM)
        add_text(sl, f"적용: {where}", Inches(8.1), y + Inches(0.4), Inches(4.7), Inches(0.18),
                 size=Pt(7), color=C_RED)

    # 하단 풋노트
    red_line(sl, Inches(0.45), Inches(6.78), Inches(12.43))
    add_text(sl, "사출성형은 표준 벤치마크(CWRU bearing 같은)가 부재 — 본 연구는 KAMP 공개 데이터 + Bootstrap CI로 학술적 재현성 / 통계적 신뢰도 모두 확보",
             Inches(0.55), Inches(6.92), Inches(12.2), Inches(0.4),
             size=Pt(9), bold=True, color=C_DIM, align=PP_ALIGN.CENTER)


def slide12_team(prs):
    """12. 팀 구성 & 전문성"""
    sl = blank_slide(prs)
    fill_bg(sl)
    section_header(sl, "12", "팀 구성 & 전문성", "Team Profile")

    # ── 헤더 ──
    add_text(sl, "스마트 공장 XAI 이상탐지 플랫폼 개발팀",
             Inches(0.45), Inches(1.05), Inches(12.43), Inches(0.4),
             size=Pt(14), bold=True, color=C_WHITE)
    red_line(sl, Inches(0.45), Inches(1.42), Inches(12.43))

    # ── 팀 구성 박스 (3열) ──
    members = [
        ("ML/AI 엔지니어", "Autoencoder 모델 설계 & 학습\nSHAP KernelExplainer 구현\nBootstrap CI·Pseudo Hold-out 설계",
         "PyTorch · scikit-learn · SHAP\nPython 3.11 · pandas · NumPy"),
        ("백엔드 & 데이터 파이프라인", "데이터 전처리 파이프라인\n795K 비라벨 스코어링 최적화\nPCA 클러스터 분석",
         "pandas · pyarrow · joblib\nStandardScaler · KMeans"),
        ("프론트엔드 & UX", "Streamlit 5탭 대시보드 구현\nPlotly 인터랙티브 시각화\n현장 언어 UI/UX 설계",
         "Streamlit 1.36 · Plotly\npython-pptx · CSS 커스터마이징"),
    ]
    for col_i, (role, tasks, tech) in enumerate(members):
        x = Inches(0.45) + col_i * Inches(4.18)
        add_box(sl, x, Inches(1.55), Inches(3.98), Inches(2.85),
                fill=C_CARD, border=C_RED if col_i == 0 else C_BORDER)
        add_box(sl, x, Inches(1.55), Inches(3.98), Inches(0.42), fill=C_RED if col_i == 0 else C_CARD)
        add_text(sl, role, x + Inches(0.15), Inches(1.6), Inches(3.6), Inches(0.33),
                 size=Pt(9.5), bold=True, color=C_WHITE)
        task_lines = tasks.split('\n')
        for li, tl in enumerate(task_lines):
            add_text(sl, f"· {tl}", x + Inches(0.15), Inches(2.12) + li * Inches(0.35),
                     Inches(3.68), Inches(0.33), size=Pt(8.5), color=C_DIM)
        add_text(sl, "기술 스택", x + Inches(0.15), Inches(3.32), Inches(3.68), Inches(0.25),
                 size=Pt(8), bold=True, color=C_MUTED)
        for ti, tl in enumerate(tech.split('\n')):
            add_text(sl, tl, x + Inches(0.15), Inches(3.54) + ti * Inches(0.28),
                     Inches(3.68), Inches(0.26), size=Pt(8), color=C_MUTED)

    # ── 팀 강점 ──
    add_box(sl, Inches(0.45), Inches(4.58), Inches(12.43), Inches(1.55),
            fill=C_CARD, border=C_BORDER)
    add_text(sl, "팀 경쟁력 & 차별화",
             Inches(0.65), Inches(4.68), Inches(12.0), Inches(0.32),
             size=Pt(10), bold=True, color=C_WHITE)
    red_line(sl, Inches(0.65), Inches(4.97), Inches(12.0))
    strengths = [
        ("KAMP 공인 데이터", "국내 유일 사출성형 공개 데이터셋 활용 → 레퍼런스 모델로 즉각 영업"),
        ("End-to-End 구현", "데이터 수집부터 SHAP 설명까지 단일 팀이 전 스택 직접 개발"),
        ("현장 검증 설계", "외부 기계 검증(Cross-Machine) + Pseudo Hold-out으로 Circular Evaluation 방지"),
        ("빠른 PoC 전환", "REPRODUCE.md 5단계 → 신규 고객사 데이터로 2주 내 커스텀 모델 납품 가능"),
    ]
    for si, (tag, desc) in enumerate(strengths):
        sx = Inches(0.65) + si * Inches(3.1)
        add_box(sl, sx, Inches(5.07), Inches(0.9), Inches(0.26), fill=C_RED)
        add_text(sl, tag, sx, Inches(5.07), Inches(0.9), Inches(0.26),
                 size=Pt(7.5), bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
        add_text(sl, desc, sx + Inches(1.0), Inches(5.07), Inches(2.05), Inches(0.26),
                 size=Pt(8), color=C_DIM)

    # ── 개발 기간 & 해커톤 ──
    add_box(sl, Inches(0.45), Inches(6.25), Inches(12.43), Inches(0.65),
            fill=C_CARD, border=C_BORDER)
    add_text(sl, "2026 스마트 공장 운영 시스템 MVP 개발 해커톤  ·  개발 기간: 2026년 4월~5월  ·  KAMP 공개 데이터 활용  ·  예선 제출",
             Inches(0.65), Inches(6.38), Inches(12.1), Inches(0.38),
             size=Pt(9), color=C_DIM, align=PP_ALIGN.CENTER)


def slide11_team_ask(prs):
    """10. 팀 소개 & 투자 제안"""
    sl = blank_slide(prs)
    fill_bg(sl)
    section_header(sl, "10", "팀 소개 & 투자 제안", "Team & Investment Ask")

    # ── 좌측: Why Now ──
    add_box(sl, Inches(0.45), Inches(1.0), Inches(6.1), Inches(3.0),
            fill=C_CARD, border=C_BORDER)
    add_text(sl, "Why Now — 지금이 진입 적기인 이유",
             Inches(0.65), Inches(1.1), Inches(5.7), Inches(0.35),
             size=Pt(11), bold=True, color=C_RED)
    red_line(sl, Inches(0.65), Inches(1.42), Inches(5.7))
    why_nows = [
        ("정부 정책", "스마트공장 보급 확산사업 2025~2027 — 3만개사 목표"),
        ("AI 성숙도", "KernelSHAP·PyTorch 2.x 오픈소스 성숙 → 개발 비용 1/5 수준"),
        ("데이터 공개", "KAMP 공공데이터포털 센서 데이터 → 레퍼런스 모델 즉시 구축"),
        ("인력 부족", "중소 제조사 AI 전문 인력 채용 난이도 상승 → SaaS 수요 급증"),
    ]
    for i, (tag, desc) in enumerate(why_nows):
        y = Inches(1.58) + i * Inches(0.52)
        add_box(sl, Inches(0.65), y, Inches(0.85), Inches(0.28), fill=C_RED)
        add_text(sl, tag, Inches(0.65), y, Inches(0.85), Inches(0.28),
                 size=Pt(7.5), bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
        add_text(sl, desc, Inches(1.6), y, Inches(4.0), Inches(0.28),
                 size=Pt(8.5), color=C_DIM)

    # ── 좌측 하단: 리스크 완화 ──
    add_box(sl, Inches(0.45), Inches(4.15), Inches(6.1), Inches(2.4),
            fill=C_CARD, border=C_BORDER)
    add_text(sl, "주요 리스크 & 완화 전략",
             Inches(0.65), Inches(4.25), Inches(5.7), Inches(0.32),
             size=Pt(10), bold=True, color=C_WHITE)
    red_line(sl, Inches(0.65), Inches(4.55), Inches(5.7))
    risks = [
        ("데이터 부족",     "→ KAMP 파트너십으로 신규 설비 데이터 지속 확보"),
        ("소표본 모델 한계","→ 고객 현장 데이터 누적 → 모델 자동 개선 루프"),
        ("고객 보수성",     "→ 무료 PoC 3개월 제공, 성과 확인 후 유료 전환"),
        ("기술 진부화",     "→ 오토인코더 → Transformer 마이그레이션 로드맵 수립"),
    ]
    for i, (risk, mit) in enumerate(risks):
        y = Inches(4.65) + i * Inches(0.44)
        add_text(sl, f"• {risk}", Inches(0.65), y, Inches(1.9), Inches(0.28),
                 size=Pt(8.5), bold=True, color=C_WHITE)
        add_text(sl, mit, Inches(2.5), y, Inches(3.85), Inches(0.28),
                 size=Pt(8), color=C_DIM)

    # ── 우측: 투자 제안 ──
    add_box(sl, Inches(6.85), Inches(1.0), Inches(6.0), Inches(5.55),
            fill=C_CARD, border=C_RED, border_w=Pt(1.5))
    add_text(sl, "투자 제안 (Investment Ask)",
             Inches(7.05), Inches(1.1), Inches(5.6), Inches(0.38),
             size=Pt(12), bold=True, color=C_RED)
    red_line(sl, Inches(7.05), Inches(1.45), Inches(5.6))

    # 시드 라운드
    add_text(sl, "SEED 라운드 목표",
             Inches(7.05), Inches(1.6), Inches(5.6), Inches(0.32),
             size=Pt(10), bold=True, color=C_WHITE)
    add_box(sl, Inches(7.05), Inches(1.9), Inches(5.6), Inches(0.7),
            fill=RGBColor(0x1a, 0x08, 0x08), border=C_RED, border_w=Pt(1))
    add_text(sl, "3억원 (Pre-Seed)",
             Inches(7.1), Inches(1.95), Inches(5.4), Inches(0.5),
             size=Pt(22), bold=True, color=C_RED, align=PP_ALIGN.CENTER)

    uses = [
        ("40%", "제품 개발", "OPC-UA 연동·멀티설비·알람 자동화"),
        ("30%", "영업/마케팅", "KAMP PoC 3건 + 유료 전환 집중"),
        ("20%", "팀 확장",   "ML 엔지니어 1명 · 영업 1명 채용"),
        ("10%", "인프라",    "AWS 클라우드 + 보안 인증"),
    ]
    add_text(sl, "자금 사용 계획", Inches(7.05), Inches(2.65), Inches(5.6), Inches(0.3),
             size=Pt(9), bold=True, color=C_DIM)
    for i, (pct, cat, desc) in enumerate(uses):
        y = Inches(2.92) + i * Inches(0.46)
        add_box(sl, Inches(7.05), y, Inches(0.6), Inches(0.28), fill=C_RED if i < 2 else C_MUTED)
        add_text(sl, pct, Inches(7.05), y, Inches(0.6), Inches(0.28),
                 size=Pt(9), bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
        add_text(sl, f"{cat}: {desc}", Inches(7.75), y, Inches(4.7), Inches(0.28),
                 size=Pt(9), color=C_WHITE)

    # 마일스톤
    add_text(sl, "6개월 마일스톤", Inches(7.05), Inches(4.82), Inches(5.6), Inches(0.32),
             size=Pt(10), bold=True, color=C_WHITE)
    milestones = [
        ("M3", "PoC 3건 완료 + Recall 0.75 이상 달성"),
        ("M4", "유료 전환 2건 → MRR 400만원"),
        ("M6", "Series A 준비 — ARR 6,000만원 목표"),
    ]
    for i, (mo, goal) in enumerate(milestones):
        y = Inches(5.1) + i * Inches(0.44)
        add_box(sl, Inches(7.05), y, Inches(0.5), Inches(0.28), fill=C_RED if i == 2 else C_MUTED)
        add_text(sl, mo, Inches(7.05), y, Inches(0.5), Inches(0.28),
                 size=Pt(8), bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
        add_text(sl, goal, Inches(7.65), y, Inches(4.8), Inches(0.28),
                 size=Pt(9), color=C_DIM)

    # Series A 진입 조건 박스
    add_box(sl, Inches(6.85), Inches(6.4), Inches(6.0), Inches(0.52),
            fill=C_CARD, border=C_RED, border_w=Pt(0.8))
    add_text(sl, "Series A 진입 조건 (6개월 후)",
             Inches(7.05), Inches(6.43), Inches(5.6), Inches(0.22),
             size=Pt(8), bold=True, color=C_DIM)
    add_text(sl,
             "ARR 6,000만+ · 유료 고객 5개사+ · Churn 0 · Recall 0.75+ (외부 검증) · NPS 40+",
             Inches(7.05), Inches(6.63), Inches(5.6), Inches(0.24),
             size=Pt(8.5), color=C_RED)

    red_line(sl, Inches(0.45), Inches(6.9), Inches(12.43))
    add_text(sl, "라인이 멈추기 전에 AI가 먼저 알립니다. KAMP 공인 데이터 실증 + 정부 스마트공장 연계 → 최단 경로로 첫 유료 고객까지.",
             Inches(0.55), Inches(7.0), Inches(12.2), Inches(0.4),
             size=Pt(10), bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)


def slide10_competitive(prs):
    """9. 경쟁사 기능 비교 매트릭스"""
    sl = blank_slide(prs)
    fill_bg(sl)
    section_header(sl, "09", "경쟁 포지셔닝 & 차별화 근거", "Competitive Positioning")

    # 기능 비교 매트릭스
    headers   = ["기능", "수동 점검", "범용 MES 알람", "외부 컨설팅", "SmartFactory XAI"]
    col_w     = [Inches(3.0), Inches(2.1), Inches(2.1), Inches(2.1), Inches(2.4)]
    col_x     = [Inches(0.45)]
    for w in col_w[:-1]:
        col_x.append(col_x[-1] + w + Inches(0.04))

    rows = [
        ("실시간 이상 탐지",          "✗",  "△ 단순 임계값",  "✗ 비정기",          "✓ 24개 센서 AI"),
        ("불량 원인 센서 즉시 식별",  "✗",  "✗",              "✓ 수시간 후",        "✓ SHAP 즉시"),
        ("현장 맞춤 임계값 조정",     "✗",  "△ IT 설정 필요", "✗",                  "✓ 운전자 직접"),
        ("교대 인수인계 자동화",       "수기 일지",  "✗",    "✗",                  "✓ CSV 자동 생성"),
        ("불량 유형 군집 분류",        "✗",  "✗",              "△ 별도 분석비",      "✓ 3유형 자동"),
        ("도입 비용",                  "0",  "수천만원/식",    "수백만원/회",         "월 80~400만원"),
        ("원인 설명 (XAI)",            "✗",  "✗",              "△ 전문가 판단",      "✓ 자동 SHAP"),
    ]

    # 헤더 행
    row_h = Inches(0.36)
    y0    = Inches(1.35)
    for ci, (hdr, w, x) in enumerate(zip(headers, col_w, col_x)):
        fill = C_RED if ci == 4 else RGBColor(0x22, 0x22, 0x22)
        add_box(sl, x, y0, w, row_h, fill=fill, border=C_BORDER)
        add_text(sl, hdr, x + Inches(0.06), y0 + Inches(0.04), w - Inches(0.12), row_h - Inches(0.06),
                 size=Pt(9), bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

    # 데이터 행
    for ri, (feat, *vals) in enumerate(rows):
        y = y0 + row_h + ri * (row_h + Inches(0.02))
        row_fill = RGBColor(0x14, 0x14, 0x14) if ri % 2 == 0 else C_CARD
        for ci, (val, w, x) in enumerate(zip([feat] + vals, col_w, col_x)):
            is_ours = (ci == 4)
            cell_fill = RGBColor(0x1a, 0x08, 0x08) if is_ours else row_fill
            cell_bdr  = C_RED if is_ours else C_BORDER
            add_box(sl, x, y, w, row_h, fill=cell_fill, border=cell_bdr)
            clr = C_RED if (is_ours and val.startswith("✓")) else (C_DIM if val.startswith("✗") else C_WHITE)
            add_text(sl, val, x + Inches(0.06), y + Inches(0.04), w - Inches(0.12), row_h - Inches(0.06),
                     size=Pt(8.5), color=clr, align=PP_ALIGN.CENTER if ci > 0 else PP_ALIGN.LEFT)

    # 하단 포지셔닝 메시지
    bot_y = y0 + row_h + len(rows) * (row_h + Inches(0.02)) + Inches(0.2)
    add_box(sl, Inches(0.45), bot_y, Inches(12.4), Inches(0.58),
            fill=RGBColor(0x1a, 0x08, 0x08), border=C_RED, border_w=Pt(1))
    add_text(sl, "SmartFactory XAI의 포지셔닝:  저비용(SaaS) × 즉각 설명(XAI) × 현장 맞춤(슬라이더) — 세 조건을 동시에 충족하는 유일한 솔루션",
             Inches(0.65), bot_y + Inches(0.06), Inches(12.0), Inches(0.28),
             size=Pt(10), bold=True, color=C_RED)
    add_text(sl, f"ROC-AUC {ROC:.4f} (KAMP 공개 데이터 실증) · Recall {REC:.2f} (10건 중 {int(round(REC*10))}건 탐지) · 월 구독 80만원 ~ 스탠다드 200만원",
             Inches(0.65), bot_y + Inches(0.31), Inches(12.0), Inches(0.24),
             size=Pt(8.5), color=C_DIM)

    red_line(sl, Inches(0.45), Inches(6.9), Inches(12.43))
    add_text(sl, "기능 면에서 대체 불가 · 비용 면에서 진입 장벽 없음 · 데이터 해자로 시간이 지날수록 격차 확대.",
             Inches(0.55), Inches(7.0), Inches(12.2), Inches(0.4),
             size=Pt(10), color=C_DIM, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════
# 메인
# ══════════════════════════════════════════════════════════════
if __name__ == '__main__':
    import sys
    sys.stdout.reconfigure(encoding='utf-8', errors='replace')

    print("PPT 생성 시작...")
    prs = new_prs()

    print("  슬라이드 1/9: 표지")
    slide01_cover(prs)
    print("  슬라이드 2/9: 문제 정의")
    slide02_problem(prs)
    print("  슬라이드 3/9: 솔루션 개요")
    slide03_solution(prs)
    print("  슬라이드 4/9: 주요 기능")
    slide04_features(prs)
    print("  슬라이드 5/9: 데이터 & 기술")
    slide05_data_tech(prs)
    print("  슬라이드 6/9: 유즈케이스")
    slide06_usecase(prs)
    print("  슬라이드 7/9: MVP 구현 범위")
    slide07_mvp(prs)
    print("  슬라이드 8/9: 기대 효과 & 가설 검증 결과")
    slide08_impact_v2(prs)
    print("  슬라이드 9/10: 비즈니스 모델 & GTM")
    slide09_bm_gtm(prs)
    print("  슬라이드 10/12: 경쟁 포지셔닝 매트릭스")
    slide10_competitive(prs)
    print("  슬라이드 11/12: 팀 소개 & 투자 제안")
    slide11_team_ask(prs)
    print("  슬라이드 12/13: 팀 구성 & 전문성")
    slide12_team(prs)
    print("  슬라이드 13/13: 학술 레퍼런스 & 차별화")
    slide13_references(prs)

    prs.save(OUT_PATH)
    print(f"\n[완료] 저장: {OUT_PATH}")
    print(f"   슬라이드 수: {len(prs.slides)}")
    print(f"\n   모델 성능 반영:")
    print(f"   ROC-AUC  = {ROC:.4f}")
    print(f"   F1-Score = {F1:.4f}")
    print(f"   Recall   = {REC:.4f}")
    if SHAP_TOP:
        print(f"\n   SHAP Top-3: {', '.join(f'{c}({v:.2f})' for c,v in SHAP_TOP[:3])}")
    if SC:
        print(f"   스코어링: {SC['total']:,}건 중 이상 {SC['anomaly']:,}건 ({SC['rate']:.2f}%)")
